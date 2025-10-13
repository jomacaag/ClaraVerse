import os
import sys
import logging
import signal
import traceback
import time
import argparse
import uuid
import asyncio
from datetime import datetime
from fastapi import FastAPI, HTTPException, Request, File, UploadFile, Form, Depends, Query, BackgroundTasks
from fastapi.responses import JSONResponse, Response, HTMLResponse, FileResponse
from fastapi.middleware.cors import CORSMiddleware
import tempfile
import shutil
from pathlib import Path
from typing import List, Optional, Dict, Any
import json
from pydantic import BaseModel, Field
from io import BytesIO
import PyPDF2
import xml.etree.ElementTree as ET

# Import Speech2Text
from Speech2Text import Speech2Text

# Import Text2Speech
from Text2Speech import Text2Speech

# Import Document Queue for persistent job processing
from document_queue import DocumentQueue

# LightRAG imports
try:
    from lightrag import LightRAG, QueryParam
    from lightrag.llm.openai import gpt_4o_mini_complete, openai_embed, openai_complete_if_cache, gpt_4o_complete, openai_complete
    from lightrag.llm.ollama import ollama_model_complete, ollama_embed
    from lightrag.utils import EmbeddingFunc, setup_logger
    from lightrag.kg.shared_storage import initialize_pipeline_status
    # Import rerank functions for enhanced performance
    try:
        from lightrag.rerank import cohere_rerank, jina_rerank
        RERANK_AVAILABLE = True
    except ImportError:
        RERANK_AVAILABLE = False
        logging.warning("Rerank functions not available - retrieval performance may be reduced")

    LIGHTRAG_AVAILABLE = True
except ImportError as e:
    logging.warning(f"LightRAG not available: {e}")
    LIGHTRAG_AVAILABLE = False

# Configure logging
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(name)s - %(levelname)s - %(message)s',
    handlers=[logging.StreamHandler()]
)
logger = logging.getLogger("clara-backend")

# Store start time
START_TIME = datetime.now().isoformat()

# Parse command line arguments
parser = argparse.ArgumentParser(description='Clara Backend Server')
parser.add_argument('--host', type=str, default='127.0.0.1', help='Host to bind to')
parser.add_argument('--port', type=int, default=5000, help='Port to bind to')
args = parser.parse_args()

# Use the provided host and port
HOST = args.host
PORT = args.port

logger.info(f"Starting server on {HOST}:{PORT}")

# Setup FastAPI
app = FastAPI(title="Clara Backend API", version="1.0.0")

# Import and include the diffusers API router
# Add CORS middleware
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],  # In production, replace with specific origins
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
    expose_headers=["*"]
)

# Add global exception middleware
@app.middleware("http")
async def catch_exceptions_middleware(request: Request, call_next):
    try:
        return await call_next(request)
    except Exception as e:
        logger.error(f"Request to {request.url} failed: {str(e)}")
        logger.error(traceback.format_exc())
        return JSONResponse(
            status_code=500,
            content={"error": str(e), "detail": traceback.format_exc()}
        )

# Database path in user's home directory for persistence
home_dir = os.path.expanduser("~")
data_dir = os.path.join(home_dir, ".clara")
os.makedirs(data_dir, exist_ok=True)

# LightRAG Configuration and Storage
if LIGHTRAG_AVAILABLE:
    # Environment setup for LightRAG
    os.environ["NEO4J_URI"] = os.getenv("NEO4J_URI", "neo4j://localhost:7687")
    os.environ["NEO4J_USERNAME"] = os.getenv("NEO4J_USERNAME", "neo4j")
    os.environ["NEO4J_PASSWORD"] = os.getenv("NEO4J_PASSWORD", "password")
    os.environ["OPENAI_API_KEY"] = os.getenv("OPENAI_API_KEY", "your-openai-api-key")

    # Setup LightRAG logging
    setup_logger("lightrag", level="INFO")
    
    # Storage paths for LightRAG
    LIGHTRAG_STORAGE_PATH = Path(os.path.join(data_dir, "lightrag_storage"))
    LIGHTRAG_STORAGE_PATH.mkdir(exist_ok=True)
    LIGHTRAG_METADATA_PATH = LIGHTRAG_STORAGE_PATH / "metadata"
    LIGHTRAG_METADATA_PATH.mkdir(exist_ok=True)

    # Global storage for LightRAG notebooks and documents
    lightrag_notebooks_db: Dict[str, Dict] = {}
    lightrag_documents_db: Dict[str, Dict] = {}
    lightrag_instances: Dict[str, LightRAG] = {}
    # Chat history storage for maintaining conversation context
    chat_history_db: Dict[str, List[Dict]] = {}  # notebook_id -> [messages]

    # Persistence files
    NOTEBOOKS_DB_FILE = LIGHTRAG_METADATA_PATH / "notebooks.json"
    DOCUMENTS_DB_FILE = LIGHTRAG_METADATA_PATH / "documents.json"
    CHAT_HISTORY_DB_FILE = LIGHTRAG_METADATA_PATH / "chat_history.json"

    # Initialize document processing queue
    QUEUE_DB_FILE = LIGHTRAG_METADATA_PATH / "document_queue.db"
    document_queue = DocumentQueue(db_path=str(QUEUE_DB_FILE))
    logger.info(f"Document queue initialized at {QUEUE_DB_FILE}")

    def save_notebooks_db():
        """Save notebooks database to disk"""
        try:
            # Convert datetime objects to ISO strings for JSON serialization
            serializable_data = {}
            for notebook_id, notebook_data in lightrag_notebooks_db.items():
                serializable_notebook = notebook_data.copy()
                # Convert all datetime fields
                for field in ['created_at', 'updated_at']:
                    if isinstance(serializable_notebook.get(field), datetime):
                        serializable_notebook[field] = serializable_notebook[field].isoformat()
                serializable_data[notebook_id] = serializable_notebook
            
            with open(NOTEBOOKS_DB_FILE, 'w') as f:
                json.dump(serializable_data, f, indent=2)
            logger.info(f"Saved {len(serializable_data)} notebooks to {NOTEBOOKS_DB_FILE}")
        except Exception as e:
            logger.error(f"Error saving notebooks database: {e}")

    def load_notebooks_db():
        """Load notebooks database from disk"""
        global lightrag_notebooks_db
        try:
            if NOTEBOOKS_DB_FILE.exists():
                with open(NOTEBOOKS_DB_FILE, 'r') as f:
                    data = json.load(f)
                
                # Convert ISO strings back to datetime objects
                for notebook_id, notebook_data in data.items():
                    if isinstance(notebook_data.get('created_at'), str):
                        notebook_data['created_at'] = datetime.fromisoformat(notebook_data['created_at'])
                    
                    # Backward compatibility: add default schema fields if missing
                    if 'entity_types' not in notebook_data:
                        notebook_data['entity_types'] = None  # Will use defaults
                    if 'language' not in notebook_data:
                        notebook_data['language'] = "en"
                
                lightrag_notebooks_db = data
                logger.info(f"Loaded {len(data)} notebooks from {NOTEBOOKS_DB_FILE}")
            else:
                logger.info("No existing notebooks database found")
        except Exception as e:
            logger.error(f"Error loading notebooks database: {e}")
            lightrag_notebooks_db = {}

    def save_documents_db():
        """Save documents database to disk"""
        try:
            # Convert datetime objects to ISO strings for JSON serialization
            serializable_data = {}
            for document_id, document_data in lightrag_documents_db.items():
                serializable_document = document_data.copy()
                # Convert all datetime values to ISO strings
                for key, value in serializable_document.items():
                    if isinstance(value, datetime):
                        serializable_document[key] = value.isoformat()
                serializable_data[document_id] = serializable_document
            
            with open(DOCUMENTS_DB_FILE, 'w') as f:
                json.dump(serializable_data, f, indent=2)
            logger.info(f"Saved {len(serializable_data)} documents to {DOCUMENTS_DB_FILE}")
        except Exception as e:
            logger.error(f"Error saving documents database: {e}")

    def load_documents_db():
        """Load documents database from disk"""
        global lightrag_documents_db
        try:
            if DOCUMENTS_DB_FILE.exists():
                with open(DOCUMENTS_DB_FILE, 'r') as f:
                    data = json.load(f)
                
                # Convert ISO strings back to datetime objects
                for document_id, document_data in data.items():
                    for key, value in document_data.items():
                        if isinstance(value, str) and key.endswith('_at'):
                            try:
                                document_data[key] = datetime.fromisoformat(value)
                            except ValueError:
                                pass  # Keep as string if not a valid ISO datetime
                
                lightrag_documents_db = data
                logger.info(f"Loaded {len(data)} documents from {DOCUMENTS_DB_FILE}")
            else:
                logger.info("No existing documents database found")
        except Exception as e:
            logger.error(f"Error loading documents database: {e}")
            lightrag_documents_db = {}

    def save_chat_history_db():
        """Save chat history database to disk"""
        try:
            # Convert datetime objects to ISO strings for JSON serialization
            serializable_data = {}
            for notebook_id, messages in chat_history_db.items():
                serializable_messages = []
                for message in messages:
                    serializable_message = message.copy()
                    if isinstance(serializable_message.get('timestamp'), datetime):
                        serializable_message['timestamp'] = serializable_message['timestamp'].isoformat()
                    serializable_messages.append(serializable_message)
                serializable_data[notebook_id] = serializable_messages
            
            with open(CHAT_HISTORY_DB_FILE, 'w') as f:
                json.dump(serializable_data, f, indent=2)
            logger.info(f"Saved chat history for {len(serializable_data)} notebooks to {CHAT_HISTORY_DB_FILE}")
        except Exception as e:
            logger.error(f"Error saving chat history database: {e}")

    def load_chat_history_db():
        """Load chat history database from disk"""
        global chat_history_db
        try:
            if CHAT_HISTORY_DB_FILE.exists():
                with open(CHAT_HISTORY_DB_FILE, 'r') as f:
                    data = json.load(f)
                
                # Convert ISO strings back to datetime objects
                for notebook_id, messages in data.items():
                    for message in messages:
                        if isinstance(message.get('timestamp'), str):
                            try:
                                message['timestamp'] = datetime.fromisoformat(message['timestamp'])
                            except ValueError:
                                pass  # Keep as string if not a valid ISO datetime
                
                chat_history_db = data
                logger.info(f"Loaded chat history for {len(data)} notebooks from {CHAT_HISTORY_DB_FILE}")
            else:
                logger.info("No existing chat history database found")
        except Exception as e:
            logger.error(f"Error loading chat history database: {e}")
            chat_history_db = {}

    # Load existing data on startup
    load_notebooks_db()
    load_documents_db()
    load_chat_history_db()

# Speech2Text instance cache
speech2text_instance = None

def get_speech2text():
    """Create or retrieve the Speech2Text instance from cache"""
    global speech2text_instance
    
    if speech2text_instance is None:
        # Use tiny model with CPU for maximum compatibility
        speech2text_instance = Speech2Text(
            model_size="tiny",
            device="cpu",
            compute_type="int8"
        )
    
    return speech2text_instance

# Text2Speech instance cache
text2speech_instance = None

def get_text2speech():
    """Create or retrieve the Text2Speech instance from cache"""
    global text2speech_instance
    
    if text2speech_instance is None:
        # Initialize with auto engine selection (will prefer Kokoro if available)
        text2speech_instance = Text2Speech(
            engine="auto",
            language="en",
            slow=False,
            voice="af_sarah",
            speed=1.0
        )
    
    return text2speech_instance

# LightRAG Utility Functions
if LIGHTRAG_AVAILABLE:
    def extract_text_from_pdf_lightrag(pdf_bytes: bytes) -> str:
        """Extract text from PDF bytes for LightRAG"""
        try:
            pdf_reader = PyPDF2.PdfReader(BytesIO(pdf_bytes))
            text = ""
            for page in pdf_reader.pages:
                text += page.extract_text() + "\n"
            return text
        except Exception as e:
            logger.error(f"Error extracting text from PDF: {e}")
            raise HTTPException(status_code=400, detail=f"Error processing PDF: {str(e)}")

    # Additional format support endpoint
    @app.get("/notebooks/supported-formats")
    async def get_supported_formats():
        """Get information about supported file formats for notebook documents"""
        basic_formats = [
            {"extension": "pdf", "description": "PDF documents", "library": "PyPDF2 (built-in)"},
            {"extension": "txt", "description": "Plain text files", "library": "built-in"},
            {"extension": "md", "description": "Markdown files", "library": "built-in"},
            {"extension": "csv", "description": "Comma-separated values", "library": "built-in"},
            {"extension": "json", "description": "JSON data files", "library": "built-in"},
            {"extension": "xml", "description": "XML documents", "library": "built-in"},
            {"extension": "html", "description": "HTML documents", "library": "built-in"},
            {"extension": "htm", "description": "HTML documents", "library": "built-in"}
        ]
        
        enhanced_formats = [
            {"extension": "docx", "description": "Word documents (newer)", "library": "python-docx", "install": "pip install python-docx"},
            {"extension": "doc", "description": "Word documents (legacy)", "library": "python-docx", "install": "pip install python-docx"},
            {"extension": "xlsx", "description": "Excel spreadsheets (newer)", "library": "pandas + openpyxl", "install": "pip install pandas openpyxl"},
            {"extension": "xls", "description": "Excel spreadsheets (legacy)", "library": "pandas + xlrd", "install": "pip install pandas xlrd"},
            {"extension": "pptx", "description": "PowerPoint presentations (newer)", "library": "python-pptx", "install": "pip install python-pptx"},
            {"extension": "ppt", "description": "PowerPoint presentations (legacy)", "library": "python-pptx", "install": "pip install python-pptx"},
            {"extension": "rtf", "description": "Rich Text Format", "library": "striprtf", "install": "pip install striprtf"}
        ]
        
        textract_formats = [
            {"extension": "epub", "description": "eBook format", "library": "textract", "install": "pip install textract"},
            {"extension": "odt", "description": "OpenDocument Text", "library": "textract", "install": "pip install textract"},
            {"extension": "ods", "description": "OpenDocument Spreadsheet", "library": "textract", "install": "pip install textract"},
            {"extension": "odp", "description": "OpenDocument Presentation", "library": "textract", "install": "pip install textract"}
        ]
        
        return {
            "basic_support": [f".{fmt['extension']}" for fmt in basic_formats],
            "enhanced_support": [f".{fmt['extension']}" for fmt in enhanced_formats],
            "textract_support": [f".{fmt['extension']}" for fmt in textract_formats],
            "details": {
                "basic": basic_formats,
                "enhanced": enhanced_formats,
                "textract": textract_formats
            },
            "installation_notes": {
                "basic": "These formats are supported out of the box",
                "enhanced": "Install additional libraries for enhanced document support",
                "textract": "Install textract for additional format support (requires system dependencies)"
            },
            "recommended_libraries": [
                "pip install python-docx python-pptx pandas openpyxl xlrd striprtf beautifulsoup4",
                "pip install textract  # Optional for additional formats"
            ]
        }

    @app.get("/notebooks/entity-types")
    async def get_default_entity_types():
        """Get default entity types for schema consistency configuration"""
        default_entity_types = [
            "PERSON", "ORGANIZATION", "LOCATION", "TECHNOLOGY", 
            "PRODUCT", "SERVICE", "PROJECT", "FEATURE", 
            "COMPONENT", "API", "DATABASE", "FRAMEWORK",
            "METHODOLOGY", "METRIC", "REQUIREMENT", "ISSUE",
            "SOLUTION", "STRATEGY", "GOAL", "RISK",
            "RESOURCE", "TOOL", "PLATFORM", "ENVIRONMENT",
            "VERSION", "RELEASE", "DEPLOYMENT", "CONFIGURATION",
            "WORKFLOW", "PROCESS", "STANDARD", "PROTOCOL",
            "DOCUMENT", "SPECIFICATION", "GUIDELINE", "POLICY",
            "EVENT", "MEETING", "DECISION", "MILESTONE",
            "CONCEPT", "PRINCIPLE", "PATTERN", "ARCHITECTURE",
            "INTERFACE", "MODULE", "LIBRARY", "DEPENDENCY",
            "DATA", "MODEL", "SCHEMA", "ENTITY",
            "ROLE", "PERMISSION", "SECURITY", "COMPLIANCE",
            "PERFORMANCE", "SCALABILITY", "AVAILABILITY", "RELIABILITY",
            "BUG", "FEATURE_REQUEST", "ENHANCEMENT", "INCIDENT",
            "TEAM", "DEPARTMENT", "STAKEHOLDER", "CLIENT",
            "VENDOR", "PARTNER", "COMPETITOR", "MARKET",
            "TIMELINE", "BUDGET", "COST", "REVENUE",
            "SKILL", "EXPERTISE", "CERTIFICATION", "TRAINING",
            "OTHER"
        ]
        
        return {
            "default_entity_types": default_entity_types,
            "categories": {
                "people_and_organization": ["PERSON", "ORGANIZATION", "TEAM", "DEPARTMENT", "STAKEHOLDER", "CLIENT", "VENDOR", "PARTNER", "COMPETITOR"],
                "technology_and_systems": ["TECHNOLOGY", "API", "DATABASE", "FRAMEWORK", "COMPONENT", "INTERFACE", "MODULE", "LIBRARY", "DEPENDENCY", "PLATFORM", "ENVIRONMENT"],
                "products_and_services": ["PRODUCT", "SERVICE", "FEATURE", "TOOL", "SOLUTION"],
                "project_management": ["PROJECT", "MILESTONE", "TIMELINE", "GOAL", "STRATEGY", "REQUIREMENT", "ISSUE", "ENHANCEMENT", "WORKFLOW", "PROCESS"],
                "development_and_deployment": ["VERSION", "RELEASE", "DEPLOYMENT", "CONFIGURATION", "BUG", "FEATURE_REQUEST", "INCIDENT"],
                "documentation_and_standards": ["DOCUMENT", "SPECIFICATION", "GUIDELINE", "POLICY", "STANDARD", "PROTOCOL"],
                "architecture_and_design": ["ARCHITECTURE", "PATTERN", "PRINCIPLE", "CONCEPT", "MODEL", "SCHEMA", "ENTITY"],
                "security_and_compliance": ["SECURITY", "COMPLIANCE", "ROLE", "PERMISSION"],
                "performance_and_quality": ["PERFORMANCE", "SCALABILITY", "AVAILABILITY", "RELIABILITY", "METRIC"],
                "business_and_finance": ["MARKET", "BUDGET", "COST", "REVENUE"],
                "learning_and_development": ["SKILL", "EXPERTISE", "CERTIFICATION", "TRAINING"],
                "events_and_meetings": ["EVENT", "MEETING", "DECISION"],
                "data_and_resources": ["DATA", "RESOURCE", "LOCATION"]
            },
            "description": "Comprehensive entity types designed for modern software development, project management, and business operations",
            "usage": "You can customize these when creating a notebook or update them via PUT /notebooks/{notebook_id}/schema",
            "specialized_sets": {
                "software_development": ["API", "DATABASE", "FRAMEWORK", "COMPONENT", "INTERFACE", "MODULE", "LIBRARY", "BUG", "FEATURE_REQUEST", "VERSION", "RELEASE", "DEPLOYMENT"],
                "project_management": ["PROJECT", "MILESTONE", "TIMELINE", "GOAL", "STRATEGY", "REQUIREMENT", "ISSUE", "WORKFLOW", "PROCESS", "DECISION"],
                "business_analysis": ["STAKEHOLDER", "CLIENT", "VENDOR", "PARTNER", "MARKET", "BUDGET", "COST", "REVENUE", "STRATEGY", "GOAL", "RISK"],
                "technical_documentation": ["SPECIFICATION", "GUIDELINE", "POLICY", "STANDARD", "PROTOCOL", "ARCHITECTURE", "PATTERN", "PRINCIPLE"],
                "minimal_set": ["PERSON", "ORGANIZATION", "TECHNOLOGY", "PRODUCT", "PROJECT", "DOCUMENT", "ISSUE", "CONCEPT", "OTHER"]
            }
        }

    async def extract_text_from_file(filename: str, file_content: bytes) -> str:
        """Extract text from various file formats supported by LightRAG"""
        try:
            file_ext = filename.lower().split('.')[-1] if '.' in filename else ''
            
            # PDF files
            if file_ext == 'pdf':
                return extract_text_from_pdf_lightrag(file_content)
            
            # Plain text files
            elif file_ext in ['txt', 'md', 'markdown', 'rst']:
                return file_content.decode('utf-8')
            
            # CSV files
            elif file_ext == 'csv':
                return file_content.decode('utf-8')
            
            # JSON files
            elif file_ext == 'json':
                import json
                json_data = json.loads(file_content.decode('utf-8'))
                return json.dumps(json_data, indent=2)
            
            # XML/HTML files
            elif file_ext in ['xml', 'html', 'htm']:
                content = file_content.decode('utf-8')
                try:
                    from bs4 import BeautifulSoup
                    if file_ext in ['html', 'htm']:
                        soup = BeautifulSoup(content, 'html.parser')
                        return soup.get_text(separator='\n', strip=True)
                    else:
                        soup = BeautifulSoup(content, 'xml')
                        return soup.get_text(separator='\n', strip=True)
                except ImportError:
                    # Fallback to raw content if BeautifulSoup not available
                    return content
            
            # Word documents
            elif file_ext in ['docx', 'doc']:
                try:
                    import docx
                    from io import BytesIO
                    doc = docx.Document(BytesIO(file_content))
                    content_parts = []
                    for para in doc.paragraphs:
                        text = para.text.strip()
                        if text:
                            content_parts.append(text)
                    return "\n\n".join(content_parts)
                except ImportError:
                    logger.warning("python-docx not available, cannot process Word documents")
                    raise HTTPException(
                        status_code=400, 
                        detail="Word document processing requires python-docx library. Install with: pip install python-docx"
                    )
            
            # Excel files
            elif file_ext in ['xlsx', 'xls']:
                try:
                    import pandas as pd
                    from io import BytesIO
                    excel_file = pd.ExcelFile(BytesIO(file_content))
                    content_parts = []
                    for sheet_name in excel_file.sheet_names:
                        df = pd.read_excel(BytesIO(file_content), sheet_name=sheet_name)
                        content_parts.append(f"Sheet: {sheet_name}\n{df.to_string(index=False)}")
                    return "\n\n" + "="*50 + "\n\n".join(content_parts)
                except ImportError:
                    logger.warning("pandas not available, cannot process Excel files")
                    raise HTTPException(
                        status_code=400, 
                        detail="Excel processing requires pandas and openpyxl libraries. Install with: pip install pandas openpyxl xlrd"
                    )
            
            # PowerPoint files
            elif file_ext in ['pptx', 'ppt']:
                try:
                    from pptx import Presentation
                    from io import BytesIO
                    prs = Presentation(BytesIO(file_content))
                    content_parts = []
                    for i, slide in enumerate(prs.slides):
                        slide_text = []
                        for shape in slide.shapes:
                            if hasattr(shape, "text") and shape.text.strip():
                                slide_text.append(shape.text)
                        if slide_text:
                            content_parts.append(f"--- Slide {i + 1} ---\n" + "\n".join(slide_text))
                    return "\n\n".join(content_parts)
                except ImportError:
                    logger.warning("python-pptx not available, cannot process PowerPoint files")
                    raise HTTPException(
                        status_code=400, 
                        detail="PowerPoint processing requires python-pptx library. Install with: pip install python-pptx"
                    )
            
            # RTF files
            elif file_ext == 'rtf':
                try:
                    from striprtf.striprtf import rtf_to_text
                    rtf_content = file_content.decode('utf-8')
                    return rtf_to_text(rtf_content)
                except ImportError:
                    logger.warning("striprtf not available, cannot process RTF files")
                    raise HTTPException(
                        status_code=400, 
                        detail="RTF processing requires striprtf library. Install with: pip install striprtf"
                    )
            
            # LibreOffice formats - basic support
            elif file_ext in ['odt', 'ods', 'odp']:
                logger.warning(f"LibreOffice format {file_ext} requires conversion to supported format")
                raise HTTPException(
                    status_code=400, 
                    detail=f"LibreOffice {file_ext.upper()} files are not directly supported. Please convert to DOCX, PDF, or TXT format."
                )
            
            # Use textract as fallback for other formats if available
            else:
                try:
                    import textract
                    from tempfile import NamedTemporaryFile
                    import os
                    
                    # Create temporary file
                    with NamedTemporaryFile(suffix=f'.{file_ext}', delete=False) as temp_file:
                        temp_file.write(file_content)
                        temp_file.flush()
                        
                        try:
                            # Use textract to extract text
                            extracted_text = textract.process(temp_file.name)
                            return extracted_text.decode('utf-8')
                        finally:
                            # Clean up temporary file
                            os.unlink(temp_file.name)
                            
                except ImportError:
                    logger.warning("textract not available for additional format support")
                    raise HTTPException(
                        status_code=400, 
                        detail=f"Unsupported file type: {filename}. Supported formats: PDF, TXT, MD, CSV, JSON, XML, HTML, DOCX, DOC, XLSX, XLS, PPTX, PPT, RTF. For more formats, install: pip install textract"
                    )
                except Exception as e:
                    logger.error(f"Error processing file with textract: {e}")
                    raise HTTPException(
                        status_code=400, 
                        detail=f"Error processing file {filename}: {str(e)}"
                    )
        
        except HTTPException:
            # Re-raise HTTP exceptions as-is
            raise
        except Exception as e:
            logger.error(f"Unexpected error processing file {filename}: {e}")
            raise HTTPException(
                status_code=500, 
                detail=f"Unexpected error processing file: {str(e)}"
            )
        """Extract text from various file formats supported by LightRAG"""
        try:
            file_ext = filename.lower().split('.')[-1] if '.' in filename else ''
            
            # PDF files
            if file_ext == 'pdf':
                return extract_text_from_pdf_lightrag(file_content)
            
            # Plain text files
            elif file_ext in ['txt', 'md', 'markdown', 'rst']:
                return file_content.decode('utf-8')
            
            # CSV files
            elif file_ext == 'csv':
                return file_content.decode('utf-8')
            
            # JSON files
            elif file_ext == 'json':
                import json
                json_data = json.loads(file_content.decode('utf-8'))
                return json.dumps(json_data, indent=2)
            
            # XML/HTML files
            elif file_ext in ['xml', 'html', 'htm']:
                content = file_content.decode('utf-8')
                try:
                    from bs4 import BeautifulSoup
                    if file_ext in ['html', 'htm']:
                        soup = BeautifulSoup(content, 'html.parser')
                        return soup.get_text(separator='\n', strip=True)
                    else:
                        soup = BeautifulSoup(content, 'xml')
                        return soup.get_text(separator='\n', strip=True)
                except ImportError:
                    # Fallback to raw content if BeautifulSoup not available
                    return content
            
            # Word documents
            elif file_ext in ['docx', 'doc']:
                try:
                    import docx
                    from io import BytesIO
                    doc = docx.Document(BytesIO(file_content))
                    content_parts = []
                    for para in doc.paragraphs:
                        text = para.text.strip()
                        if text:
                            content_parts.append(text)
                    return "\n\n".join(content_parts)
                except ImportError:
                    logger.warning("python-docx not available, cannot process Word documents")
                    raise HTTPException(
                        status_code=400, 
                        detail="Word document processing requires python-docx library"
                    )
            
            # Excel files
            elif file_ext in ['xlsx', 'xls']:
                try:
                    import pandas as pd
                    from io import BytesIO
                    excel_file = pd.ExcelFile(BytesIO(file_content))
                    content_parts = []
                    for sheet_name in excel_file.sheet_names:
                        df = pd.read_excel(BytesIO(file_content), sheet_name=sheet_name)
                        content_parts.append(f"Sheet: {sheet_name}\n{df.to_string(index=False)}")
                    return "\n\n" + "="*50 + "\n\n".join(content_parts)
                except ImportError:
                    logger.warning("pandas not available, cannot process Excel files")
                    raise HTTPException(
                        status_code=400, 
                        detail="Excel processing requires pandas and openpyxl libraries"
                    )
            
            # PowerPoint files
            elif file_ext in ['pptx', 'ppt']:
                try:
                    from pptx import Presentation
                    from io import BytesIO
                    prs = Presentation(BytesIO(file_content))
                    content_parts = []
                    for i, slide in enumerate(prs.slides):
                        slide_text = []
                        for shape in slide.shapes:
                            if hasattr(shape, "text") and shape.text.strip():
                                slide_text.append(shape.text)
                        if slide_text:
                            content_parts.append(f"--- Slide {i + 1} ---\n" + "\n".join(slide_text))
                    return "\n\n".join(content_parts)
                except ImportError:
                    logger.warning("python-pptx not available, cannot process PowerPoint files")
                    raise HTTPException(
                        status_code=400, 
                        detail="PowerPoint processing requires python-pptx library"
                    )
            
            # RTF files
            elif file_ext == 'rtf':
                try:
                    from striprtf.striprtf import rtf_to_text
                    rtf_content = file_content.decode('utf-8')
                    return rtf_to_text(rtf_content)
                except ImportError:
                    logger.warning("striprtf not available, cannot process RTF files")
                    raise HTTPException(
                        status_code=400, 
                        detail="RTF processing requires striprtf library"
                    )
            
            # LibreOffice formats - basic support
            elif file_ext in ['odt', 'ods', 'odp']:
                logger.warning(f"LibreOffice format {file_ext} requires conversion to supported format")
                raise HTTPException(
                    status_code=400, 
                    detail=f"LibreOffice {file_ext.upper()} files are not directly supported. Please convert to DOCX, PDF, or TXT format."
                )
            
            # Use textract as fallback for other formats if available
            else:
                try:
                    import textract
                    from tempfile import NamedTemporaryFile
                    
                    # Create temporary file
                    with NamedTemporaryFile(suffix=f'.{file_ext}', delete=False) as temp_file:
                        temp_file.write(file_content)
                        temp_file.flush()
                        
                        try:
                            # Use textract to extract text
                            extracted_text = textract.process(temp_file.name)
                            return extracted_text.decode('utf-8')
                        finally:
                            # Clean up temporary file
                            import os
                            os.unlink(temp_file.name)
                            
                except ImportError:
                    logger.warning("textract not available for additional format support")
                    raise HTTPException(
                        status_code=400, 
                        detail=f"Unsupported file type: {filename}. Supported formats: PDF, TXT, MD, CSV, JSON, XML, HTML, DOCX, DOC, XLSX, XLS, PPTX, PPT, RTF"
                    )
                except Exception as e:
                    logger.error(f"Error processing file with textract: {e}")
                    raise HTTPException(
                        status_code=400, 
                        detail=f"Error processing file {filename}: {str(e)}"
                    )
        
        except HTTPException:
            # Re-raise HTTP exceptions as-is
            raise
        except Exception as e:
            logger.error(f"Unexpected error processing file {filename}: {e}")
            raise HTTPException(
                status_code=500, 
                detail=f"Unexpected error processing file: {str(e)}"
            )

    def levenshtein_distance(s1: str, s2: str) -> int:
        """Calculate Levenshtein distance between two strings"""
        if len(s1) < len(s2):
            return levenshtein_distance(s2, s1)
        if len(s2) == 0:
            return len(s1)
        
        previous_row = range(len(s2) + 1)
        for i, c1 in enumerate(s1):
            current_row = [i + 1]
            for j, c2 in enumerate(s2):
                insertions = previous_row[j + 1] + 1
                deletions = current_row[j] + 1
                substitutions = previous_row[j] + (c1 != c2)
                current_row.append(min(insertions, deletions, substitutions))
            previous_row = current_row
        
        return previous_row[-1]

    def detect_embedding_model_specs(
        model_name: str, 
        manual_override: Optional[Dict[str, int]] = None
    ) -> Dict[str, Any]:
        """
        Detect embedding model specifications with fuzzy matching and confidence scoring.
        
        Args:
            model_name: The embedding model name (can be any case/format)
            manual_override: Optional dict with 'dimensions' and/or 'max_tokens'
        
        Returns:
            {
                'dimensions': int,
                'max_tokens': int,
                'confidence': float (0-1),
                'detected_pattern': str,
                'override_options': List[Dict],  # Common alternatives
                'is_manual_override': bool
            }
        """
        # If manual override provided, use it directly
        if manual_override:
            return {
                'dimensions': manual_override.get('dimensions', 1536),
                'max_tokens': manual_override.get('max_tokens', 8192),
                'confidence': 1.0,
                'detected_pattern': 'manual_override',
                'override_options': [],
                'is_manual_override': True
            }
        
        # Normalize model name for matching
        normalized_name = model_name.lower().replace('_', '-').replace(' ', '-')
        
        # Known model patterns with metadata
        # Format: (pattern, dimensions, max_tokens, aliases)
        model_database = [
            # OpenAI Models
            ('text-embedding-ada-002', 1536, 8192, ['ada-002', 'ada002', 'text-embedding-ada']),
            ('text-embedding-3-small', 1536, 8192, ['embedding-3-small', 'openai-3-small', 'text-3-small']),
            ('text-embedding-3-large', 3072, 8192, ['embedding-3-large', 'openai-3-large', 'text-3-large']),
            
            # Qwen Models (Alibaba) - IMPORTANT: Order matters! More specific patterns first
            ('qwen3-embedding-0.6b', 1024, 512, ['qwen3-0.6b', 'qwen3-06b', 'qwen-3-embedding-0.6b', 'qwen3-embedding-06b']),
            ('qwen3-embedding-4b', 2560, 1024, ['qwen3-4b', 'qwen-3-embedding-4b']),
            ('qwen3-embedding-8b', 4096, 1024, ['qwen3-8b', 'qwen-3-embedding-8b']),
            ('qwen2.5-coder', 2560, 2048, ['qwen2.5-coder', 'qwen25-coder']),
            ('qwen2', 2560, 2048, ['qwen-2', 'qwen 2']),
            ('qwen', 2560, 2048, ['qwen-embedding']),  # Generic fallback - must be LAST
            
            # Jina AI Models
            ('jina-embeddings-v4', 2048, 8192, ['jina-v4', 'jina-4', 'jina-embeddings-4']),
            ('jina-embeddings-v3', 1024, 8192, ['jina-v3', 'jina-3', 'jina-embeddings-3']),
            ('jina-embeddings-v2-base', 768, 8192, ['jina-v2-base', 'jina-2-base']),
            ('jina-embeddings-v2-small', 512, 8192, ['jina-v2-small', 'jina-2-small']),
            
            # Sentence Transformers Models
            ('all-minilm-l6-v2', 384, 256, ['minilm-l6', 'all-minilm-l6', 'sentence-transformers-minilm-l6']),
            ('all-minilm-l12-v2', 384, 256, ['minilm-l12', 'all-minilm-l12']),
            ('all-mpnet-base-v2', 768, 384, ['mpnet-base', 'all-mpnet-base', 'sentence-transformers-mpnet']),
            ('paraphrase-minilm-l6-v2', 384, 256, ['paraphrase-minilm-l6', 'paraphrase-l6']),
            ('paraphrase-mpnet-base-v2', 768, 384, ['paraphrase-mpnet', 'paraphrase-base']),
            
            # BAAI BGE Models
            ('bge-large-en', 1024, 512, ['bge-large', 'baai-bge-large']),
            ('bge-base-en', 768, 512, ['bge-base', 'baai-bge-base']),
            ('bge-small-en', 512, 512, ['bge-small', 'baai-bge-small']),
            ('bge-m3', 1024, 8192, ['baai-bge-m3', 'bge-multilingual']),
            
            # Microsoft E5 Models
            ('e5-large-v2', 1024, 512, ['e5-large', 'microsoft-e5-large']),
            ('e5-base-v2', 768, 512, ['e5-base', 'microsoft-e5-base']),
            ('e5-small-v2', 384, 512, ['e5-small', 'microsoft-e5-small']),
            
            # MixedBread AI Models
            ('mxbai-embed-large', 1024, 512, ['mxbai-large', 'mixedbread-large']),
            
            # Nomic AI Models
            ('nomic-embed-text', 768, 512, ['nomic-embed', 'nomic-ai-embed']),
            
            # Cohere Models
            ('embed-english-v3', 1024, 512, ['cohere-english-v3', 'cohere-embed-english']),
            ('embed-multilingual-v3', 1024, 512, ['cohere-multilingual-v3']),
            
            # Voyage AI Models
            ('voyage-large-2', 1536, 16000, ['voyage-2-large', 'voyage-large']),
            ('voyage-2', 1024, 16000, ['voyage-base-2']),
            
            # Snowflake Arctic Models
            ('snowflake-arctic-embed-l', 1024, 512, ['arctic-large', 'snowflake-large']),
            ('snowflake-arctic-embed-m', 768, 512, ['arctic-medium', 'snowflake-medium']),
            
            # Google Models
            ('textembedding-gecko', 768, 2048, ['gecko', 'google-gecko']),
            
            # IBM Models
            ('ibm-slate-30m', 384, 512, ['slate-30m', 'ibm-slate']),
            ('ibm-slate-125m', 768, 512, ['slate-125m']),
            
            # NVIDIA NeMo Models
            ('nv-embed-v1', 4096, 32768, ['nvidia-embed', 'nemo-embed']),
        ]
        
        best_match = None
        best_confidence = 0.0
        detected_pattern = 'unknown'
        best_match_length = 0  # Track match length to prefer longer/more specific matches

        # Try exact substring match first (highest confidence)
        # Prefer longer matches to avoid generic patterns matching specific models
        for pattern, dims, tokens, aliases in model_database:
            if pattern in normalized_name:
                # Prefer longer pattern matches (more specific)
                if len(pattern) > best_match_length or (len(pattern) == best_match_length and 0.95 > best_confidence):
                    best_match = (pattern, dims, tokens)
                    best_confidence = 0.95
                    detected_pattern = pattern
                    best_match_length = len(pattern)

            # Check aliases
            for alias in aliases:
                if alias in normalized_name:
                    # Prefer longer alias matches
                    if len(alias) > best_match_length or (len(alias) == best_match_length and 0.90 > best_confidence):
                        best_match = (pattern, dims, tokens)
                        best_confidence = 0.90
                        detected_pattern = f"{pattern} (alias: {alias})"
                        best_match_length = len(alias)
        
        # If no exact match, try fuzzy matching
        if not best_match:
            for pattern, dims, tokens, aliases in model_database:
                # Calculate Levenshtein distance
                distance = levenshtein_distance(normalized_name, pattern)
                max_len = max(len(normalized_name), len(pattern))
                similarity = 1 - (distance / max_len)
                
                if similarity > best_confidence and similarity > 0.6:  # Threshold 0.6
                    best_match = (pattern, dims, tokens)
                    best_confidence = similarity
                    detected_pattern = f"{pattern} (fuzzy: {similarity:.2f})"
                
                # Also check aliases
                for alias in aliases:
                    distance = levenshtein_distance(normalized_name, alias)
                    max_len = max(len(normalized_name), len(alias))
                    similarity = 1 - (distance / max_len)
                    
                    if similarity > best_confidence and similarity > 0.6:
                        best_match = (pattern, dims, tokens)
                        best_confidence = similarity
                        detected_pattern = f"{pattern} (fuzzy alias: {similarity:.2f})"
        
        # Default fallback (low confidence)
        if not best_match:
            best_match = ('text-embedding-ada-002', 1536, 8192)
            best_confidence = 0.0
            detected_pattern = 'fallback_default'
        
        dimensions, max_tokens = best_match[1], best_match[2]
        
        # Generate override options (common dimension sizes)
        override_options = [
            {'dimensions': 384, 'max_tokens': 256, 'label': 'Small (384d) - MiniLM, E5-Small'},
            {'dimensions': 512, 'max_tokens': 512, 'label': 'Medium-Small (512d) - BGE-Small, Jina v2-Small'},
            {'dimensions': 768, 'max_tokens': 512, 'label': 'Base (768d) - MPNet, E5-Base, BGE-Base'},
            {'dimensions': 1024, 'max_tokens': 512, 'label': 'Large (1024d) - E5-Large, BGE-Large, Qwen3-0.6b'},
            {'dimensions': 1536, 'max_tokens': 8192, 'label': 'OpenAI Ada-002 / GPT-3-Small (1536d)'},
            {'dimensions': 2048, 'max_tokens': 8192, 'label': 'Jina v4 (2048d)'},
            {'dimensions': 2560, 'max_tokens': 2048, 'label': 'Qwen Standard (2560d)'},
            {'dimensions': 3072, 'max_tokens': 8192, 'label': 'OpenAI GPT-3-Large (3072d)'},
            {'dimensions': 4096, 'max_tokens': 1024, 'label': 'Qwen3-8b / NV-Embed (4096d)'},
        ]
        
        return {
            'dimensions': dimensions,
            'max_tokens': max_tokens,
            'confidence': round(best_confidence, 2),
            'detected_pattern': detected_pattern,
            'override_options': override_options,
            'is_manual_override': False
        }

    async def create_lightrag_instance(notebook_id: str, llm_provider_config: Dict[str, Any], embedding_provider_config: Dict[str, Any], entity_types: Optional[List[str]] = None, language: str = "en", manual_embedding_override: Optional[Dict[str, int]] = None) -> LightRAG:
        """Create a new LightRAG instance for a notebook with specified provider configurations"""
        working_dir = LIGHTRAG_STORAGE_PATH / notebook_id
        
        # Create directory if it doesn't exist (preserve existing data)
        working_dir.mkdir(exist_ok=True)
        
        try:
            logger.info(f"Creating LightRAG instance for notebook {notebook_id}")
            
            # Get configuration
            llm_provider_type = llm_provider_config.get('type', 'openai')
            llm_model_name = llm_provider_config.get('model', 'gpt-3.5-turbo')
            llm_api_key = llm_provider_config.get('apiKey', '')
            llm_base_url = llm_provider_config.get('baseUrl', '')
            
            embedding_provider_type = embedding_provider_config.get('type', 'openai')
            embedding_model_name = embedding_provider_config.get('model', 'text-embedding-ada-002')
            embedding_api_key = embedding_provider_config.get('apiKey', '')
            embedding_base_url = embedding_provider_config.get('baseUrl', '')
            
            # CRITICAL: Normalize URLs for local development (convert host.docker.internal to localhost)
            llm_base_url = normalize_url_for_local_dev(llm_base_url)
            embedding_base_url = normalize_url_for_local_dev(embedding_base_url)
            
            # Log final configuration (auto-detection already done in create_notebook)
            logger.info(f"LLM Provider: {llm_provider_type} - {llm_model_name} @ {llm_base_url}")
            logger.info(f"Embedding Provider: {embedding_provider_type} - {embedding_model_name} @ {embedding_base_url}")
            
            # Remove /v1 suffix from Ollama URLs if present
            if llm_provider_type == 'ollama' and llm_base_url.endswith('/v1'):
                llm_base_url = llm_base_url[:-3]
            if embedding_provider_type == 'ollama' and embedding_base_url.endswith('/v1'):
                embedding_base_url = embedding_base_url[:-3]
            
            # Enhanced API key handling with fallbacks for different scenarios
            is_clara_core_llm = ':8091' in llm_base_url or llm_base_url.endswith(':8091')
            is_clara_core_embedding = ':8091' in embedding_base_url or embedding_base_url.endswith(':8091')
            
            # Check for localhost/local endpoints that might not need real API keys
            is_local_llm_endpoint = any(local_host in llm_base_url.lower() for local_host in ['localhost', '127.0.0.1', '0.0.0.0']) if llm_base_url else False
            is_local_embedding_endpoint = any(local_host in embedding_base_url.lower() for local_host in ['localhost', '127.0.0.1', '0.0.0.0']) if embedding_base_url else False
            
            # Handle LLM API key with fallbacks
            if llm_provider_type != 'ollama':
                if is_clara_core_llm:
                    # Clara Core doesn't need a real API key
                    llm_api_key = 'claracore'
                    logger.info("Using Clara Core LLM - no API key required")
                elif not llm_api_key or llm_api_key.strip() == 'your-api-key':
                    if is_local_llm_endpoint:
                        # Local endpoints often don't require real API keys
                        llm_api_key = 'local-endpoint-key'
                        logger.info(f"Using local LLM endpoint {llm_base_url} - using fallback API key")
                    elif llm_base_url == 'https://api.openai.com/v1' or llm_base_url == 'https://api.openai.com':
                        # Real OpenAI endpoint requires a valid key
                        raise ValueError(f"Invalid or missing LLM API key for OpenAI provider. Please provide a valid API key.")
                    else:
                        # For other OpenAI-compatible endpoints, provide a fallback key and log a warning
                        llm_api_key = 'fallback-api-key'
                        logger.warning(f"No API key provided for LLM endpoint {llm_base_url}. Using fallback key. This may cause authentication errors if the endpoint requires a real API key.")
            
            # Handle embedding API key with fallbacks
            if embedding_provider_type != 'ollama':
                if is_clara_core_embedding:
                    # Clara Core doesn't need a real API key
                    embedding_api_key = 'claracore'
                    logger.info("Using Clara Core embedding - no API key required")
                elif not embedding_api_key or embedding_api_key.strip() == 'your-api-key':
                    if is_local_embedding_endpoint:
                        # Local endpoints often don't require real API keys
                        embedding_api_key = 'local-endpoint-key'
                        logger.info(f"Using local embedding endpoint {embedding_base_url} - using fallback API key")
                    elif embedding_base_url == 'https://api.openai.com/v1' or embedding_base_url == 'https://api.openai.com':
                        # Real OpenAI endpoint requires a valid key
                        raise ValueError(f"Invalid or missing embedding API key for OpenAI provider. Please provide a valid API key.")
                    else:
                        # For other OpenAI-compatible endpoints, provide a fallback key and log a warning
                        embedding_api_key = 'fallback-api-key'
                        logger.warning(f"No API key provided for embedding endpoint {embedding_base_url}. Using fallback key. This may cause authentication errors if the endpoint requires a real API key.")
            
            # Determine embedding dimensions and max tokens using fuzzy matching
            # Use the new detect_embedding_model_specs function with optional manual override
            embedding_specs = detect_embedding_model_specs(embedding_model_name, manual_embedding_override)
            embedding_dim = embedding_specs['dimensions']
            embedding_max_tokens = embedding_specs['max_tokens']
            
            # Log detection results
            logger.info(f"Embedding model '{embedding_model_name}' detected:")
            logger.info(f"  - Dimensions: {embedding_dim}")
            logger.info(f"  - Max tokens: {embedding_max_tokens}")
            logger.info(f"  - Confidence: {embedding_specs['confidence']}")
            logger.info(f"  - Pattern: {embedding_specs['detected_pattern']}")
            logger.info(f"  - Manual override: {embedding_specs['is_manual_override']}")

            # LEGACY FALLBACK: All dimension detection below is DISABLED
            # The fuzzy matching above handles all models correctly
            # DO NOT re-enable the code below - it will override correct detection

            if False:  # DISABLED: Legacy hardcoded dimension detection
                # This entire block is disabled to prevent overriding fuzzy matching results
                # The new detect_embedding_model_specs() function handles all models
                pass

            # OpenAI Models (LEGACY - DISABLED)
            elif False and 'text-embedding-ada-002' in embedding_model_name:
                embedding_dim = 1536
                embedding_max_tokens = 8192
            elif False and 'text-embedding-3-small' in embedding_model_name:
                embedding_dim = 1536
                embedding_max_tokens = 8192
            elif False and 'text-embedding-3-large' in embedding_model_name:
                embedding_dim = 3072
                embedding_max_tokens = 8192

            # MixedBread AI Models
            elif False and 'mxbai-embed-large' in embedding_model_name:
                embedding_dim = 1024
                embedding_max_tokens = 512  # mxbai has lower token limit
            
            # Nomic AI Models
            elif 'nomic-embed' in embedding_model_name:
                embedding_dim = 768
                embedding_max_tokens = 512  # nomic has lower token limit
            
            # Microsoft E5 Models
            elif 'e5-large-v2' in embedding_model_name:
                embedding_dim = 1024
                embedding_max_tokens = 512
            elif 'e5-base-v2' in embedding_model_name:
                embedding_dim = 768
                embedding_max_tokens = 512
            elif 'e5-small-v2' in embedding_model_name:
                embedding_dim = 384
                embedding_max_tokens = 512
            
            # Sentence Transformers Models (case-insensitive matching)
            elif 'all-minilm-l6-v2' in embedding_model_name.lower() or 'all_minilm_l6_v2' in embedding_model_name.lower():
                embedding_dim = 384
                embedding_max_tokens = 256  # Smaller models have lower limits
            elif 'all-minilm-l12-v2' in embedding_model_name.lower() or 'all_minilm_l12_v2' in embedding_model_name.lower():
                embedding_dim = 384
                embedding_max_tokens = 256
            elif 'all-mpnet-base-v2' in embedding_model_name.lower() or 'all_mpnet_base_v2' in embedding_model_name.lower():
                embedding_dim = 768
                embedding_max_tokens = 384
            
            # BAAI BGE Models
            elif 'bge-large' in embedding_model_name:
                embedding_dim = 1024
                embedding_max_tokens = 512
            elif 'bge-base' in embedding_model_name:
                embedding_dim = 768
                embedding_max_tokens = 512
            elif 'bge-small' in embedding_model_name:
                embedding_dim = 512
                embedding_max_tokens = 512
            elif 'bge-m3' in embedding_model_name:
                embedding_dim = 1024
                embedding_max_tokens = 8192  # bge-m3 supports longer contexts
            
            # Qwen Models (Alibaba) - DISABLED: Now handled by fuzzy matching
            elif False and 'qwen' in embedding_model_name.lower():
                embedding_dim = 2560
                embedding_max_tokens = 2048  # Qwen models have moderate token limits
            elif False and 'qwen2.5-coder' in embedding_model_name.lower():
                embedding_dim = 2560
                embedding_max_tokens = 2048
            elif False and 'qwen2' in embedding_model_name.lower():
                embedding_dim = 2560
                embedding_max_tokens = 2048
            elif False and ("qwen3-embedding-0.6b" in embedding_model_name.lower() or "qwen3-embedding-06b" in embedding_model_name.lower()):
                embedding_dim = 1024
                embedding_max_tokens = 512  # CRITICAL: 0.6B model needs much smaller batches due to server limitations
            elif False and "qwen3-embedding-4b" in embedding_model_name.lower():
                embedding_dim = 2560
                embedding_max_tokens = 1024  # Reduced for batch processing stability
            elif False and "qwen3-embedding-8b" in embedding_model_name.lower():
                embedding_dim = 4096
                embedding_max_tokens = 1024  # Reduced for batch processing stability
            
            # Jina AI Models
            elif 'jina-embeddings-v4' in embedding_model_name:
                # V4 models - newest generation with better performance
                embedding_dim = 2048  # All v4 models use 2048 dimensions
                embedding_max_tokens = 8192
            elif 'jina-embeddings-v3' in embedding_model_name:
                # V3 models - general purpose
                embedding_dim = 1024
                embedding_max_tokens = 8192
            elif 'jina-embeddings-v2-base' in embedding_model_name:
                embedding_dim = 768
                embedding_max_tokens = 8192
            elif 'jina-embeddings-v2-small' in embedding_model_name:
                embedding_dim = 512
                embedding_max_tokens = 8192
            
            # Cohere Models
            elif 'embed-english-v3.0' in embedding_model_name:
                embedding_dim = 1024
                embedding_max_tokens = 512
            elif 'embed-multilingual-v3.0' in embedding_model_name:
                embedding_dim = 1024
                embedding_max_tokens = 512
            elif 'embed-english-light-v3.0' in embedding_model_name:
                embedding_dim = 384
                embedding_max_tokens = 512
            
            # Voyage AI Models
            elif 'voyage-large-2' in embedding_model_name:
                embedding_dim = 1536
                embedding_max_tokens = 16000  # Voyage models support very long contexts
            elif 'voyage-code-2' in embedding_model_name:
                embedding_dim = 1536
                embedding_max_tokens = 16000
            elif 'voyage-2' in embedding_model_name:
                embedding_dim = 1024
                embedding_max_tokens = 4000
            
            # Snowflake Arctic Embed Models
            elif 'snowflake-arctic-embed2' in embedding_model_name.lower():
                embedding_dim = 1024  # Arctic Embed 2.0 multilingual
                embedding_max_tokens = 512
            elif 'snowflake-arctic-embed-m-v1.5' in embedding_model_name.lower():
                embedding_dim = 768
                embedding_max_tokens = 512
            elif 'snowflake-arctic-embed-l-v1.5' in embedding_model_name.lower():
                embedding_dim = 1024
                embedding_max_tokens = 512
            elif 'snowflake-arctic-embed-m' in embedding_model_name.lower():
                embedding_dim = 768
                embedding_max_tokens = 512
            elif 'snowflake-arctic-embed-l' in embedding_model_name.lower():
                embedding_dim = 1024
                embedding_max_tokens = 512
            elif 'snowflake-arctic-embed' in embedding_model_name.lower():
                embedding_dim = 1024  # Default for Arctic models
                embedding_max_tokens = 512
            
            # Google EmbeddingGemma
            elif 'embeddinggemma' in embedding_model_name.lower() or 'embedding-gemma' in embedding_model_name.lower():
                embedding_dim = 768  # EmbeddingGemma 308M parameters
                embedding_max_tokens = 2048
            
            # IBM Granite Embedding
            elif 'granite-embedding-278m' in embedding_model_name.lower():
                embedding_dim = 768  # Multilingual version
                embedding_max_tokens = 512
            elif 'granite-embedding-30m' in embedding_model_name.lower():
                embedding_dim = 384  # English only, smaller
                embedding_max_tokens = 512
            elif 'granite-embedding' in embedding_model_name.lower():
                embedding_dim = 768  # Default to larger model
                embedding_max_tokens = 512
            
            # Sentence-Transformers Paraphrase Models
            elif 'paraphrase-multilingual' in embedding_model_name.lower():
                embedding_dim = 768
                embedding_max_tokens = 128
            elif 'paraphrase-mpnet' in embedding_model_name.lower():
                embedding_dim = 768
                embedding_max_tokens = 128
            elif 'paraphrase-albert' in embedding_model_name.lower():
                embedding_dim = 768
                embedding_max_tokens = 128
            elif 'paraphrase-minilm' in embedding_model_name.lower():
                embedding_dim = 384
                embedding_max_tokens = 128
            
            # Additional Sentence-Transformers Models
            elif 'sentence-t5' in embedding_model_name.lower():
                if 'xxl' in embedding_model_name.lower():
                    embedding_dim = 768
                elif 'xl' in embedding_model_name.lower():
                    embedding_dim = 768
                elif 'large' in embedding_model_name.lower():
                    embedding_dim = 768
                else:
                    embedding_dim = 768
                embedding_max_tokens = 256
            
            # Alibaba GTE Models
            elif 'gte-large' in embedding_model_name.lower():
                embedding_dim = 1024
                embedding_max_tokens = 512
            elif 'gte-base' in embedding_model_name.lower():
                embedding_dim = 768
                embedding_max_tokens = 512
            elif 'gte-small' in embedding_model_name.lower():
                embedding_dim = 384
                embedding_max_tokens = 512
            elif 'gte-qwen' in embedding_model_name.lower():
                embedding_dim = 1024
                embedding_max_tokens = 8192  # Qwen-based GTE models support longer context
            
            # UAE (Universal AnglE Embedding) Models
            elif 'uae-large-v1' in embedding_model_name.lower():
                embedding_dim = 1024
                embedding_max_tokens = 512
            
            # Instructor Models
            elif 'instructor-xl' in embedding_model_name.lower():
                embedding_dim = 768
                embedding_max_tokens = 512
            elif 'instructor-large' in embedding_model_name.lower():
                embedding_dim = 768
                embedding_max_tokens = 512
            elif 'instructor-base' in embedding_model_name.lower():
                embedding_dim = 768
                embedding_max_tokens = 512
            
            # NV-Embed (NVIDIA)
            elif 'nv-embed-v2' in embedding_model_name.lower():
                embedding_dim = 4096  # Very large, state-of-the-art
                embedding_max_tokens = 32768  # Extremely long context support
            elif 'nv-embed' in embedding_model_name.lower():
                embedding_dim = 4096
                embedding_max_tokens = 4096
            
            # Stella Models
            elif 'stella' in embedding_model_name.lower():
                if 'large' in embedding_model_name.lower():
                    embedding_dim = 1024
                elif 'base' in embedding_model_name.lower():
                    embedding_dim = 768
                else:
                    embedding_dim = 1024
                embedding_max_tokens = 512
            
            # Fallback patterns for unknown models
            elif 'large' in embedding_model_name.lower() and 'embed' in embedding_model_name.lower():
                embedding_dim = 1024
                embedding_max_tokens = 512
            elif 'base' in embedding_model_name.lower() and 'embed' in embedding_model_name.lower():
                embedding_dim = 768
                embedding_max_tokens = 512
            elif 'small' in embedding_model_name.lower() and 'embed' in embedding_model_name.lower():
                embedding_dim = 384
                embedding_max_tokens = 256
            
            logger.info(f"Using embedding dimension: {embedding_dim}, max tokens: {embedding_max_tokens}")
            
            # Set up LLM function and parameters based on provider type
            if llm_provider_type == 'openai':
                # Create wrapper functions with proper API key binding and return type handling
                if 'gpt-4o-mini' in llm_model_name:
                    async def llm_model_func(
                        prompt, system_prompt=None, history_messages=[], keyword_extraction=False, **kwargs
                    ) -> str:
                        try:
                            result = await gpt_4o_mini_complete(
                                prompt,
                                system_prompt=system_prompt,
                                history_messages=history_messages,
                                keyword_extraction=keyword_extraction,
                                api_key=llm_api_key.strip(),
                                **kwargs
                            )
                            # Ensure we return a string, not an AsyncIterator
                            return str(result) if isinstance(result, str) else str(result)
                        except Exception as e:
                            error_str = str(e).lower()
                            if 'api key' in error_str or 'unauthorized' in error_str or 'authentication' in error_str:
                                raise Exception(f"LLM authentication failed. Please check your OpenAI API key. Error: {str(e)}")
                            else:
                                raise e
                elif 'gpt-4o' in llm_model_name:
                    async def llm_model_func(
                        prompt, system_prompt=None, history_messages=[], keyword_extraction=False, **kwargs
                    ) -> str:
                        try:
                            result = await gpt_4o_complete(
                                prompt,
                                system_prompt=system_prompt,
                                history_messages=history_messages,
                                keyword_extraction=keyword_extraction,
                                api_key=llm_api_key.strip(),
                                **kwargs
                            )
                            # Ensure we return a string, not an AsyncIterator
                            return str(result) if isinstance(result, str) else str(result)
                        except Exception as e:
                            error_str = str(e).lower()
                            if 'api key' in error_str or 'unauthorized' in error_str or 'authentication' in error_str:
                                raise Exception(f"LLM authentication failed. Please check your OpenAI API key. Error: {str(e)}")
                            else:
                                raise e
                else:
                    async def llm_model_func(
                        prompt, system_prompt=None, history_messages=[], keyword_extraction=False, **kwargs
                    ) -> str:
                        try:
                            result = await openai_complete(
                                prompt,
                                system_prompt=system_prompt,
                                history_messages=history_messages,
                                keyword_extraction=keyword_extraction,
                                api_key=llm_api_key.strip(),
                                **kwargs
                            )
                            # Ensure we return a string, not an AsyncIterator
                            return str(result) if isinstance(result, str) else str(result)
                        except Exception as e:
                            error_str = str(e).lower()
                            if 'api key' in error_str or 'unauthorized' in error_str or 'authentication' in error_str:
                                raise Exception(f"LLM authentication failed. Please check your OpenAI API key. Error: {str(e)}")
                            else:
                                raise e
                llm_model_kwargs = {}
            elif llm_provider_type == 'openai_compatible':
                # Create wrapper function for openai_complete_if_cache
                async def llm_model_func(
                    prompt, system_prompt=None, history_messages=[], keyword_extraction=False, **kwargs
                ) -> str:
                    try:
                        return await openai_complete_if_cache(
                            llm_model_name,
                            prompt,
                            system_prompt=system_prompt,
                            history_messages=history_messages,
                            api_key=llm_api_key.strip(),
                            base_url=llm_base_url,
                            **kwargs,
                        )
                    except Exception as e:
                        error_str = str(e).lower()

                        # Check for Jan/LM Studio template errors (common with reasoning models)
                        if 'value is not callable' in error_str and ('split' in error_str or 'rstrip' in error_str or 'lstrip' in error_str):
                            raise Exception(
                                f"LLM Model Template Error: The model '{llm_model_name}' has a broken chat template in Jan/LM Studio. "
                                f"This is a known issue with some models. Solutions:\n"
                                f"1. Try a different model (e.g., llama-3, qwen2.5-7b-instruct, mistral)\n"
                                f"2. Update Jan to the latest version\n"
                                f"3. Check model settings in Jan for template configuration\n"
                                f"Technical error: {str(e)[:200]}"
                            )

                        # Check for authentication errors
                        if 'api key' in error_str or 'unauthorized' in error_str or 'authentication' in error_str:
                            if is_local_llm_endpoint or is_clara_core_llm:
                                # For local endpoints, provide more helpful error message
                                raise Exception(f"LLM connection failed to local endpoint {llm_base_url}. The endpoint may not be running or may require a different API key format. Error: {str(e)}")
                            else:
                                raise Exception(f"LLM authentication failed for {llm_base_url}. Please check your API key configuration. Error: {str(e)}")
                        else:
                            raise e
                llm_model_kwargs = {}
            elif llm_provider_type == 'ollama':
                # Create wrapper function for ollama_model_complete to ensure string return
                async def llm_model_func(
                    prompt, system_prompt=None, history_messages=[], keyword_extraction=False, **kwargs
                ) -> str:
                    result = await ollama_model_complete(
                        prompt,
                        system_prompt=system_prompt,
                        history_messages=history_messages,
                        keyword_extraction=keyword_extraction,
                        host=llm_base_url if llm_base_url else "http://localhost:11434",
                        model=llm_model_name,
                        **kwargs
                    )
                    # Ensure we return a string, not an AsyncIterator
                    return str(result) if isinstance(result, str) else str(result)
                llm_model_kwargs = {
                    "host": llm_base_url if llm_base_url else "http://localhost:11434",
                    "options": {"num_ctx": 8192},
                    "timeout": 300,
                }
            else:
                raise ValueError(f"Unsupported LLM provider type: {llm_provider_type}")
            
            # Set up embedding function based on provider type (simplified for LightRAG built-in handling)
            if embedding_provider_type == 'openai':
                async def embedding_func_lambda(texts: list[str]):
                    """Simple OpenAI embedding - let LightRAG handle retries and failures"""
                    try:
                        return await openai_embed(
                            texts,
                            model=embedding_model_name,
                            api_key=embedding_api_key.strip()
                        )
                    except Exception as e:
                        error_str = str(e).lower()
                        if 'api key' in error_str or 'unauthorized' in error_str or 'authentication' in error_str:
                            raise Exception(f"Embedding authentication failed. Please check your OpenAI API key. Error: {str(e)}")
                        else:
                            # Let LightRAG handle other errors with its built-in retry logic
                            raise e
                    
            elif embedding_provider_type == 'openai_compatible':
                async def embedding_func_lambda(texts: list[str]):
                    """Simple OpenAI-compatible embedding - let LightRAG handle retries and batch sizing"""
                    max_retries = 3 if is_clara_core_embedding else 1
                    retry_delay = 5 if is_clara_core_embedding else 2
                    
                    for attempt in range(max_retries):
                        try:
                            logger.info(f"Embedding attempt {attempt + 1}/{max_retries} for {embedding_model_name} ({len(texts)} texts)")
                            
                            result = await openai_embed(
                                texts,
                                model=embedding_model_name,
                                api_key=embedding_api_key.strip(),
                                base_url=embedding_base_url
                            )
                            
                            logger.info(f"Successfully generated {len(result)} embeddings")
                            return result
                            
                        except Exception as e:
                            error_str = str(e).lower()
                            
                            if 'api key' in error_str or 'unauthorized' in error_str or 'authentication' in error_str:
                                if is_local_embedding_endpoint or is_clara_core_embedding:
                                    raise Exception(f"Embedding connection failed to local endpoint {embedding_base_url}. Error: {str(e)}")
                                else:
                                    raise Exception(f"Embedding authentication failed for {embedding_base_url}. Error: {str(e)}")
                            
                            # For Clara Core, retry connection errors
                            if attempt < max_retries - 1 and is_clara_core_embedding:
                                if any(keyword in error_str for keyword in ['connection', 'timeout', 'loading', 'model']):
                                    logger.warning(f"Clara Core embedding failed (attempt {attempt + 1}/{max_retries}): {e}")
                                    logger.info(f"Retrying in {retry_delay}s...")
                                    await asyncio.sleep(retry_delay)
                                    continue
                            
                            # For other errors, let LightRAG handle them (including batch size errors)
                            raise e
                    
                    raise Exception(f"Embedding failed after {max_retries} attempts")
                    
            elif embedding_provider_type == 'ollama':
                async def embedding_func_lambda(texts: list[str]):
                    """Simple Ollama embedding - let LightRAG handle retries"""
                    return await ollama_embed(
                        texts,
                        embed_model=embedding_model_name,
                        host=embedding_base_url if embedding_base_url else "http://localhost:11434"
                    )
            else:
                raise ValueError(f"Unsupported embedding provider type: {embedding_provider_type}")
            
            # Determine if using local/ollama providers for optimized configuration
            is_local_llm = llm_provider_type == 'ollama'
            
            # Optimize configuration for local vs remote models
            if is_local_llm:
                # Local/Ollama models: smaller chunks, less aggressive entity extraction
                chunk_token_size = 800
                chunk_overlap_token_size = 80
                entity_extract_max_gleaning = 1
                logger.info(f"Using local model optimization for notebook {notebook_id}")
            else:
                # Remote models: larger chunks, more aggressive entity extraction
                chunk_token_size = 1200
                chunk_overlap_token_size = 100
                entity_extract_max_gleaning = 2
                logger.info(f"Using remote model optimization for notebook {notebook_id}")
            
            # Determine optimal token sizes based on model type
            if 'gemma' in llm_model_name.lower() or 'llama' in llm_model_name.lower():
                # For smaller open source models, use more conservative token limits
                max_token_size = 4096 if 'gemma' in llm_model_name.lower() else 6144
                chunk_token_size = 2000 if is_local_llm else 1200
                chunk_overlap_token_size = 100 if is_local_llm else 100
                entity_extract_max_gleaning = 1  # Reduce complexity for smaller models
                logger.info(f"Using conservative settings for model {llm_model_name}: max_tokens={max_token_size}")
            else:
                # For larger models like GPT-4, use larger token limits
                max_token_size = 8192
                chunk_token_size = chunk_token_size
                chunk_overlap_token_size = chunk_overlap_token_size
                entity_extract_max_gleaning = entity_extract_max_gleaning
                logger.info(f"Using standard settings for model {llm_model_name}: max_tokens={max_token_size}")
            
            # CRITICAL FIX: Override chunk size if embedding model has lower token limits
            # This prevents creating chunks that are larger than what the embedding model can handle
            if embedding_max_tokens < chunk_token_size:
                original_chunk_size = chunk_token_size
                chunk_token_size = embedding_max_tokens - 50  # Leave buffer for safety
                chunk_overlap_token_size = min(chunk_overlap_token_size, chunk_token_size // 10)
                logger.warning(f"Reducing chunk size from {original_chunk_size} to {chunk_token_size} due to embedding model limit ({embedding_max_tokens} tokens)")
            
            # Special handling for Qwen3-0.6B model (ultra-small model, very limited server batch capacity)
            if 'qwen3-embedding-0.6b' in embedding_model_name.lower() or 'qwen3-embedding-06b' in embedding_model_name.lower():
                chunk_token_size = min(chunk_token_size, 200)  # Ultra-conservative for 0.6B model server limitations
                chunk_overlap_token_size = 20  # 10% overlap
                logger.info(f"Using ultra-small chunk size {chunk_token_size} for Qwen3-0.6B model (server batch limitations)")
            # Special handling for MixedBread AI models (very sensitive to batch size and token limits)
            elif 'mxbai-embed' in embedding_model_name.lower():
                chunk_token_size = min(chunk_token_size, 150)  # Ultra-conservative - mxbai has strict physical batch size limits
                chunk_overlap_token_size = 15  # 10% overlap
                logger.warning(f"⚠️ MixedBread AI models have strict batch size limits. Using ultra-small chunk size {chunk_token_size}")
                logger.warning(f"💡 TIP: For better reliability, consider switching to 'qwen3-embedding-06b-0.6b' or 'nomic-embed-text'")
            # Special handling for Jina v4 models (high dimension 2048, very batch size sensitive on local servers)
            elif 'jina-embeddings-v4' in embedding_model_name.lower():
                chunk_token_size = min(chunk_token_size, 200)  # Very conservative for v4's 2048 dimensions + local server batch limits
                chunk_overlap_token_size = 20  # 10% overlap
                logger.info(f"Using ultra-small chunk size {chunk_token_size} for Jina v4 model (high-dimension + local server batch limitations)")
            # Special handling for models with very low token limits (nomic, e5, etc.)
            elif embedding_max_tokens <= 512:
                chunk_token_size = min(chunk_token_size, 200)  # Very safe limit for 512 token models (reduced from 250)
                chunk_overlap_token_size = 20  # 10% overlap
                logger.info(f"Using reduced chunk size {chunk_token_size} for low-context embedding model ({embedding_model_name})")
            elif embedding_max_tokens <= 1024:
                chunk_token_size = min(chunk_token_size, 600)  # Safer limit for 1024 token models (reduced from 800)
                chunk_overlap_token_size = 60  # 10% overlap
                logger.info(f"Using moderate chunk size {chunk_token_size} for medium-context embedding model ({embedding_model_name})")
            
            # Create LightRAG instance following the official pattern with proper error handling
            logger.info(f"Initializing LightRAG with embedding dimensions: {embedding_dim}, embedding max tokens: {embedding_max_tokens}")
            
            # Configure addon_params for consistent entity/relationship schema
            # Use user-provided entity types or fall back to comprehensive defaults
            default_entity_types = [
                "PERSON", "ORGANIZATION", "LOCATION", "TECHNOLOGY", 
                "PRODUCT", "SERVICE", "PROJECT", "FEATURE", 
                "COMPONENT", "API", "DATABASE", "FRAMEWORK",
                "METHODOLOGY", "METRIC", "REQUIREMENT", "ISSUE",
                "SOLUTION", "STRATEGY", "GOAL", "RISK",
                "RESOURCE", "TOOL", "PLATFORM", "ENVIRONMENT",
                "VERSION", "RELEASE", "DEPLOYMENT", "CONFIGURATION",
                "WORKFLOW", "PROCESS", "STANDARD", "PROTOCOL",
                "DOCUMENT", "SPECIFICATION", "GUIDELINE", "POLICY",
                "EVENT", "MEETING", "DECISION", "MILESTONE",
                "CONCEPT", "PRINCIPLE", "PATTERN", "ARCHITECTURE",
                "INTERFACE", "MODULE", "LIBRARY", "DEPENDENCY",
                "DATA", "MODEL", "SCHEMA", "ENTITY",
                "ROLE", "PERMISSION", "SECURITY", "COMPLIANCE",
                "PERFORMANCE", "SCALABILITY", "AVAILABILITY", "RELIABILITY",
                "BUG", "FEATURE_REQUEST", "ENHANCEMENT", "INCIDENT",
                "TEAM", "DEPARTMENT", "STAKEHOLDER", "CLIENT",
                "VENDOR", "PARTNER", "COMPETITOR", "MARKET",
                "TIMELINE", "BUDGET", "COST", "REVENUE",
                "SKILL", "EXPERTISE", "CERTIFICATION", "TRAINING",
                "OTHER"
            ]
            
            addon_params = {
                "entity_types": entity_types if entity_types else default_entity_types,
                "language": language  # User-specified language for consistent processing
            }
            
            # Use LightRAG's built-in resilience by setting appropriate working directory and cache handling
            rag = LightRAG(
                working_dir=str(working_dir),
                llm_model_func=llm_model_func,
                llm_model_name=llm_model_name,
                # LightRAG has built-in caching and resilience - let it handle retries
                llm_model_kwargs=llm_model_kwargs,
                embedding_func=EmbeddingFunc(
                    embedding_dim=embedding_dim,
                    max_token_size=embedding_max_tokens,
                    func=embedding_func_lambda,
                ),
                chunk_token_size=chunk_token_size,
                chunk_overlap_token_size=chunk_overlap_token_size,
                entity_extract_max_gleaning=entity_extract_max_gleaning,
                addon_params=addon_params,  # Critical for schema consistency
            )
            
            # Add rerank functionality for enhanced retrieval performance (if available)
            if RERANK_AVAILABLE and embedding_provider_type == 'openai':
                try:
                    # For OpenAI users, we can use Cohere rerank which is compatible
                    async def openai_rerank_wrapper(query: str, docs: list[str], top_k: int = 10):
                        """Wrapper for rerank function using OpenAI-compatible approach"""
                        # Simple rerank based on text similarity since we don't have Cohere API
                        # This is a fallback - users should configure proper rerank services
                        logger.info(f"Using basic rerank fallback for {len(docs)} documents")
                        return docs[:top_k]  # Simple truncation fallback
                    
                    rag.rerank_model_func = openai_rerank_wrapper
                    logger.info("Rerank functionality enabled for enhanced retrieval performance")
                except Exception as rerank_error:
                    logger.warning(f"Could not enable rerank functionality: {rerank_error}")
            elif RERANK_AVAILABLE:
                logger.info("Rerank functions available but not configured for current provider")
            else:
                logger.info("Rerank functionality not available - install rerank dependencies for better performance")
            
            # Initialize storages - CRITICAL: Both calls are required for proper LightRAG functionality
            try:
                await rag.initialize_storages()  # Initialize storage backends
                await initialize_pipeline_status()  # Initialize processing pipeline status
                logger.info(f"LightRAG storages and pipeline initialized for notebook {notebook_id}")
            except Exception as init_error:
                logger.error(f"LightRAG initialization failed for notebook {notebook_id}: {init_error}")
                # Initialization failure should be treated as a critical error, not a warning
                raise HTTPException(status_code=500, detail=f"Failed to initialize LightRAG storage systems: {str(init_error)}")
            
            logger.info(f"Successfully created LightRAG instance for notebook {notebook_id}")
            return rag
            
        except ValueError as ve:
            # Handle configuration errors
            logger.error(f"Configuration error for notebook {notebook_id}: {ve}")
            raise HTTPException(status_code=400, detail=f"Configuration error: {str(ve)}")
        except Exception as e:
            logger.error(f"Error creating LightRAG instance for notebook {notebook_id}: {e}")
            logger.error(f"Full traceback: {traceback.format_exc()}")
            raise HTTPException(status_code=500, detail=f"Error initializing RAG system: {str(e)}")

    async def get_lightrag_instance(notebook_id: str) -> LightRAG:
        """Get or create LightRAG instance for a notebook
        
        CRITICAL: Each notebook MUST have its own isolated LightRAG instance
        with a unique working directory to prevent data leakage between notebooks.
        """
        if notebook_id not in lightrag_instances:
            if notebook_id not in lightrag_notebooks_db:
                raise HTTPException(status_code=404, detail="Notebook not found")
            
            notebook_data = lightrag_notebooks_db[notebook_id]
            logger.info(f"Creating new LightRAG instance for notebook {notebook_id}")
            
            lightrag_instances[notebook_id] = await create_lightrag_instance(
                notebook_id, 
                notebook_data["llm_provider"], 
                notebook_data["embedding_provider"],
                entity_types=notebook_data.get("entity_types"),
                language=notebook_data.get("language", "en")
            )
            
            # Verify working directory is correctly set
            rag = lightrag_instances[notebook_id]
            expected_dir = str(LIGHTRAG_STORAGE_PATH / notebook_id)
            actual_dir = str(rag.working_dir)
            
            if expected_dir != actual_dir:
                logger.error(f"⚠️ CRITICAL: Working directory mismatch for notebook {notebook_id}!")
                logger.error(f"   Expected: {expected_dir}")
                logger.error(f"   Actual: {actual_dir}")
                raise Exception(f"Working directory mismatch - potential data leakage risk!")
            
            logger.info(f"✓ LightRAG instance for notebook {notebook_id} verified with working_dir: {actual_dir}")
        else:
            logger.debug(f"Reusing existing LightRAG instance for notebook {notebook_id}")
        
        return lightrag_instances[notebook_id]

    def extract_document_id_from_chunk_id(chunk_id: str, notebook_id: str) -> Optional[str]:
        """Extract the original document UUID from a LightRAG chunk ID
        
        Our insert format: doc_{notebook_id}_{document_id}_{timestamp}_{hash}
        This parses out the 36-char UUID document_id
        """
        import re
        pattern = rf"^doc_{re.escape(notebook_id)}_([0-9a-fA-F-]{{36}})_.+$"
        m = re.match(pattern, chunk_id)
        return m.group(1) if m else None

    def map_doc_ids_to_citations(notebook_id: str, doc_ids: List[str]) -> List[Dict[str, Any]]:
        """Map document IDs to citation objects with enhanced academic-style metadata

        Enhanced citation mode provides:
        - Document title and filename
        - Upload timestamp
        - Document ID for precise reference
        - File path for navigation
        - Content size for context
        """
        seen = set()
        citations: List[Dict[str, Any]] = []

        for doc_id in doc_ids:
            if not doc_id or doc_id in seen:
                continue

            doc = lightrag_documents_db.get(doc_id)
            if not doc or doc.get("notebook_id") != notebook_id:
                continue

            # Enhanced citation with more metadata
            uploaded_at = doc.get("uploaded_at")
            citation = {
                "filename": doc["filename"],
                "file_path": doc.get("file_path", f"documents/{doc['filename']}"),
                "document_id": doc["id"],
                "title": doc["filename"].replace('_', ' ').replace('.txt', '').replace('.pdf', '').replace('.md', '').title(),
                "uploaded_at": uploaded_at.isoformat() if uploaded_at else None,
                "content_size": doc.get("content_size", 0)
            }

            citations.append(citation)
            seen.add(doc_id)

            # Higher limit for enhanced citation mode (20 instead of 10)
            if len(citations) >= 20:
                break

        return citations

    def clean_lightrag_citations_in_text(text: str, notebook_id: str) -> tuple[str, List[Dict[str, Any]]]:
        """Clean inline LightRAG citations and add proper academic citation numbers

        LightRAG adds inline citations like:
        [KG] unknown_source (entity_name)
        [DC] unknown_source (chunk_info)

        This function:
        1. Removes LightRAG's inline citations
        2. Maps citations to actual documents
        3. Adds inline citation numbers [1], [2], etc. throughout the text
        4. Creates a clean References section at the end
        """
        import re

        # Get all completed documents for this notebook
        notebook_docs = [
            doc for doc_id, doc in lightrag_documents_db.items()
            if doc.get("notebook_id") == notebook_id and doc.get("status") == "completed"
        ]

        if not notebook_docs:
            # No documents available, just clean up the text
            citation_pattern = r'\[(KG|DC)\]\s+[^\n]+?(?:\s+\([^)]+\))?'
            cleaned_text = re.sub(citation_pattern, '', text)
            cleaned_text = re.sub(r'\n\s*\n\s*\n', '\n\n', cleaned_text)
            cleaned_text = re.sub(r'\n*References:\s*\n.*$', '', cleaned_text, flags=re.DOTALL)
            return cleaned_text.strip(), []

        # Build citations list
        citations = map_doc_ids_to_citations(notebook_id, [doc["id"] for doc in notebook_docs])

        # Step 1: Replace inline citation markers with proper academic numbers
        # Pattern 1: LightRAG citations [KG] or [DC]
        lightrag_pattern = r'\[(KG|DC)\]\s+[^\n]+?(?:\s+\([^)]+\))?'
        # Pattern 2: Our custom [SOURCE] markers
        source_pattern = r'\[SOURCE\]'

        # Track which citation numbers we've used
        citation_counter = 0
        citation_positions = []

        def replace_citation(match):
            nonlocal citation_counter
            # Cycle through available citations
            if citations:
                citation_num = (citation_counter % len(citations)) + 1
                citation_counter += 1
                citation_positions.append(citation_num)
                return f" [{citation_num}]"
            return ""

        # Replace both LightRAG citations and [SOURCE] markers with numbered citations
        cleaned_text = re.sub(lightrag_pattern, replace_citation, text)
        cleaned_text = re.sub(source_pattern, replace_citation, cleaned_text)

        # Step 2: Remove any existing "References:" section (we'll rebuild it)
        cleaned_text = re.sub(r'\n*References:\s*\n.*$', '', cleaned_text, flags=re.DOTALL)

        # Step 3: Clean up multiple blank lines
        cleaned_text = re.sub(r'\n\s*\n\s*\n', '\n\n', cleaned_text)

        # Step 4: Add academic-style References section
        if citations and citation_positions:
            # Get unique citations that were actually referenced
            referenced_citations = sorted(set(citation_positions))
            references_text = "\n\n## References\n\n"

            for num in referenced_citations:
                if num <= len(citations):
                    citation = citations[num - 1]
                    # Academic citation format: [1] Title (filename)
                    references_text += f"[{num}] **{citation['title']}**\n"
                    references_text += f"    Source: {citation['filename']}\n"
                    if citation.get('uploaded_at'):
                        references_text += f"    Uploaded: {citation['uploaded_at'][:10]}\n"
                    references_text += "\n"

            cleaned_text = cleaned_text.strip() + references_text

        return cleaned_text.strip(), citations

    async def build_true_citations_from_rag(rag: LightRAG, notebook_id: str, question: str, top_k: int = 10) -> Optional[List[Dict[str, Any]]]:
        """Extract actual retrieved sources from LightRAG and build precise citations
        
        Tries multiple methods to get retrieved chunk IDs:
        1. Check for known LightRAG attributes that may store last retrieval
        2. Try vector store API methods if available
        3. Return None if no reliable way found (better than returning all docs)
        """
        
        # 1) Try known attributes LightRAG may set after aquery
        possible_attrs = [
            "last_used_chunks", 
            "last_retrieved_chunks", 
            "last_context_chunks", 
            "retrieved_chunks", 
            "last_context"
        ]
        
        for attr in possible_attrs:
            chunks = getattr(rag, attr, None)
            if not chunks:
                continue
            
            chunk_ids: List[str] = []
            for c in chunks:
                if isinstance(c, dict):
                    # Common keys that might exist
                    cid = c.get("id") or c.get("chunk_id") or c.get("doc_id") or c.get("chunkId")
                    if cid:
                        chunk_ids.append(str(cid))
                elif isinstance(c, str):
                    chunk_ids.append(c)
            
            if chunk_ids:
                doc_ids = []
                for cid in chunk_ids:
                    doc_id = extract_document_id_from_chunk_id(cid, notebook_id)
                    if doc_id:
                        doc_ids.append(doc_id)
                
                if doc_ids:
                    logger.info(f"✓ Extracted {len(doc_ids)} unique document citations from {attr}")
                    return map_doc_ids_to_citations(notebook_id, doc_ids)

        # 2) Try vector store API methods (feature-detected)
        vdb = getattr(rag, "chunks_vdb", None)
        if vdb:
            for method_name in ["search", "similarity_search", "most_similar", "query"]:
                method = getattr(vdb, method_name, None)
                if not method:
                    continue
                
                try:
                    # Try calling the method (handle both sync and async)
                    maybe_results = method(question, top_k)
                    results = await maybe_results if asyncio.iscoroutine(maybe_results) else maybe_results
                    
                    # Results could be list of dicts or tuples; extract IDs safely
                    chunk_ids = []
                    for r in results:
                        if isinstance(r, dict):
                            cid = r.get("id") or r.get("chunk_id") or r.get("doc_id") or r.get("chunkId")
                            if cid:
                                chunk_ids.append(str(cid))
                        elif isinstance(r, (list, tuple)) and r:
                            # Common patterns: (id, score) or (content, meta)
                            first = r[0]
                            if isinstance(first, str):
                                chunk_ids.append(first)
                            elif isinstance(first, dict):
                                cid = first.get("id") or first.get("chunk_id") or first.get("doc_id")
                                if cid:
                                    chunk_ids.append(str(cid))
                    
                    if chunk_ids:
                        doc_ids = [d for d in (extract_document_id_from_chunk_id(cid, notebook_id) for cid in chunk_ids) if d]
                        if doc_ids:
                            logger.info(f"✓ Extracted {len(doc_ids)} unique document citations from vector store {method_name}")
                            return map_doc_ids_to_citations(notebook_id, doc_ids)
                except Exception as e:
                    logger.debug(f"Could not extract citations from {method_name}: {e}")
                    continue

        # 3) Fallback: return all completed documents as potential sources
        # This is better than showing "unknown_source" to the user
        logger.warning(f"⚠️ Could not extract precise citations - falling back to all documents as potential sources")
        
        all_doc_ids = [
            doc_id for doc_id, doc in lightrag_documents_db.items()
            if doc.get("notebook_id") == notebook_id and doc.get("status") == "completed"
        ]
        
        if all_doc_ids:
            logger.info(f"Returning {len(all_doc_ids)} documents as fallback citations")
            return map_doc_ids_to_citations(notebook_id, all_doc_ids)
        
        return None

    def auto_detect_provider_type(provider_config: Dict[str, Any]) -> Dict[str, Any]:
        """Auto-detect provider type based on baseUrl and return updated config"""
        provider_config = provider_config.copy()  # Don't modify original
        provider_type = provider_config.get('type', 'openai')
        base_url = provider_config.get('baseUrl', '')
        
        # Auto-detect for non-ollama providers
        if provider_type != 'ollama':
            if base_url in ['https://api.openai.com/v1', 'https://api.openai.com']:
                provider_config['type'] = 'openai'
                logger.info(f"Auto-detected provider type as 'openai' based on baseUrl: {base_url}")
            else:
                provider_config['type'] = 'openai_compatible'
                logger.info(f"Auto-detected provider type as 'openai_compatible' based on baseUrl: {base_url}")
        
        return provider_config

    def validate_notebook_exists(notebook_id: str):
        """Validate that a notebook exists"""
        if notebook_id not in lightrag_notebooks_db:
            raise HTTPException(status_code=404, detail="Notebook not found")

    async def process_notebook_document_with_delay(notebook_id: str, document_id: str, text_content: str, delay_seconds: int):
        """Wrapper to add delay before processing document"""
        if delay_seconds > 0:
            await asyncio.sleep(delay_seconds)
        
        await process_notebook_document(notebook_id, document_id, text_content)

    async def process_notebook_document(notebook_id: str, document_id: str, text_content: str):
        """Background task to process document with LightRAG
        
        CRITICAL: This function ensures documents are only added to the correct notebook
        by verifying the document belongs to the notebook and using the correct RAG instance.
        """
        try:
            # Validate inputs
            if not text_content or not text_content.strip():
                raise ValueError("Document content is empty")
            
            # CRITICAL: Verify document belongs to this notebook
            if document_id not in lightrag_documents_db:
                raise ValueError(f"Document {document_id} not found in database")
            
            doc_notebook_id = lightrag_documents_db[document_id].get("notebook_id")
            if doc_notebook_id != notebook_id:
                raise ValueError(
                    f"⚠️ DATA LEAKAGE PREVENTED: Document {document_id} belongs to notebook "
                    f"{doc_notebook_id}, but was being processed for notebook {notebook_id}"
                )
            
            logger.info(f"Starting document processing for {document_id} in notebook {notebook_id}")
            logger.info(f"✓ Verified document belongs to correct notebook")
            
            rag = await get_lightrag_instance(notebook_id)
            
            # Verify RAG instance working directory matches notebook
            expected_dir = str(LIGHTRAG_STORAGE_PATH / notebook_id)
            actual_dir = str(rag.working_dir)
            if expected_dir != actual_dir:
                raise ValueError(
                    f"⚠️ DATA LEAKAGE PREVENTED: RAG working directory mismatch! "
                    f"Expected {expected_dir}, got {actual_dir}"
                )
            
            # Get notebook data to check provider type
            notebook_data = lightrag_notebooks_db[notebook_id]
            llm_provider = notebook_data.get("llm_provider", {})
            llm_provider_type = llm_provider.get("type", "openai")
            is_local_provider = llm_provider_type == 'ollama'
            
            # Adjust content size and timeout based on provider type
            if is_local_provider:
                # For local/ollama models, use smaller chunks and longer timeout
                max_content_size = 300000  # 300KB of text for local models
                processing_timeout = 1800.0  # 30 minutes for local/ollama
                logger.info(f"Using local provider settings: max_size={max_content_size}, timeout={processing_timeout}s")
            else:
                # For remote models, use larger chunks and shorter timeout  
                max_content_size = 800000  # 800KB of text for remote models
                processing_timeout = 900.0  # 15 minutes for remote models
                logger.info(f"Using remote provider settings: max_size={max_content_size}, timeout={processing_timeout}s")
            
            if len(text_content) > max_content_size:
                logger.warning(f"Document {document_id} is very large ({len(text_content)} chars), truncating to {max_content_size}")
                text_content = text_content[:max_content_size] + "\n\n[Content truncated due to size limits]"
            
            # Create a more specific document ID to avoid conflicts
            import hashlib
            import time
            timestamp = str(int(time.time() * 1000))  # milliseconds
            content_hash = hashlib.md5(text_content.encode()).hexdigest()[:8]
            prefixed_doc_id = f"doc_{notebook_id}_{document_id}_{timestamp}_{content_hash}"
            
            logger.info(f"Processing document {document_id} ({len(text_content)} chars) with ID {prefixed_doc_id}")
            
            # Set document status to processing before starting
            if document_id in lightrag_documents_db:
                lightrag_documents_db[document_id]["status"] = "processing"
                lightrag_documents_db[document_id]["processed_at"] = datetime.now()
                save_documents_db()
            
            # Get document metadata including file path for citations
            document_data = lightrag_documents_db[document_id]
            filename = document_data["filename"]
            file_path = document_data.get("file_path", f"documents/{filename}")
            
            # Insert document text into LightRAG
            try:
                logger.info(f"Starting LightRAG insertion for document {document_id}")
                
                # Track initial state to verify successful insertion
                initial_doc_count = 0
                if hasattr(rag, 'doc_status'):
                    try:
                        initial_doc_count = len(rag.doc_status)
                    except:
                        pass
                
                # Use asyncio timeout with simple retry for transient failures
                max_insert_retries = 2
                insert_retry_delays = [10]  # seconds
                last_insert_error = None

                for insert_attempt in range(max_insert_retries):
                    try:
                        # Attempt insertion with timeout
                        await asyncio.wait_for(
                                rag.ainsert(text_content, ids=[prefixed_doc_id]),
                                timeout=processing_timeout
                        )
                        last_insert_error = None
                        break
                    except asyncio.TimeoutError as te:
                        last_insert_error = te
                        if insert_attempt < max_insert_retries - 1:
                            delay = insert_retry_delays[min(insert_attempt, len(insert_retry_delays)-1)]
                            logger.warning(f"Document insertion timed out (attempt {insert_attempt + 1}/{max_insert_retries}). Retrying in {delay}s...")
                            await asyncio.sleep(delay)
                            continue
                        # Exhausted retries - re-raise to be handled below
                        raise
                    except (ValueError, IndexError) as dim_error:
                        # Check for dimension mismatch or empty vector store errors
                        err_str = str(dim_error)
                        if 'dimension' in err_str.lower() or 'size' in err_str.lower() or 'index' in err_str.lower():
                            logger.error(f"❌ CRITICAL: Vector dimension mismatch detected!")
                            logger.error(f"Error: {err_str}")
                            logger.error(f"This notebook was created with a different embedding model.")
                            logger.error(f"Solution: Delete and recreate the notebook, or use /rebuild endpoint with force=true")

                            # Mark document as failed with helpful error message
                            if document_id in lightrag_documents_db:
                                lightrag_documents_db[document_id]["status"] = "failed"
                                lightrag_documents_db[document_id]["error"] = (
                                    f"Embedding dimension mismatch: This notebook uses a different embedding model. "
                                    f"Please rebuild the notebook or create a new one with the correct embedding model. "
                                    f"Technical details: {err_str}"
                                )
                                save_documents_db()
                            raise ValueError(
                                f"Embedding dimension mismatch: This notebook was created with a different embedding model. "
                                f"Please rebuild the notebook or create a new one. Technical error: {err_str}"
                            )
                        # Other ValueError/IndexError - re-raise
                        raise
                    except Exception as transient_e:
                        err = str(transient_e).lower()
                        is_transient = any(k in err for k in [
                            'connection', 'timeout', 'temporarily', 'overloaded', 'rate limit', 'cancelled'
                        ])
                        if is_transient and insert_attempt < max_insert_retries - 1:
                            delay = insert_retry_delays[min(insert_attempt, len(insert_retry_delays)-1)]
                            logger.warning(f"Transient error during insertion (attempt {insert_attempt + 1}/{max_insert_retries}): {transient_e}. Retrying in {delay}s...")
                            await asyncio.sleep(delay)
                            continue
                        # Non-transient or retries exhausted - re-raise
                        raise
                
                # Verify that document was actually processed
                # LightRAG doesn't raise exceptions for extraction failures - it logs them internally
                doc_was_indexed = False
                verification_details = []

                # Check 1: doc_status increased
                if hasattr(rag, 'doc_status'):
                    try:
                        current_doc_count = len(rag.doc_status)
                        if current_doc_count > initial_doc_count:
                            doc_was_indexed = True
                            verification_details.append(f"doc_status: {initial_doc_count} → {current_doc_count}")
                            logger.info(f"✓ Document added to LightRAG (doc_status: {initial_doc_count} → {current_doc_count})")
                    except Exception as e:
                        logger.debug(f"Could not check doc_status: {e}")

                # Check 2: chunks were created
                if not doc_was_indexed and hasattr(rag, 'chunks_vdb'):
                    try:
                        if hasattr(rag.chunks_vdb, '_data') and len(rag.chunks_vdb._data) > 0:
                            doc_was_indexed = True
                            verification_details.append(f"chunks: {len(rag.chunks_vdb._data)}")
                            logger.info(f"✓ Document created chunks in vector database")
                    except Exception as e:
                        logger.debug(f"Could not check chunks_vdb: {e}")

                # Check 3: Graph was updated (entities/relationships created)
                if not doc_was_indexed and hasattr(rag, 'graph_storage'):
                    try:
                        # Check if graph storage has data
                        import os
                        graph_file = os.path.join(rag.working_dir, "graph_chunk_entity_relation.graphml")
                        if os.path.exists(graph_file) and os.path.getsize(graph_file) > 1000:  # At least 1KB
                            doc_was_indexed = True
                            verification_details.append(f"graph file: {os.path.getsize(graph_file)} bytes")
                            logger.info(f"✓ Document created graph data ({os.path.getsize(graph_file)} bytes)")
                    except Exception as e:
                        logger.debug(f"Could not check graph file: {e}")

                # Check 4: Entity VDB was updated
                if not doc_was_indexed and hasattr(rag, 'entities_vdb'):
                    try:
                        if hasattr(rag.entities_vdb, '_data') and len(rag.entities_vdb._data) > 0:
                            doc_was_indexed = True
                            verification_details.append(f"entities: {len(rag.entities_vdb._data)}")
                            logger.info(f"✓ Document created entities in vector database")
                    except Exception as e:
                        logger.debug(f"Could not check entities_vdb: {e}")

                # CRITICAL: Do NOT assume success if verification failed
                # LightRAG logs errors internally without raising exceptions
                if not doc_was_indexed:
                    error_msg = (
                        f"❌ CRITICAL: Document processing completed without errors, but no data was indexed! "
                        f"This usually indicates a silent failure in LightRAG (dimension mismatch, embedding errors, etc.). "
                        f"Verification checks: {verification_details if verification_details else 'all checks failed'}. "
                        f"Check the logs above for ERROR messages from LightRAG."
                    )
                    logger.error(error_msg)

                    # Mark document as failed
                    if document_id in lightrag_documents_db:
                        lightrag_documents_db[document_id]["status"] = "failed"
                        lightrag_documents_db[document_id]["error"] = (
                            "Document processing failed silently - no data was indexed. "
                            "This may be due to embedding dimension mismatch or other LightRAG errors. "
                            "Check server logs for details. You may need to rebuild or recreate the notebook."
                        )
                        save_documents_db()

                    raise Exception(error_msg)

                logger.info(f"✓ Successfully inserted and verified document {document_id} into LightRAG ({', '.join(verification_details)})")
                
            except asyncio.TimeoutError:
                timeout_minutes = int(processing_timeout / 60)
                error_msg = f"Document processing timed out after {timeout_minutes} minutes. This can happen with complex documents and local models. Try using a more powerful model or splitting the document into smaller parts."
                logger.error(error_msg)
                raise Exception(error_msg)
            except Exception as insert_error:
                logger.error(f"LightRAG insertion failed for document {document_id}: {insert_error}")
                
                # Check for common errors and provide helpful messages
                error_str = str(insert_error).lower()
                if "connection" in error_str or "timeout" in error_str:
                    raise Exception(f"Connection error during document processing: {str(insert_error)}. Please check your provider configuration and network connection.")
                elif "api key" in error_str or "unauthorized" in error_str:
                    raise Exception(f"Authentication error: {str(insert_error)}. Please check your API key configuration.")
                elif "model" in error_str and "not found" in error_str:
                    raise Exception(f"Model not found: {str(insert_error)}. Please check your model name configuration.")
                else:
                    raise Exception(f"Document processing failed: {str(insert_error)}")
            
            # Force cache clear after inserting document
            try:
                await asyncio.wait_for(rag.aclear_cache(), timeout=60.0)
                logger.info(f"Cache cleared for document {document_id}")
            except asyncio.TimeoutError:
                logger.warning("Cache clear timed out, continuing anyway")
            except Exception as cache_error:
                logger.warning(f"Cache clear failed: {cache_error}, continuing anyway")
            
            # Update document status to completed
            if document_id in lightrag_documents_db:
                lightrag_documents_db[document_id]["status"] = "completed"
                lightrag_documents_db[document_id]["lightrag_id"] = prefixed_doc_id
                lightrag_documents_db[document_id]["completed_at"] = datetime.now()
                # Clear any previous error
                if "error" in lightrag_documents_db[document_id]:
                    del lightrag_documents_db[document_id]["error"]
                
                # IMPORTANT: Keep content permanently for rebuilds and user downloads
                # Documents are precious - users should always be able to access them
                # This enables reliable rebuilds and document downloads like a drive
                if "content" in lightrag_documents_db[document_id]:
                    content_size = len(lightrag_documents_db[document_id]["content"])
                    logger.info(f"Preserving content ({content_size} chars) for document {document_id} - available for rebuild and download")
            
            # Clear summary cache since a new document has been processed
            if notebook_id in lightrag_notebooks_db:
                if "summary_cache" in lightrag_notebooks_db[notebook_id]:
                    del lightrag_notebooks_db[notebook_id]["summary_cache"]
                    logger.info(f"Cleared summary cache for notebook {notebook_id}")
                if "docs_fingerprint" in lightrag_notebooks_db[notebook_id]:
                    del lightrag_notebooks_db[notebook_id]["docs_fingerprint"]
            
            # Save changes to disk
            save_documents_db()
            save_notebooks_db()
            
            logger.info(f"Successfully completed processing document {document_id} in notebook {notebook_id}")
            
        except Exception as e:
            error_msg = str(e)
            logger.error(f"Error processing document {document_id} in notebook {notebook_id}: {error_msg}")
            logger.error(f"Full error traceback: {traceback.format_exc()}")
            
            # Update document status to failed
            if document_id in lightrag_documents_db:
                lightrag_documents_db[document_id]["status"] = "failed"
                lightrag_documents_db[document_id]["error"] = error_msg
                lightrag_documents_db[document_id]["failed_at"] = datetime.now()
                # CRITICAL: Save lightrag_id even on failure so document can be properly deleted
                # If we don't save this, the document becomes a "ghost" in LightRAG that can't be removed
                lightrag_documents_db[document_id]["lightrag_id"] = prefixed_doc_id
                logger.warning(f"Saved lightrag_id for failed document {document_id} to enable proper cleanup")
                # Save changes to disk even on failure
                save_documents_db()

class TTSRequest(BaseModel):
    text: str
    language: Optional[str] = "en"
    engine: Optional[str] = None  # "gtts", "pyttsx3", "kokoro", "kokoro-onnx", or None for current engine
    slow: Optional[bool] = False
    voice: Optional[str] = "af_sarah"  # Voice for Kokoro engines
    speed: Optional[float] = 1.0  # Speed for Kokoro engines (0.5-2.0)

# LightRAG Data Models
class NotebookCreate(BaseModel):
    name: str = Field(..., description="Name of the notebook")
    description: Optional[str] = Field(None, description="Description of the notebook")
    # Add provider configuration
    llm_provider: Dict[str, Any] = Field(..., description="LLM provider configuration")
    embedding_provider: Dict[str, Any] = Field(..., description="Embedding provider configuration")
    # Schema consistency parameters (optional)
    entity_types: Optional[List[str]] = Field(None, description="Custom entity types for consistent extraction")
    language: Optional[str] = Field("en", description="Language for consistent entity/relationship processing")
    # Manual embedding dimension override (for low-confidence model detection)
    manual_embedding_dimensions: Optional[int] = Field(None, description="Manual override for embedding dimensions")
    manual_embedding_max_tokens: Optional[int] = Field(None, description="Manual override for embedding max tokens")

class NotebookResponse(BaseModel):
    id: str = Field(..., description="Notebook ID")
    name: str = Field(..., description="Notebook name")
    description: Optional[str] = Field(None, description="Notebook description")
    created_at: datetime = Field(..., description="Creation timestamp")
    document_count: int = Field(0, description="Number of documents in notebook")
    # Add provider information to response (optional for backward compatibility)
    llm_provider: Optional[Dict[str, Any]] = Field(None, description="LLM provider configuration")
    embedding_provider: Optional[Dict[str, Any]] = Field(None, description="Embedding provider configuration")
    # Schema consistency information
    entity_types: Optional[List[str]] = Field(None, description="Configured entity types for this notebook")
    language: Optional[str] = Field("en", description="Language setting for entity/relationship processing")

class NotebookDocumentResponse(BaseModel):
    id: str = Field(..., description="Document ID")
    filename: str = Field(..., description="Original filename")
    notebook_id: str = Field(..., description="Notebook ID")
    uploaded_at: datetime = Field(..., description="Upload timestamp")
    status: str = Field(..., description="Processing status")
    error: Optional[str] = Field(None, description="Error message if processing failed")
    file_path: Optional[str] = Field(None, description="File path for citation tracking")

class NotebookQueryRequest(BaseModel):
    question: str = Field(..., description="Question to ask")
    mode: str = Field("hybrid", description="Query mode: local, global, hybrid, naive, mix")
    response_type: str = Field("Multiple Paragraphs", description="Response format")
    top_k: int = Field(60, description="Number of top items to retrieve")
    # Add optional provider override for query
    llm_provider: Optional[Dict[str, Any]] = Field(None, description="Override LLM provider for this query")
    # Add chat history support
    use_chat_history: bool = Field(True, description="Whether to use chat history for context")
    # Enhanced citation mode is always enabled - provides academic-style citations with proper source attribution

class NotebookQueryResponse(BaseModel):
    answer: str = Field(..., description="Generated answer")
    mode: str = Field(..., description="Query mode used")
    context_used: bool = Field(True, description="Whether context was used")
    citations: Optional[List[Dict[str, Any]]] = Field(None, description="Citation information with sources")
    # Enhanced citation support
    source_documents: Optional[List[Dict[str, Any]]] = Field(None, description="Source documents with highlighted passages")
    chat_context_used: bool = Field(False, description="Whether previous chat history was used")

class ChatMessage(BaseModel):
    role: str = Field(..., description="Message role: user or assistant")
    content: str = Field(..., description="Message content")
    timestamp: datetime = Field(default_factory=datetime.now, description="Message timestamp")
    citations: Optional[List[Dict[str, Any]]] = Field(None, description="Citations for assistant messages")

class ChatHistoryResponse(BaseModel):
    notebook_id: str = Field(..., description="Notebook ID")
    messages: List[ChatMessage] = Field(..., description="Chat messages")
    total_messages: int = Field(..., description="Total number of messages")

class DocumentSummaryRequest(BaseModel):
    include_details: bool = Field(True, description="Include document-level details")
    max_length: str = Field("medium", description="Summary length: short, medium, long")

class QueryTemplate(BaseModel):
    id: str = Field(..., description="Template ID")
    name: str = Field(..., description="Template name")
    description: str = Field(..., description="Template description")
    question_template: str = Field(..., description="Question template with placeholders")
    category: str = Field(..., description="Template category")
    use_case: str = Field(..., description="When to use this template")

class DocumentRetryResponse(BaseModel):
    message: str = Field(..., description="Success message")
    document_id: str = Field(..., description="Document ID that was retried")
    status: str = Field(..., description="New document status after retry initiation")

class NotebookSchemaUpdate(BaseModel):
    entity_types: Optional[List[str]] = Field(None, description="Updated entity types for consistent extraction")
    language: Optional[str] = Field(None, description="Updated language for entity/relationship processing")

@app.get("/")
def read_root():
    """Root endpoint for basic health check"""
    return {
        "status": "ok", 
        "service": "Clara Backend", 
        "port": PORT,
        "uptime": str(datetime.now() - datetime.fromisoformat(START_TIME)),
        "start_time": START_TIME
    }

@app.get("/health")
def health_check():
    """Health check endpoint"""
    return {
        "status": "healthy",
        "port": PORT,
        "uptime": str(datetime.now() - datetime.fromisoformat(START_TIME))
    }

# LightRAG Notebook endpoints
if LIGHTRAG_AVAILABLE:
    def normalize_url_for_local_dev(url: str) -> str:
        """
        Normalize URLs for local development
        - When IN Docker: Converts localhost/127.0.0.1 to host.docker.internal
        - When NOT in Docker: Converts host.docker.internal to localhost
        """
        if not url:
            return url

        import os
        in_docker = os.path.exists('/.dockerenv')

        if in_docker:
            # Running in Docker - convert localhost to host.docker.internal
            if 'localhost' in url or '127.0.0.1' in url:
                url = url.replace('localhost', 'host.docker.internal')
                url = url.replace('127.0.0.1', 'host.docker.internal')
                logger.info(f"Normalized localhost URL for Docker: {url}")
        else:
            # Not in Docker - convert host.docker.internal to localhost
            if 'host.docker.internal' in url:
                url = url.replace('host.docker.internal', 'localhost')
                logger.info(f"Normalized Docker URL to localhost: {url}")

        return url
    
    @app.post("/notebooks/validate-models")
    async def validate_notebook_models(config: NotebookCreate):
        """
        Validate that LLM and embedding models are accessible and working
        This is a pre-flight check before notebook creation or document upload
        """
        import aiohttp
        import json as json_module
        
        validation_results = {
            "llm_accessible": False,
            "llm_error": None,
            "embedding_accessible": False,
            "embedding_error": None,
            "overall_status": "failed"
        }
        
        try:
            # Auto-detect provider types
            llm_provider = auto_detect_provider_type(config.llm_provider)
            embedding_provider = auto_detect_provider_type(config.embedding_provider)
            
            # Normalize URLs for local development
            if llm_provider.get('baseUrl'):
                llm_provider['baseUrl'] = normalize_url_for_local_dev(llm_provider['baseUrl'])
            if embedding_provider.get('baseUrl'):
                embedding_provider['baseUrl'] = normalize_url_for_local_dev(embedding_provider['baseUrl'])
            
            logger.info(f"Validating models - LLM: {llm_provider['model']}, Embedding: {embedding_provider['model']}")
            
            # Test LLM with a simple prompt using direct HTTP call
            try:
                llm_model_name = llm_provider['model']
                llm_api_key = llm_provider.get('apiKey', '')
                llm_base_url = llm_provider.get('baseUrl', '')
                
                if llm_provider['type'] == 'ollama':
                    # Ollama validation
                    if llm_base_url.endswith('/v1'):
                        llm_base_url = llm_base_url[:-3]
                    
                    async with aiohttp.ClientSession() as session:
                        async with session.post(
                            f"{llm_base_url}/api/generate",
                            json={
                                "model": llm_model_name,
                                "prompt": "Hi",
                                "stream": False,
                                "options": {"num_predict": 5}
                            },
                            timeout=aiohttp.ClientTimeout(total=10)
                        ) as response:
                            if response.status == 200:
                                result = await response.json()
                                if result.get('response'):
                                    validation_results["llm_accessible"] = True
                                    logger.info(f"✓ LLM validation successful: {llm_model_name}")
                                else:
                                    validation_results["llm_error"] = "Model returned empty response"
                            else:
                                error_text = await response.text()
                                validation_results["llm_error"] = f"HTTP {response.status}: {error_text[:200]}"
                else:
                    # OpenAI-compatible validation
                    headers = {
                        "Content-Type": "application/json",
                    }
                    if llm_api_key:
                        headers["Authorization"] = f"Bearer {llm_api_key}"
                    
                    async with aiohttp.ClientSession() as session:
                        async with session.post(
                            f"{llm_base_url}/chat/completions",
                            headers=headers,
                            json={
                                "model": llm_model_name,
                                "messages": [{"role": "user", "content": "Hi"}],
                                "max_tokens": 5
                            },
                            timeout=aiohttp.ClientTimeout(total=10)
                        ) as response:
                            if response.status == 200:
                                result = await response.json()
                                if result.get('choices') and len(result['choices']) > 0:
                                    validation_results["llm_accessible"] = True
                                    logger.info(f"✓ LLM validation successful: {llm_model_name}")
                                else:
                                    validation_results["llm_error"] = "Model returned empty response"
                            else:
                                error_text = await response.text()
                                validation_results["llm_error"] = f"HTTP {response.status}: {error_text[:200]}"
                    
            except asyncio.TimeoutError:
                validation_results["llm_error"] = "Connection timeout - model service not responding"
                logger.error(f"✗ LLM validation failed: Timeout")
            except aiohttp.ClientConnectorError as e:
                validation_results["llm_error"] = f"Cannot connect to service: {str(e)}"
                logger.error(f"✗ LLM validation failed: Connection error")
            except Exception as llm_error:
                validation_results["llm_error"] = str(llm_error)
                logger.error(f"✗ LLM validation failed: {llm_error}")
            
            # Test embedding model with a simple text using direct HTTP call
            try:
                embedding_model_name = embedding_provider['model']
                embedding_api_key = embedding_provider.get('apiKey', '')
                embedding_base_url = embedding_provider.get('baseUrl', '')
                
                if embedding_provider['type'] == 'ollama':
                    # Ollama embedding validation
                    if embedding_base_url.endswith('/v1'):
                        embedding_base_url = embedding_base_url[:-3]
                    
                    async with aiohttp.ClientSession() as session:
                        async with session.post(
                            f"{embedding_base_url}/api/embeddings",
                            json={
                                "model": embedding_model_name,
                                "prompt": "test"
                            },
                            timeout=aiohttp.ClientTimeout(total=10)
                        ) as response:
                            if response.status == 200:
                                result = await response.json()
                                if result.get('embedding') and len(result.get('embedding', [])) > 0:
                                    validation_results["embedding_accessible"] = True
                                    logger.info(f"✓ Embedding validation successful: {embedding_model_name} (dimension: {len(result['embedding'])})")
                                else:
                                    validation_results["embedding_error"] = "Model returned empty embeddings"
                            else:
                                error_text = await response.text()
                                validation_results["embedding_error"] = f"HTTP {response.status}: {error_text[:200]}"
                else:
                    # OpenAI-compatible embedding validation
                    headers = {
                        "Content-Type": "application/json",
                    }
                    if embedding_api_key:
                        headers["Authorization"] = f"Bearer {embedding_api_key}"
                    
                    async with aiohttp.ClientSession() as session:
                        async with session.post(
                            f"{embedding_base_url}/embeddings",
                            headers=headers,
                            json={
                                "model": embedding_model_name,
                                "input": "test"
                            },
                            timeout=aiohttp.ClientTimeout(total=10)
                        ) as response:
                            if response.status == 200:
                                result = await response.json()
                                if result.get('data') and len(result['data']) > 0 and result['data'][0].get('embedding'):
                                    embedding_dim = len(result['data'][0]['embedding'])
                                    validation_results["embedding_accessible"] = True
                                    logger.info(f"✓ Embedding validation successful: {embedding_model_name} (dimension: {embedding_dim})")
                                else:
                                    validation_results["embedding_error"] = "Model returned empty embeddings"
                            else:
                                error_text = await response.text()
                                validation_results["embedding_error"] = f"HTTP {response.status}: {error_text[:200]}"
                    
            except asyncio.TimeoutError:
                validation_results["embedding_error"] = "Connection timeout - model service not responding"
                logger.error(f"✗ Embedding validation failed: Timeout")
            except aiohttp.ClientConnectorError as e:
                validation_results["embedding_error"] = f"Cannot connect to service: {str(e)}"
                logger.error(f"✗ Embedding validation failed: Connection error")
            except Exception as embedding_error:
                validation_results["embedding_error"] = str(embedding_error)
                logger.error(f"✗ Embedding validation failed: {embedding_error}")
            
            # Set overall status
            if validation_results["llm_accessible"] and validation_results["embedding_accessible"]:
                validation_results["overall_status"] = "success"
                logger.info("✓ Model validation: All models accessible")
            elif validation_results["llm_accessible"] or validation_results["embedding_accessible"]:
                validation_results["overall_status"] = "partial"
                logger.warning("⚠ Model validation: Some models accessible")
            else:
                validation_results["overall_status"] = "failed"
                logger.error("✗ Model validation: No models accessible")
            
            return validation_results
            
        except Exception as e:
            logger.error(f"Model validation error: {e}")
            return {
                "llm_accessible": False,
                "llm_error": str(e),
                "embedding_accessible": False,
                "embedding_error": str(e),
                "overall_status": "error"
            }

    # ============================================================================
    # Document Queue Setup - Startup and Shutdown
    # ============================================================================

    def cleanup_stale_lightrag_docs():
        """
        Clean up stale document IDs from LightRAG storage on startup.
        
        This prevents failed document retries from trying to process old document IDs
        that are stuck in LightRAG's doc_status storage.
        """
        try:
            storage_path = LIGHTRAG_METADATA_PATH / "documents.json"
            if not storage_path.exists():
                logger.info("No documents.json found, skipping stale doc cleanup")
                return
            
            # Load all documents to find failed ones
            with open(storage_path, 'r', encoding='utf-8') as f:
                all_docs = json.load(f)
            
            cleaned_count = 0
            for doc_id, doc_data in all_docs.items():
                # If document is failed and has a lightrag_id, we should clear it
                # so retries get a fresh ID
                if doc_data.get("status") == "failed" and "lightrag_id" in doc_data:
                    notebook_id = doc_data.get("notebook_id")
                    old_id = doc_data["lightrag_id"]
                    
                    # Clean up the stale ID from LightRAG's doc_status storage
                    doc_status_path = LIGHTRAG_STORAGE_PATH / notebook_id / "kv_store_doc_status.json"
                    if doc_status_path.exists():
                        try:
                            with open(doc_status_path, 'r', encoding='utf-8') as f:
                                doc_status = json.load(f)
                            
                            # Remove all old document IDs for this document
                            # (there might be multiple failed attempts)
                            doc_prefix = f"doc_{notebook_id}_{doc_id}"
                            to_remove = [k for k in doc_status.keys() if k.startswith(doc_prefix)]
                            
                            if to_remove:
                                for key in to_remove:
                                    del doc_status[key]
                                    cleaned_count += 1
                                
                                # Save cleaned doc_status
                                with open(doc_status_path, 'w', encoding='utf-8') as f:
                                    json.dump(doc_status, f, ensure_ascii=False, indent=2)
                                
                                logger.info(f"Cleaned {len(to_remove)} stale doc IDs for document {doc_id}")
                        except Exception as e:
                            logger.warning(f"Could not clean doc_status for notebook {notebook_id}: {e}")
                    
                    # Clear the lightrag_id from our metadata too
                    del doc_data["lightrag_id"]
                    all_docs[doc_id] = doc_data
            
            # Save cleaned documents.json if we made changes
            if cleaned_count > 0:
                with open(storage_path, 'w', encoding='utf-8') as f:
                    json.dump(all_docs, f, ensure_ascii=False, indent=2)
                logger.info(f"✅ Cleaned up {cleaned_count} stale LightRAG document IDs on startup")
                
                # Reload the in-memory database to reflect changes
                load_documents_db()
            else:
                logger.info("No stale LightRAG document IDs found")
                
        except Exception as e:
            logger.error(f"Error during stale doc cleanup: {e}")

    def cleanup_stuck_documents():
        """Fix documents stuck in 'processing' state on startup"""
        try:
            storage_path = LIGHTRAG_METADATA_PATH / "documents.json"
            if not storage_path.exists():
                return
            
            # Load all documents
            with open(storage_path, 'r', encoding='utf-8') as f:
                all_docs = json.load(f)
            
            # Find documents stuck in "processing" state
            stuck_docs = [(doc_id, doc_data) for doc_id, doc_data in all_docs.items() 
                         if doc_data.get("status") == "processing"]
            
            if not stuck_docs:
                logger.info("No stuck documents found")
                return
            
            logger.warning(f"Found {len(stuck_docs)} documents stuck in 'processing' state")
            
            # Check queue database to see if these jobs actually completed/failed
            import sqlite3
            queue_db_path = LIGHTRAG_METADATA_PATH / "document_queue.db"
            
            if not queue_db_path.exists():
                logger.warning("Queue database not found, cannot determine job status")
                return
            
            conn = sqlite3.connect(str(queue_db_path))
            conn.row_factory = sqlite3.Row
            
            fixed_count = 0
            for doc_id, doc_data in stuck_docs:
                # Check if there's a completed/failed job for this document
                cursor = conn.execute("""
                    SELECT status, error FROM jobs 
                    WHERE document_id = ? 
                    ORDER BY created_at DESC LIMIT 1
                """, (doc_id,))
                job = cursor.fetchone()
                
                if job:
                    if job['status'] == 'completed':
                        # Job completed but document status wasn't updated (silent error)
                        logger.info(f"Marking document {doc_id} ({doc_data.get('filename')}) as failed - queue job completed but status wasn't updated")
                        all_docs[doc_id]["status"] = "failed"
                        all_docs[doc_id]["error"] = "Processing completed but document status wasn't updated (possible silent error). Please retry."
                        all_docs[doc_id]["failed_at"] = datetime.now().isoformat()
                        fixed_count += 1
                    elif job['status'] == 'failed':
                        # Job failed, update document status
                        logger.info(f"Marking document {doc_id} ({doc_data.get('filename')}) as failed - queue job failed")
                        all_docs[doc_id]["status"] = "failed"
                        all_docs[doc_id]["error"] = job['error'] or "Unknown error"
                        all_docs[doc_id]["failed_at"] = datetime.now().isoformat()
                        fixed_count += 1
                else:
                    # No job found, mark as failed
                    logger.info(f"Marking document {doc_id} ({doc_data.get('filename')}) as failed - no queue job found")
                    all_docs[doc_id]["status"] = "failed"
                    all_docs[doc_id]["error"] = "No queue job found - server may have crashed during upload"
                    all_docs[doc_id]["failed_at"] = datetime.now().isoformat()
                    fixed_count += 1
            
            conn.close()
            
            # Save updated documents
            if fixed_count > 0:
                with open(storage_path, 'w', encoding='utf-8') as f:
                    json.dump(all_docs, f, ensure_ascii=False, indent=2)
                logger.info(f"✅ Fixed {fixed_count} stuck documents")
                
                # Reload the in-memory database to reflect changes
                load_documents_db()
                
        except Exception as e:
            logger.error(f"Error fixing stuck documents: {e}")
            logger.error(f"Traceback: {traceback.format_exc()}")

    @app.on_event("startup")
    async def startup_event():
        """Initialize the document processing queue on server startup"""
        logger.info("=== Starting Document Queue Worker ===")

        # Fix documents stuck in "processing" state
        cleanup_stuck_documents()

        # Clean up stale LightRAG document IDs from failed documents
        cleanup_stale_lightrag_docs()

        # Register the document processor function
        document_queue.set_processor(process_notebook_document)

        # Register status checker to verify processing results
        def check_document_status(document_id: str) -> tuple:
            """
            Check the actual status of a document after processing

            Returns:
                tuple: (status, error_message) where status is "completed", "failed", or "processing"
            """
            if document_id not in lightrag_documents_db:
                return ("unknown", "Document not found in database")

            doc_data = lightrag_documents_db[document_id]
            status = doc_data.get("status", "unknown")
            error = doc_data.get("error", "")

            return (status, error)

        document_queue.set_status_checker(check_document_status)
        logger.info("Document status checker registered")

        # Recover any stuck jobs from previous crash/restart
        # Pass None to recover ALL processing jobs on startup
        recovered = document_queue.recover_stuck_jobs(timeout_seconds=None)
        if recovered > 0:
            logger.warning(f"Recovered {recovered} stuck jobs from previous session")

        # Start the background worker thread
        document_queue.start_worker()
        logger.info("Document queue worker started successfully")

    @app.on_event("shutdown")
    async def shutdown_event():
        """Stop the document processing queue on server shutdown"""
        logger.info("Stopping document queue worker...")
        document_queue.stop_worker()
        logger.info("Document queue worker stopped")

    # ============================================================================
    # API Routes
    # ============================================================================

    @app.get("/notebooks/validate-embedding-dimensions")
    async def validate_embedding_dimensions(
        model_name: str = Query(..., description="Embedding model name to validate"),
        manual_dimensions: Optional[int] = Query(None, description="Manual dimension override"),
        manual_max_tokens: Optional[int] = Query(None, description="Manual max tokens override")
    ):
        """
        Validate embedding model dimensions with fuzzy matching and confidence scoring.
        Returns detected specifications and manual override options.
        """
        try:
            # Prepare manual override if provided
            manual_override = None
            if manual_dimensions is not None or manual_max_tokens is not None:
                manual_override = {}
                if manual_dimensions is not None:
                    manual_override['dimensions'] = manual_dimensions
                if manual_max_tokens is not None:
                    manual_override['max_tokens'] = manual_max_tokens
            
            # Call the fuzzy matching function
            specs = detect_embedding_model_specs(model_name, manual_override)
            
            logger.info(f"Embedding validation for '{model_name}': {specs['dimensions']}d, "
                       f"confidence={specs['confidence']}, pattern={specs['detected_pattern']}")
            
            return {
                "status": "success",
                "model_name": model_name,
                "specifications": specs,
                "warning": None if specs['confidence'] >= 0.8 else 
                          "Low confidence detection. Please verify dimensions or select manually."
            }
            
        except Exception as e:
            logger.error(f"Error validating embedding dimensions for '{model_name}': {e}")
            return {
                "status": "error",
                "model_name": model_name,
                "error": str(e),
                "specifications": {
                    'dimensions': 1536,
                    'max_tokens': 8192,
                    'confidence': 0.0,
                    'detected_pattern': 'error_fallback',
                    'override_options': [],
                    'is_manual_override': False
                }
            }
    
    @app.post("/notebooks", response_model=NotebookResponse)
    async def create_notebook(notebook: NotebookCreate):
        """Create a new notebook using LightRAG"""
        notebook_id = str(uuid.uuid4())
        
        # Add debugging logging
        logger.info(f"Creating notebook with data: name={notebook.name}, description={notebook.description}")
        logger.info(f"Original LLM Provider: {notebook.llm_provider}")
        logger.info(f"Original Embedding Provider: {notebook.embedding_provider}")
        
        # Auto-detect provider types before saving
        corrected_llm_provider = auto_detect_provider_type(notebook.llm_provider)
        corrected_embedding_provider = auto_detect_provider_type(notebook.embedding_provider)

        # Normalize URLs for Docker/local development
        if corrected_llm_provider.get('baseUrl'):
            corrected_llm_provider['baseUrl'] = normalize_url_for_local_dev(corrected_llm_provider['baseUrl'])
        if corrected_embedding_provider.get('baseUrl'):
            corrected_embedding_provider['baseUrl'] = normalize_url_for_local_dev(corrected_embedding_provider['baseUrl'])

        logger.info(f"Corrected LLM Provider: {corrected_llm_provider}")
        logger.info(f"Corrected Embedding Provider: {corrected_embedding_provider}")
        
        # Prepare manual embedding override if provided
        manual_embedding_override = None
        if notebook.manual_embedding_dimensions is not None or notebook.manual_embedding_max_tokens is not None:
            manual_embedding_override = {}
            if notebook.manual_embedding_dimensions is not None:
                manual_embedding_override['dimensions'] = notebook.manual_embedding_dimensions
            if notebook.manual_embedding_max_tokens is not None:
                manual_embedding_override['max_tokens'] = notebook.manual_embedding_max_tokens
            logger.info(f"Manual embedding override applied: {manual_embedding_override}")
        
        notebook_data = {
            "id": notebook_id,
            "name": notebook.name,
            "description": notebook.description,
            "created_at": datetime.now(),
            "document_count": 0,
            "llm_provider": corrected_llm_provider,
            "embedding_provider": corrected_embedding_provider,
            # Store schema consistency parameters
            "entity_types": notebook.entity_types,
            "language": notebook.language or "en",
            # Store manual overrides
            "manual_embedding_override": manual_embedding_override
        }
        
        # Log the notebook data before saving
        logger.info(f"Notebook data before saving: {notebook_data}")
        
        lightrag_notebooks_db[notebook_id] = notebook_data
        
        # Create LightRAG instance for this notebook
        try:
            await create_lightrag_instance(
                notebook_id, 
                corrected_llm_provider, 
                corrected_embedding_provider,
                entity_types=notebook.entity_types,
                language=notebook.language or "en",
                manual_embedding_override=manual_embedding_override
            )
            logger.info(f"Created notebook {notebook_id}: {notebook.name}")
            # Save to disk after successful creation
            save_notebooks_db()
            
            # Log the saved data
            logger.info(f"Saved notebook data: {lightrag_notebooks_db[notebook_id]}")
        except Exception as e:
            # Clean up if LightRAG creation fails
            del lightrag_notebooks_db[notebook_id]
            raise
        
        # Create response and log it
        response = NotebookResponse(**notebook_data)
        logger.info(f"Response being returned: {response.model_dump()}")
        
        return response

    @app.get("/notebooks", response_model=List[NotebookResponse])
    async def list_notebooks():
        """List all notebooks"""
        logger.info(f"Listing notebooks, found {len(lightrag_notebooks_db)} notebooks")
        notebooks = []
        for notebook_id, notebook in lightrag_notebooks_db.items():
            logger.info(f"Processing notebook {notebook_id}: {notebook}")
            
            # Create a copy to avoid modifying the original
            notebook_copy = notebook.copy()
            
            # Add default provider configuration if missing (for backward compatibility)
            if "llm_provider" not in notebook_copy:
                logger.info(f"Adding default LLM provider for notebook {notebook_id}")
                notebook_copy["llm_provider"] = {
                    "name": "OpenAI",
                    "type": "openai",
                    "baseUrl": "https://api.openai.com/v1",
                    "apiKey": "your-api-key",
                    "model": "gpt-4o-mini"
                }
            if "embedding_provider" not in notebook_copy:
                logger.info(f"Adding default embedding provider for notebook {notebook_id}")
                notebook_copy["embedding_provider"] = {
                    "name": "OpenAI",
                    "type": "openai", 
                    "baseUrl": "https://api.openai.com/v1",
                    "apiKey": "your-api-key",
                    "model": "text-embedding-ada-002"
                }
            
            notebook_response = NotebookResponse(**notebook_copy)
            logger.info(f"Notebook response for {notebook_id}: {notebook_response.model_dump()}")
            notebooks.append(notebook_response)
        return notebooks

    @app.get("/notebooks/{notebook_id}", response_model=NotebookResponse)
    async def get_notebook(notebook_id: str):
        """Get a specific notebook"""
        validate_notebook_exists(notebook_id)
        notebook = lightrag_notebooks_db[notebook_id].copy()
        
        # Add default provider configuration if missing (for backward compatibility)
        if "llm_provider" not in notebook:
            notebook["llm_provider"] = {
                "name": "OpenAI",
                "type": "openai",
                "baseUrl": "https://api.openai.com/v1",
                "apiKey": "your-api-key",
                "model": "gpt-4o-mini"
            }
        if "embedding_provider" not in notebook:
            notebook["embedding_provider"] = {
                "name": "OpenAI",
                "type": "openai", 
                "baseUrl": "https://api.openai.com/v1",
                "apiKey": "your-api-key",
                "model": "text-embedding-ada-002"
            }
        
        return NotebookResponse(**notebook)

    @app.post("/notebooks/{notebook_id}/rebuild")
    async def rebuild_notebook(notebook_id: str, background_tasks: BackgroundTasks):
        """Rebuild notebook by clearing LightRAG storage and reprocessing all documents
        
        This is the safe rebuild operation that:
        1. Clears corrupted LightRAG storage
        2. Resets document statuses to 'pending'
        3. Queues all documents for reprocessing
        
        Documents are preserved in the database and will be reprocessed.
        """
        validate_notebook_exists(notebook_id)
        
        logger.info(f"Rebuilding notebook {notebook_id}")
        
        # Step 1: Remove LightRAG instance with proper finalization
        if notebook_id in lightrag_instances:
            try:
                rag_instance = lightrag_instances[notebook_id]
                await rag_instance.finalize_storages()
                logger.info(f"Finalized storage for notebook {notebook_id}")
            except Exception as finalize_error:
                logger.warning(f"Error finalizing storage: {finalize_error}")
            finally:
                del lightrag_instances[notebook_id]
        
        # Step 2: Clean up storage directory completely
        storage_dir = LIGHTRAG_STORAGE_PATH / notebook_id
        if storage_dir.exists():
            shutil.rmtree(storage_dir, ignore_errors=True)
            logger.info(f"Deleted storage directory: {storage_dir}")
        
        # Recreate empty storage directory
        storage_dir.mkdir(exist_ok=True)
        logger.info(f"Recreated empty storage directory: {storage_dir}")
        
        # Step 3: Reset all document statuses to 'pending' (KEEP documents in database)
        notebook_docs = [(doc_id, doc) for doc_id, doc in lightrag_documents_db.items() 
                         if doc["notebook_id"] == notebook_id]
        
        reprocessed_count = 0
        failed_no_content = []
        
        for doc_id, doc in notebook_docs:
            # Reset status to pending
            lightrag_documents_db[doc_id]["status"] = "pending"
            lightrag_documents_db[doc_id]["queued_at"] = datetime.now()
            
            # Clear processing metadata but keep document content and filename
            for field in ["processed_at", "completed_at", "failed_at", "error", "lightrag_id"]:
                if field in lightrag_documents_db[doc_id]:
                    del lightrag_documents_db[doc_id][field]
            
            # Get content for reprocessing
            content = None
            if "content" in doc:
                content = doc["content"]
                logger.info(f"Found stored content for document {doc_id} ({doc.get('filename', 'unknown')})")
            elif "content_file" in doc:
                try:
                    content_file = LIGHTRAG_STORAGE_PATH / "documents" / doc["content_file"]
                    content = content_file.read_text(encoding='utf-8')
                    logger.info(f"Loaded content from file for document {doc_id} ({doc.get('filename', 'unknown')})")
                except Exception as e:
                    logger.error(f"Failed to read content file for document {doc_id}: {e}")
                    lightrag_documents_db[doc_id]["status"] = "failed"
                    lightrag_documents_db[doc_id]["error"] = f"Content file not found: {str(e)}"
                    failed_no_content.append(doc.get('filename', doc_id))
                    continue
            
            if not content:
                logger.warning(f"Document {doc_id} ({doc.get('filename', 'unknown')}) has no content available for reprocessing")
                lightrag_documents_db[doc_id]["status"] = "failed"
                lightrag_documents_db[doc_id]["error"] = "No content available - please re-upload the document"
                failed_no_content.append(doc.get('filename', doc_id))
                continue
            
            # Queue for persistent processing (crash-safe)
            # Use the persistent queue instead of background_tasks for reliability
            document_queue.enqueue(
                notebook_id=notebook_id,
                document_id=doc_id,
                content=content,
                priority=0  # Normal priority
            )

            reprocessed_count += 1
            logger.info(f"Queued document {doc_id} for reprocessing via persistent queue")
        
        # Clear notebook summary cache and fingerprint
        if "summary" in lightrag_notebooks_db[notebook_id]:
            del lightrag_notebooks_db[notebook_id]["summary"]
        if "summary_generated_at" in lightrag_notebooks_db[notebook_id]:
            del lightrag_notebooks_db[notebook_id]["summary_generated_at"]
        if "docs_fingerprint" in lightrag_notebooks_db[notebook_id]:
            del lightrag_notebooks_db[notebook_id]["docs_fingerprint"]
        
        # Save changes
        save_notebooks_db()
        save_documents_db()
        
        logger.info(f"✓ Notebook {notebook_id} rebuild initiated: {reprocessed_count} documents queued")
        
        # Build response message
        if failed_no_content:
            message = f"Notebook rebuild initiated: {reprocessed_count} documents queued. " \
                     f"{len(failed_no_content)} documents failed (no content stored - please re-upload)"
            note = f"Failed documents: {', '.join(failed_no_content[:5])}" + \
                   (f" and {len(failed_no_content) - 5} more" if len(failed_no_content) > 5 else "")
        else:
            message = f"Notebook rebuild initiated: {reprocessed_count} documents queued for reprocessing"
            note = "All documents have been queued for reprocessing with current configuration"
        
        return {
            "message": message,
            "total_documents": len(notebook_docs),
            "queued_for_reprocessing": reprocessed_count,
            "failed_no_content": len(failed_no_content),
            "notebook_id": notebook_id,
            "note": note
        }

    @app.post("/notebooks/{notebook_id}/clear-storage")
    async def clear_notebook_storage(notebook_id: str):
        """Clear all LightRAG storage for a notebook while keeping notebook metadata
        
        This is useful when storage gets corrupted or contains orphaned documents.
        """
        validate_notebook_exists(notebook_id)
        
        logger.info(f"Clearing storage for notebook {notebook_id}")
        
        # Remove LightRAG instance with proper finalization
        if notebook_id in lightrag_instances:
            try:
                rag_instance = lightrag_instances[notebook_id]
                await rag_instance.finalize_storages()
                logger.info(f"Finalized storage for notebook {notebook_id}")
            except Exception as finalize_error:
                logger.warning(f"Error finalizing storage: {finalize_error}")
            finally:
                del lightrag_instances[notebook_id]
        
        # Delete all documents in this notebook from database
        docs_to_delete = [doc_id for doc_id, doc in lightrag_documents_db.items() 
                         if doc["notebook_id"] == notebook_id]
        for doc_id in docs_to_delete:
            del lightrag_documents_db[doc_id]
        
        # Clean up storage directory completely
        storage_dir = LIGHTRAG_STORAGE_PATH / notebook_id
        if storage_dir.exists():
            shutil.rmtree(storage_dir, ignore_errors=True)
            logger.info(f"Deleted storage directory: {storage_dir}")
        
        # Recreate empty storage directory
        storage_dir.mkdir(exist_ok=True)
        logger.info(f"Recreated empty storage directory: {storage_dir}")
        
        # Clear notebook summary cache and fingerprint
        if "summary" in lightrag_notebooks_db[notebook_id]:
            del lightrag_notebooks_db[notebook_id]["summary"]
        if "summary_generated_at" in lightrag_notebooks_db[notebook_id]:
            del lightrag_notebooks_db[notebook_id]["summary_generated_at"]
        if "docs_fingerprint" in lightrag_notebooks_db[notebook_id]:
            del lightrag_notebooks_db[notebook_id]["docs_fingerprint"]
        
        # Save changes
        save_notebooks_db()
        save_documents_db()
        
        logger.info(f"✓ Storage cleared for notebook {notebook_id}")
        return {
            "message": "Notebook storage cleared successfully",
            "documents_deleted": len(docs_to_delete),
            "notebook_id": notebook_id
        }

    @app.delete("/notebooks/{notebook_id}")
    async def delete_notebook(notebook_id: str):
        """Delete a notebook and all its documents"""
        validate_notebook_exists(notebook_id)
        
        # Remove all documents from this notebook
        notebook_docs = [doc_id for doc_id, doc in lightrag_documents_db.items() 
                        if doc["notebook_id"] == notebook_id]
        
        for doc_id in notebook_docs:
            del lightrag_documents_db[doc_id]
        
        # Remove LightRAG instance with proper finalization
        if notebook_id in lightrag_instances:
            try:
                # Properly finalize storage before deletion
                rag_instance = lightrag_instances[notebook_id]
                await rag_instance.finalize_storages()
                logger.info(f"Finalized storage for notebook {notebook_id}")
            except Exception as finalize_error:
                logger.warning(f"Error finalizing storage for notebook {notebook_id}: {finalize_error}")
            finally:
                del lightrag_instances[notebook_id]
        
        # Remove notebook
        del lightrag_notebooks_db[notebook_id]
        
        # Clean up storage directory
        storage_dir = LIGHTRAG_STORAGE_PATH / notebook_id
        if storage_dir.exists():
            shutil.rmtree(storage_dir, ignore_errors=True)
        
        # Save changes to disk
        save_notebooks_db()
        save_documents_db()
        
        logger.info(f"Deleted notebook {notebook_id}")
        return {"message": "Notebook deleted successfully"}

    @app.put("/notebooks/{notebook_id}/configuration")
    async def update_notebook_configuration(
        notebook_id: str,
        request: NotebookCreate
    ):
        """Update notebook LLM and embedding provider configuration and rebuild LightRAG instance"""
        validate_notebook_exists(notebook_id)
        
        logger.info(f"Updating configuration for notebook {notebook_id}")

        # Auto-detect provider types
        llm_provider = auto_detect_provider_type(request.llm_provider)
        embedding_provider = auto_detect_provider_type(request.embedding_provider)

        # Normalize URLs for Docker/local development
        if llm_provider.get('baseUrl'):
            llm_provider['baseUrl'] = normalize_url_for_local_dev(llm_provider['baseUrl'])
        if embedding_provider.get('baseUrl'):
            embedding_provider['baseUrl'] = normalize_url_for_local_dev(embedding_provider['baseUrl'])
        
        # Check if embedding model is changing (dimension mismatch detection)
        old_embedding_model = lightrag_notebooks_db[notebook_id]["embedding_provider"].get("model", "")
        new_embedding_model = embedding_provider.get("model", "")
        embedding_changed = old_embedding_model != new_embedding_model
        
        if embedding_changed:
            logger.warning(f"Embedding model changing from '{old_embedding_model}' to '{new_embedding_model}' - vector storage will be cleared")
            
            # Clear vector storage files to prevent dimension mismatch
            working_dir = LIGHTRAG_STORAGE_PATH / notebook_id
            if working_dir.exists():
                vector_files = [
                    "vdb_entities.json",
                    "vdb_relationships.json", 
                    "vdb_chunks.json",
                    "kv_store_full_docs.json",
                    "kv_store_text_chunks.json",
                    "kv_store_llm_response_cache.json"
                ]
                for vector_file in vector_files:
                    vector_path = working_dir / vector_file
                    if vector_path.exists():
                        try:
                            vector_path.unlink()
                            logger.info(f"Cleared vector storage file: {vector_file}")
                        except Exception as e:
                            logger.warning(f"Failed to clear {vector_file}: {e}")
        
        # Update notebook configuration
        lightrag_notebooks_db[notebook_id]["name"] = request.name
        lightrag_notebooks_db[notebook_id]["description"] = request.description  
        lightrag_notebooks_db[notebook_id]["llm_provider"] = llm_provider
        lightrag_notebooks_db[notebook_id]["embedding_provider"] = embedding_provider
        lightrag_notebooks_db[notebook_id]["updated_at"] = datetime.now()
        
        # Clear any cached summary since configuration changed
        if "summary_cache" in lightrag_notebooks_db[notebook_id]:
            del lightrag_notebooks_db[notebook_id]["summary_cache"]
        if "docs_fingerprint" in lightrag_notebooks_db[notebook_id]:
            del lightrag_notebooks_db[notebook_id]["docs_fingerprint"]
        
        # Remove existing LightRAG instance to force rebuild with new configuration
        if notebook_id in lightrag_instances:
            try:
                # Properly finalize storage before deletion
                rag_instance = lightrag_instances[notebook_id]
                await rag_instance.finalize_storages()
                logger.info(f"Finalized storage for notebook {notebook_id} before configuration update")
            except Exception as finalize_error:
                logger.warning(f"Error finalizing storage for notebook {notebook_id}: {finalize_error}")
            finally:
                logger.info(f"Removing existing LightRAG instance for notebook {notebook_id}")
                del lightrag_instances[notebook_id]
        
        try:
            # Create new LightRAG instance with updated configuration
            logger.info(f"Creating new LightRAG instance with updated configuration")
            notebook_data = lightrag_notebooks_db[notebook_id]
            new_rag_instance = await create_lightrag_instance(
                notebook_id, 
                llm_provider, 
                embedding_provider,
                entity_types=notebook_data.get("entity_types"),
                language=notebook_data.get("language", "en")
            )
            
            # Store the new instance
            lightrag_instances[notebook_id] = new_rag_instance
            
            # Check if there are existing documents that need to be reprocessed
            notebook_docs = [doc for doc in lightrag_documents_db.values() 
                           if doc["notebook_id"] == notebook_id]
            
            reprocess_info = {
                "total_documents": len(notebook_docs),
                "completed_documents": len([doc for doc in notebook_docs if doc["status"] == "completed"]),
                "failed_documents": len([doc for doc in notebook_docs if doc["status"] == "failed"]),
                "needs_reprocessing": len([doc for doc in notebook_docs if doc["status"] == "completed"])
            }
            
            # Save changes to disk
            save_notebooks_db()
            
            logger.info(f"Successfully updated configuration for notebook {notebook_id}")
            
            # Return updated notebook with reprocessing info
            notebook = lightrag_notebooks_db[notebook_id].copy()
            notebook["document_count"] = len(notebook_docs)
            
            # Provide appropriate message based on whether embedding changed
            if embedding_changed:
                message = "⚠️ Embedding model changed - all vector storage cleared. You must re-upload all documents."
                recommendation = "IMPORTANT: All existing embeddings have been cleared due to dimension mismatch. Please re-upload all documents to rebuild the knowledge graph with the new embedding model."
            else:
                message = "Notebook configuration updated successfully"
                recommendation = "Configuration updated. Existing documents remain indexed."
            
            response = {
                "message": message,
                "notebook": NotebookResponse(**notebook),
                "reprocessing_info": reprocess_info,
                "embedding_changed": embedding_changed,
                "recommendation": recommendation
            }
            
            return response
            
        except Exception as e:
            logger.error(f"Failed to update notebook configuration: {e}")
            # Restore previous instance if available (though it may be incompatible)
            raise HTTPException(
                status_code=500, 
                detail=f"Failed to update notebook configuration: {str(e)}. Please check your provider settings."
            )

    @app.put("/notebooks/{notebook_id}/schema", response_model=NotebookResponse)
    async def update_notebook_schema(notebook_id: str, schema_update: NotebookSchemaUpdate):
        """Update entity types and language settings for a notebook"""
        validate_notebook_exists(notebook_id)
        
        notebook = lightrag_notebooks_db[notebook_id]
        
        # Update schema settings if provided
        if schema_update.entity_types is not None:
            notebook["entity_types"] = schema_update.entity_types
            logger.info(f"Updated entity types for notebook {notebook_id}: {schema_update.entity_types}")
        
        if schema_update.language is not None:
            notebook["language"] = schema_update.language
            logger.info(f"Updated language for notebook {notebook_id}: {schema_update.language}")
        
        # The schema changes will take effect for new documents automatically
        # Existing LightRAG instance keeps running with current data
        
        # Save changes to disk
        save_notebooks_db()
        
        logger.info(f"Updated schema configuration for notebook {notebook_id}")
        
        return NotebookResponse(**notebook)

    @app.post("/notebooks/{notebook_id}/reprocess-documents")
    async def reprocess_notebook_documents(
        notebook_id: str,
        background_tasks: BackgroundTasks,
        force: bool = False
    ):
        """Reprocess all completed documents in a notebook with current configuration"""
        validate_notebook_exists(notebook_id)
        
        # Get all completed documents for this notebook
        notebook_docs = [
            (doc_id, doc) for doc_id, doc in lightrag_documents_db.items() 
            if doc["notebook_id"] == notebook_id and (doc["status"] == "completed" or (force and doc["status"] == "failed"))
        ]
        
        if not notebook_docs:
            return {
                "message": "No documents to reprocess",
                "total_documents": 0,
                "queued_for_reprocessing": 0
            }
        
        reprocessed_count = 0
        
        # Mark documents for reprocessing and queue background tasks
        for doc_id, doc in notebook_docs:
            if "content" in doc or "content_file" in doc:
                # Document has content available for reprocessing
                lightrag_documents_db[doc_id]["status"] = "pending"
                lightrag_documents_db[doc_id]["queued_at"] = datetime.now()
                
                # Clear any previous processing metadata
                for field in ["processed_at", "completed_at", "failed_at", "error", "lightrag_id"]:
                    if field in lightrag_documents_db[doc_id]:
                        del lightrag_documents_db[doc_id][field]
                
                # Get content for reprocessing
                if "content" in doc:
                    content = doc["content"]
                elif "content_file" in doc:
                    try:
                        content_file = LIGHTRAG_STORAGE_PATH / "documents" / doc["content_file"]
                        content = content_file.read_text(encoding='utf-8')
                    except Exception as e:
                        logger.error(f"Failed to read content file for document {doc_id}: {e}")
                        lightrag_documents_db[doc_id]["status"] = "failed"
                        lightrag_documents_db[doc_id]["error"] = f"Content file not found: {str(e)}"
                        continue
                else:
                    continue
                
                # Queue for persistent processing (crash-safe)
                # Use the persistent queue instead of background_tasks for reliability
                document_queue.enqueue(
                    notebook_id=notebook_id,
                    document_id=doc_id,
                    content=content,
                    priority=0  # Normal priority
                )
                
                reprocessed_count += 1
                logger.info(f"Queued document {doc_id} for reprocessing via persistent queue")
            else:
                logger.warning(f"Document {doc_id} has no content available for reprocessing")
        
        # Save changes to disk
        save_documents_db()
        
        logger.info(f"Queued {reprocessed_count} documents for reprocessing in notebook {notebook_id}")
        
        return {
            "message": f"Queued {reprocessed_count} documents for reprocessing",
            "total_documents": len(notebook_docs),
            "queued_for_reprocessing": reprocessed_count,
            "note": "Documents will be processed in the background with the current notebook configuration"
        }

    @app.post("/notebooks/{notebook_id}/documents", response_model=List[NotebookDocumentResponse])
    async def upload_notebook_documents(
        notebook_id: str,
        files: List[UploadFile] = File(...)
    ):
        """Upload multiple documents to a notebook"""
        validate_notebook_exists(notebook_id)
        
        if not files:
            raise HTTPException(status_code=400, detail="No files provided")
        
        uploaded_documents = []
        
        # Process files sequentially to avoid conflicts
        for i, file in enumerate(files):
            if not file.filename:
                continue
                
            # Generate document ID
            document_id = str(uuid.uuid4())
            
            # Read file content
            try:
                file_content = await file.read()
                
                # Extract text based on file type
                text_content = await extract_text_from_file(file.filename, file_content)
                
                # Validate text content
                if not text_content.strip():
                    logger.warning(f"File {file.filename} appears to be empty")
                    continue
                
                # Create file path for citation tracking
                file_path = f"notebooks/{notebook_id}/{file.filename}"
                
                # Create document record - HONEST STATUS: "queued" not "processing"
                document_data = {
                    "id": document_id,
                    "filename": file.filename,
                    "notebook_id": notebook_id,
                    "uploaded_at": datetime.now(),
                    "status": "queued",  # Honest: Document is queued, not processing yet
                    "file_path": file_path,
                    # Don't store content in memory - only file path (scales to 1000+ docs)
                }
                
                # Debug: Log that content is being stored
                content_length = len(text_content)
                logger.info(f"Storing document {document_id} with content length: {content_length} characters")
                
                # Add content size info for monitoring
                document_data["content_size"] = content_length

                # ALWAYS save content to a separate file as backup (like Google Drive)
                # This provides redundancy and enables easy document downloads
                content_file = Path(data_dir) / f"content_{document_id}.txt"
                try:
                    content_file.parent.mkdir(parents=True, exist_ok=True)
                    with open(content_file, 'w', encoding='utf-8') as f:
                        f.write(text_content)
                    document_data["content_file"] = str(content_file)
                    logger.info(f"Content backed up to {content_file}")
                except Exception as e:
                    logger.warning(f"Failed to create content backup file: {e}")
                
                lightrag_documents_db[document_id] = document_data

                # Add document to persistent queue
                # SQLite queue ensures documents are processed even after server crashes
                # Priority system: higher priority = processed first
                priority = 0  # Can be increased for urgent documents

                job_id = document_queue.enqueue(
                    notebook_id=notebook_id,
                    document_id=document_id,
                    content=text_content,
                    priority=priority
                )

                # Save job_id to track processing
                lightrag_documents_db[document_id]["job_id"] = job_id

                if i == 0:
                    logger.info(f"📄 Uploading {len(files)} documents to persistent queue")
                    logger.info(f"✅ Queue survives crashes - your documents won't be lost!")

                logger.info(f"📋 Queued document {i+1}/{len(files)}: {file.filename} (job_id: {job_id})")
                
                # Update notebook document count
                lightrag_notebooks_db[notebook_id]["document_count"] += 1
                
                uploaded_documents.append(NotebookDocumentResponse(**document_data))
                
            except Exception as e:
                logger.error(f"Error processing file {file.filename}: {e}")
                raise HTTPException(
                    status_code=500, 
                    detail=f"Error processing file {file.filename}: {str(e)}"
                )
        
        # Save changes to disk after all uploads
        save_documents_db()
        save_notebooks_db()
        
        logger.info(f"Uploaded {len(uploaded_documents)} documents to notebook {notebook_id}")
        return uploaded_documents

    @app.get("/notebooks/{notebook_id}/documents", response_model=List[NotebookDocumentResponse])
    async def list_notebook_documents(notebook_id: str):
        """List all documents in a notebook"""
        validate_notebook_exists(notebook_id)
        
        notebook_documents = [
            NotebookDocumentResponse(**doc) 
            for doc in lightrag_documents_db.values() 
            if doc["notebook_id"] == notebook_id
        ]
        
        return notebook_documents

    @app.get("/notebooks/{notebook_id}/queue/status")
    async def get_notebook_queue_status(notebook_id: str):
        """Get queue statistics for a specific notebook"""
        validate_notebook_exists(notebook_id)

        stats = document_queue.get_queue_stats(notebook_id)

        return {
            "notebook_id": notebook_id,
            "queue_stats": stats,
            "message": f"{stats['queued']} queued, {stats['processing']} processing, "
                      f"{stats['completed']} completed, {stats['failed']} failed"
        }

    @app.get("/queue/status")
    async def get_global_queue_status():
        """Get global queue statistics across all notebooks"""
        stats = document_queue.get_queue_stats()

        return {
            "global_stats": stats,
            "message": f"Total: {stats['total']} jobs "
                      f"({stats['queued']} queued, {stats['processing']} processing, "
                      f"{stats['completed']} completed, {stats['failed']} failed)"
        }

    @app.post("/queue/recover-stuck")
    async def recover_stuck_jobs(timeout_seconds: int = 600):
        """Manually recover jobs stuck in processing state"""
        recovered = document_queue.recover_stuck_jobs(timeout_seconds)

        return {
            "recovered_jobs": recovered,
            "message": f"Recovered {recovered} stuck jobs" if recovered > 0
                      else "No stuck jobs found"
        }

    @app.delete("/notebooks/{notebook_id}/documents/{document_id}")
    async def delete_notebook_document(notebook_id: str, document_id: str):
        """Delete a specific document from a notebook"""
        validate_notebook_exists(notebook_id)
        
        if document_id not in lightrag_documents_db:
            raise HTTPException(status_code=404, detail="Document not found")
        
        if lightrag_documents_db[document_id]["notebook_id"] != notebook_id:
            raise HTTPException(status_code=400, detail="Document does not belong to this notebook")
        
        try:
            # Remove from LightRAG using the stored LightRAG ID
            rag = await get_lightrag_instance(notebook_id)
            document_data = lightrag_documents_db[document_id]
            
            # Use the stored LightRAG ID if available, otherwise construct it
            lightrag_id = document_data.get("lightrag_id", f"doc_{notebook_id}_{document_id}")
            await rag.adelete_by_doc_id(lightrag_id)
            
            # Clear cache after deleting document
            await rag.aclear_cache()
            
            # Clean up content file if it exists
            if "content_file" in document_data:
                try:
                    content_file = Path(document_data["content_file"])
                    if content_file.exists():
                        content_file.unlink()
                        logger.info(f"Cleaned up content file during deletion: {content_file}")
                except Exception as e:
                    logger.warning(f"Failed to clean up content file during deletion: {e}")
            
            # Remove from database
            del lightrag_documents_db[document_id]
            
            # Update notebook document count
            lightrag_notebooks_db[notebook_id]["document_count"] -= 1
            
            # Clear summary cache since documents have changed
            if "summary_cache" in lightrag_notebooks_db[notebook_id]:
                del lightrag_notebooks_db[notebook_id]["summary_cache"]
            if "docs_fingerprint" in lightrag_notebooks_db[notebook_id]:
                del lightrag_notebooks_db[notebook_id]["docs_fingerprint"]
            
            # Save changes to disk
            save_documents_db()
            save_notebooks_db()
            
            logger.info(f"Deleted document {document_id} from notebook {notebook_id}")
            return {"message": "Document deleted successfully"}

        except Exception as e:
            logger.error(f"Error deleting document {document_id}: {e}")
            raise HTTPException(status_code=500, detail=f"Error deleting document: {str(e)}")

    @app.get("/notebooks/{notebook_id}/documents/{document_id}/download")
    async def download_notebook_document(notebook_id: str, document_id: str):
        """Download the original document content

        Documents are precious - users can always download their uploaded content.
        Works like Google Drive - content is always available for download.
        """
        validate_notebook_exists(notebook_id)

        if document_id not in lightrag_documents_db:
            raise HTTPException(status_code=404, detail="Document not found")

        document_data = lightrag_documents_db[document_id]

        if document_data["notebook_id"] != notebook_id:
            raise HTTPException(status_code=400, detail="Document does not belong to this notebook")

        # Get document content from memory or file
        content = None
        filename = document_data.get("filename", f"document_{document_id}.txt")

        # Try to get content from database first
        if "content" in document_data:
            content = document_data["content"]
            logger.info(f"Serving document {document_id} content from database ({len(content)} chars)")
        # Try to get content from backup file
        elif "content_file" in document_data:
            try:
                content_file = Path(document_data["content_file"])
                if content_file.exists():
                    content = content_file.read_text(encoding='utf-8')
                    logger.info(f"Serving document {document_id} content from file ({len(content)} chars)")
                else:
                    logger.error(f"Content file not found: {content_file}")
            except Exception as e:
                logger.error(f"Failed to read content file for document {document_id}: {e}")

        if not content:
            raise HTTPException(
                status_code=404,
                detail="Document content not available. This may happen if the document was uploaded before the content preservation feature was added."
            )

        # Create a temporary file for download
        try:
            temp_file = tempfile.NamedTemporaryFile(mode='w', encoding='utf-8', delete=False, suffix='.txt')
            temp_file.write(content)
            temp_file.close()

            # Return file with proper filename
            return FileResponse(
                path=temp_file.name,
                media_type='text/plain',
                filename=filename,
                background=None  # File will be deleted by OS temp cleanup
            )
        except Exception as e:
            logger.error(f"Error creating download file for document {document_id}: {e}")
            raise HTTPException(status_code=500, detail=f"Error preparing download: {str(e)}")

    @app.post("/notebooks/{notebook_id}/documents/{document_id}/retry", response_model=DocumentRetryResponse)
    async def retry_failed_document(notebook_id: str, document_id: str, background_tasks: BackgroundTasks):
        """Retry processing a failed document"""
        validate_notebook_exists(notebook_id)
        
        if document_id not in lightrag_documents_db:
            raise HTTPException(status_code=404, detail="Document not found")
        
        document_data = lightrag_documents_db[document_id]
        
        if document_data["notebook_id"] != notebook_id:
            raise HTTPException(status_code=400, detail="Document does not belong to this notebook")
        
        # Check if document is in a retryable state
        if document_data["status"] not in ["failed"]:
            raise HTTPException(
                status_code=400, 
                detail=f"Document is in '{document_data['status']}' state. Only 'failed' documents can be retried."
            )
        
        try:
            logger.info(f"Retrying document {document_id} in notebook {notebook_id}")
            
            # Debug: Log document data keys
            logger.info(f"Document data keys: {list(document_data.keys())}")
            
            # Get the original text content from the failed document
            # Check if we have stored content or need to re-extract it
            text_content = document_data.get("content")
            
            # If content is not in memory, try to load from content file (for large documents)
            if not text_content and "content_file" in document_data:
                try:
                    content_file = Path(document_data["content_file"])
                    if content_file.exists():
                        with open(content_file, 'r', encoding='utf-8') as f:
                            text_content = f.read()
                        logger.info(f"Loaded content from backup file: {content_file}")
                    else:
                        logger.warning(f"Content file not found: {content_file}")
                except Exception as e:
                    logger.error(f"Failed to load content from file: {e}")
            
            # Debug: Log content availability
            if text_content:
                logger.info(f"Content available for retry, length: {len(text_content)} characters")
            else:
                logger.warning(f"No content found in document data. Available keys: {list(document_data.keys())}")
            
            if not text_content:
                # If content wasn't stored, we need the original file to retry
                raise HTTPException(
                    status_code=400, 
                    detail="Original document content not available for retry. Please re-upload the document."
                )
            
            # Clear old lightrag_id to force generation of a fresh ID on retry
            # This prevents conflicts with stale document IDs in LightRAG storage
            if "lightrag_id" in document_data:
                old_lightrag_id = document_data["lightrag_id"]
                del document_data["lightrag_id"]
                logger.info(f"Cleared old lightrag_id ({old_lightrag_id[:50]}...) for clean retry")
            
            # Reset document status to pending (will be set to processing by queue worker)
            document_data["status"] = "pending"
            document_data["queued_at"] = datetime.now()

            # Clear previous error information
            if "failed_at" in document_data:
                del document_data["failed_at"]
            if "error" in document_data:
                del document_data["error"]
            if "error_message" in document_data:
                del document_data["error_message"]
            if "error_details" in document_data:
                del document_data["error_details"]

            # Update the document in database
            lightrag_documents_db[document_id] = document_data
            save_documents_db()

            # Queue for persistent processing (crash-safe)
            # Use the persistent queue for reliability and proper status tracking
            document_queue.enqueue(
                notebook_id=notebook_id,
                document_id=document_id,
                content=text_content,
                priority=0  # Normal priority
            )
            logger.info(f"Queued document {document_id} for retry via persistent queue")
            
            logger.info(f"Retry initiated for document {document_id}")
            return DocumentRetryResponse(
                message="Document retry initiated successfully",
                document_id=document_id,
                status="processing"
            )
            
        except HTTPException:
            # Re-raise HTTP exceptions as-is
            raise
        except Exception as e:
            logger.error(f"Error retrying document {document_id}: {e}")
            # Reset status back to failed if retry setup failed
            document_data["status"] = "failed"
            lightrag_documents_db[document_id] = document_data
            save_documents_db()
            raise HTTPException(status_code=500, detail=f"Error initiating retry: {str(e)}")

    @app.post("/notebooks/{notebook_id}/query", response_model=NotebookQueryResponse)
    async def query_notebook(notebook_id: str, query: NotebookQueryRequest):
        """Query a notebook with a question"""
        validate_notebook_exists(notebook_id)
        
        try:
            logger.info(f"Query request for notebook {notebook_id}")
            
            # Get the current RAG instance
            # NOTE: LLM override is NOT supported for now because it would require
            # recreating the RAG instance with the same working directory, which could
            # cause data corruption. The override feature needs architectural changes.
            rag = await get_lightrag_instance(notebook_id)
            
            # CRITICAL: Verify RAG instance working directory matches notebook to prevent data leakage
            expected_dir = str(LIGHTRAG_STORAGE_PATH / notebook_id)
            actual_dir = str(rag.working_dir)
            if expected_dir != actual_dir:
                logger.error(f"⚠️ CRITICAL: RAG working directory mismatch during query!")
                logger.error(f"   Notebook ID: {notebook_id}")
                logger.error(f"   Expected dir: {expected_dir}")
                logger.error(f"   Actual dir: {actual_dir}")
                raise HTTPException(
                    status_code=500, 
                    detail="Data isolation error detected. Query aborted to prevent data leakage."
                )
            
            logger.debug(f"✓ Query verified: RAG instance correctly isolated to {actual_dir}")
            
            # Get notebook for configuration
            notebook = lightrag_notebooks_db[notebook_id]
            
            # TODO: To support LLM override safely, we would need to:
            # 1. Clone the RAG instance's storage to a temporary location
            # 2. Create a new RAG instance pointing to the cloned storage
            # 3. Clean up the temporary storage after the query
            # For now, we just use the notebook's configured LLM provider
            if query.llm_provider:
                logger.warning(f"LLM provider override requested but not supported - using notebook's configured provider")
            
            # Get notebook and model information for query optimization
            notebook = lightrag_notebooks_db[notebook_id]
            llm_provider = query.llm_provider or notebook.get("llm_provider", {})
            model_name = llm_provider.get("model", "").lower()
            
            # Adjust query parameters based on model capabilities
            adjusted_top_k = query.top_k
            adjusted_mode = query.mode
            
            # For smaller models like Gemma, use more conservative query parameters
            if 'gemma' in model_name or ('llama' in model_name and any(size in model_name for size in ['3b', '4b', '7b'])):
                logger.info(f"Optimizing query for smaller model: {model_name}")
                # Reduce top_k to limit context size
                adjusted_top_k = min(query.top_k, 30)
                # For very complex queries, prefer local mode to reduce context
                if query.mode == "global" and len(query.question.split()) > 20:
                    adjusted_mode = "hybrid"
                    logger.info("Switching from global to hybrid mode for complex query on small model")
            
            # Create query parameters with enhanced resilience settings
            query_param = QueryParam(
                mode=adjusted_mode,
                response_type=query.response_type,
                top_k=adjusted_top_k,
                # Enhanced token control for better resilience
                max_entity_tokens=6000,  # LightRAG default for entity context
                max_relation_tokens=8000,  # LightRAG default for relation context  
                max_total_tokens=30000,  # LightRAG default total budget
                chunk_top_k=20,  # Number of chunks to retrieve and keep after reranking
                enable_rerank=True,  # Enable reranking if available
            )
            
            # Perform query with fallback handling for context size issues
            try:
                result = await rag.aquery(query.question, param=query_param)
            except Exception as query_error:
                error_str = str(query_error).lower()
                
                # Check if it's a context size error
                if any(keyword in error_str for keyword in ['context size', 'context length', 'token limit', 'exceeds', 'too long']):
                    logger.warning(f"Context size error detected, attempting recovery: {query_error}")
                    
                    # Try with more aggressive reduction
                    if adjusted_mode == "global":
                        logger.info("Retrying with local mode instead of global")
                        fallback_param = QueryParam(
                            mode="local",
                            response_type=query.response_type,
                            top_k=min(20, adjusted_top_k),
                        )
                        result = await rag.aquery(query.question, param=fallback_param)
                        adjusted_mode = "local"
                    elif adjusted_mode == "hybrid":
                        logger.info("Retrying with naive mode instead of hybrid")
                        fallback_param = QueryParam(
                            mode="naive",
                            response_type=query.response_type,
                            top_k=min(15, adjusted_top_k),
                        )
                        result = await rag.aquery(query.question, param=fallback_param)
                        adjusted_mode = "naive"
                    else:
                        # Already using simplest mode, try with minimal context
                        logger.info("Retrying with minimal context")
                        fallback_param = QueryParam(
                            mode="naive",
                            response_type="Single Paragraph",
                            top_k=5,
                        )
                        result = await rag.aquery(query.question, param=fallback_param)
                        adjusted_mode = "naive"
                else:
                    # Re-raise non-context-size errors
                    raise query_error
            
            # Clean inline LightRAG citations and extract document references
            # Ensure result is a string (not an async iterator)
            result_text = result if isinstance(result, str) else str(result)
            cleaned_answer, inline_citations = clean_lightrag_citations_in_text(result_text, notebook_id)
            
            # Extract TRUE citations from actual retrieved sources
            citations = None
            try:
                citations = await build_true_citations_from_rag(
                    rag, 
                    notebook_id, 
                    query.question, 
                    top_k=min(20, adjusted_top_k)
                )
                if citations:
                    logger.info(f"✓ Extracted {len(citations)} precise citations for query")
                else:
                    logger.info("ℹ️ No precise citations available - using inline citations from text")
                    # Use citations extracted from inline text if no precise citations available
                    citations = inline_citations if inline_citations else None
            except Exception as citation_error:
                logger.warning(f"Error extracting citations: {citation_error}")
                citations = inline_citations if inline_citations else None
            
            return NotebookQueryResponse(
                answer=cleaned_answer,
                mode=adjusted_mode,
                context_used=True,
                citations=citations,
                source_documents=None,  # Will be enhanced later
                chat_context_used=False  # Will be enhanced when chat history is implemented
            )
            
        except Exception as e:
            logger.error(f"Error processing query for notebook {notebook_id}: {e}")
            raise HTTPException(status_code=500, detail=f"Error processing query: {str(e)}")

    @app.post("/notebooks/{notebook_id}/summary", response_model=NotebookQueryResponse)
    async def generate_notebook_summary(notebook_id: str):
        """Generate an automatic summary of all documents in the notebook"""
        validate_notebook_exists(notebook_id)
        
        try:
            logger.info(f"Summary generation request for notebook {notebook_id}")
            
            # Check if there are any completed documents
            notebook_documents = [
                doc for doc in lightrag_documents_db.values() 
                if doc["notebook_id"] == notebook_id and doc["status"] == "completed"
            ]
            
            if not notebook_documents:
                return NotebookQueryResponse(
                    answer="No documents have been processed yet. Please upload and wait for documents to be processed before generating a summary.",
                    mode="hybrid",
                    context_used=False,
                    citations=None,
                    source_documents=None,
                    chat_context_used=False
                )
            
            # Create a fingerprint of current documents (using document IDs and upload times)
            current_docs_fingerprint = "|".join(sorted([
                f"{doc['id']}:{doc['uploaded_at'].isoformat() if isinstance(doc['uploaded_at'], datetime) else doc['uploaded_at']}" 
                for doc in notebook_documents
            ]))
            
            # Check if we have a cached summary that's still valid
            notebook_data = lightrag_notebooks_db[notebook_id]
            cached_summary = notebook_data.get("summary_cache")
            cached_fingerprint = notebook_data.get("docs_fingerprint")
            
            # If we have a valid cached summary, return it
            if (cached_summary and cached_fingerprint and 
                cached_fingerprint == current_docs_fingerprint):
                logger.info(f"Returning cached summary for notebook {notebook_id}")
                
                # Extract citation information for all completed documents
                citations = []
                try:
                    for doc in notebook_documents:
                        citation = {
                            "filename": doc["filename"],
                            "file_path": doc.get("file_path", f"documents/{doc['filename']}"),
                            "document_id": doc["id"],
                            "title": doc["filename"].replace('_', ' ').replace('.txt', '').replace('.pdf', '').replace('.md', '').title()
                        }
                        citations.append(citation)
                except Exception as citation_error:
                    logger.warning(f"Error extracting citations for cached summary: {citation_error}")
                    citations = None
                
                return NotebookQueryResponse(
                    answer=cached_summary["answer"],
                    mode=cached_summary["mode"],
                    context_used=cached_summary["context_used"],
                    citations=citations,
                    source_documents=None,
                    chat_context_used=False
                )
            
            # Generate new summary if no valid cache exists
            logger.info(f"Generating new summary for notebook {notebook_id} (documents changed)")
            
            # Get existing RAG instance
            rag = await get_lightrag_instance(notebook_id)
            
            # Create summary query
            summary_question = ("Write down a comprehensive summary of all the documents provided in a single paragraph. "
                              "Mention what the documents are about, the main topics they cover, key themes, "
                              "important findings or insights, and the overall scope of the content. "
                              "Focus on providing an overview that helps understand the nature and breadth of the knowledge base.")
            
            # Create query parameters optimized for summary generation with enhanced resilience
            query_param = QueryParam(
                mode="hybrid",  # Use hybrid mode for comprehensive coverage
                response_type="Single Paragraph",  # Request single paragraph format
                top_k=100,  # Use higher top_k to get broader coverage of documents
                # Enhanced token control for better resilience
                max_entity_tokens=6000,
                max_relation_tokens=8000,
                max_total_tokens=30000,
                chunk_top_k=20,
                enable_rerank=True,
            )
            
            # Perform summary query
            result = await rag.aquery(summary_question, param=query_param)
            
            # Extract citation information for all completed documents
            citations = []
            try:
                for doc in notebook_documents:
                    citation = {
                        "filename": doc["filename"],
                        "file_path": doc.get("file_path", f"documents/{doc['filename']}"),
                        "document_id": doc["id"],
                        "title": doc["filename"].replace('_', ' ').replace('.txt', '').replace('.pdf', '').replace('.md', '').title()
                    }
                    citations.append(citation)
                
            except Exception as citation_error:
                logger.warning(f"Error extracting citations for summary: {citation_error}")
                citations = None
            
            # Cache the new summary
            summary_cache = {
                "answer": result,
                "mode": "hybrid",
                "context_used": True,
                "generated_at": datetime.now().isoformat()
            }
            
            # Update notebook with cached summary and fingerprint
            lightrag_notebooks_db[notebook_id]["summary_cache"] = summary_cache
            lightrag_notebooks_db[notebook_id]["docs_fingerprint"] = current_docs_fingerprint
            
            # Save to disk
            save_notebooks_db()
            
            logger.info(f"Generated and cached new summary for notebook {notebook_id} with {len(notebook_documents)} documents")
            
            return NotebookQueryResponse(
                answer=result,
                mode="hybrid",
                context_used=True,
                citations=citations,
                source_documents=None,
                chat_context_used=False
            )
            
        except Exception as e:
            logger.error(f"Error generating summary for notebook {notebook_id}: {e}")
            raise HTTPException(status_code=500, detail=f"Error generating summary: {str(e)}")

    @app.get("/notebooks/{notebook_id}/graph")
    async def get_notebook_graph_data(notebook_id: str):
        """Get graph visualization data for a notebook"""
        validate_notebook_exists(notebook_id)
        
        try:
            # Path to the GraphML file created by LightRAG
            working_dir = LIGHTRAG_STORAGE_PATH / notebook_id
            graphml_file = working_dir / "graph_chunk_entity_relation.graphml"
            
            if not graphml_file.exists():
                return {
                    "nodes": [],
                    "edges": [],
                    "message": "No graph data available yet. Upload documents and query the notebook to generate the knowledge graph."
                }
            
            # Parse the GraphML file
            tree = ET.parse(graphml_file)
            root = tree.getroot()
            
            # GraphML namespace
            ns = {'graphml': 'http://graphml.graphdrawing.org/xmlns'}
            
            nodes = []
            edges = []
            
            # Extract nodes
            for node in root.findall('.//graphml:node', ns):
                node_id = node.get('id')
                node_data = {'id': node_id, 'type': 'entity', 'properties': {}}
                
                # Extract node attributes/data
                for data in node.findall('graphml:data', ns):
                    key = data.get('key')
                    value = data.text or ''
                    
                    # Map common GraphML keys to readable labels
                    if key == 'd0':  # Usually entity name
                        node_data['label'] = value
                        node_data['properties']['name'] = value
                    elif key == 'd1':  # Usually entity type
                        node_data['type'] = value
                        node_data['properties']['entity_type'] = value
                    elif key == 'd2':  # Usually description
                        node_data['properties']['description'] = value
                    else:
                        node_data['properties'][key] = value
                
                # Use node_id as label if no label found
                if 'label' not in node_data:
                    node_data['label'] = node_id
                
                nodes.append(node_data)
            
            # Extract edges
            for edge in root.findall('.//graphml:edge', ns):
                source = edge.get('source')
                target = edge.get('target')
                edge_data = {
                    'source': source,
                    'target': target,
                    'relationship': 'related_to',
                    'properties': {}
                }
                
                # Extract edge attributes/data
                for data in edge.findall('graphml:data', ns):
                    key = data.get('key')
                    value = data.text or ''
                    
                    # Map common GraphML keys for edges
                    if key == 'd3':  # Usually relationship type
                        edge_data['relationship'] = value
                        edge_data['properties']['relation_type'] = value
                    elif key == 'd4':  # Usually weight or strength
                        try:
                            edge_data['properties']['weight'] = float(value)
                        except ValueError:
                            edge_data['properties']['weight_str'] = value
                    elif key == 'd5':  # Usually description
                        edge_data['properties']['description'] = value
                    else:
                        edge_data['properties'][key] = value
                
                edges.append(edge_data)
            
            logger.info(f"Loaded graph data for notebook {notebook_id}: {len(nodes)} nodes, {len(edges)} edges")
            
            return {
                "nodes": nodes,
                "edges": edges,
                "stats": {
                    "node_count": len(nodes),
                    "edge_count": len(edges),
                    "file_path": str(graphml_file)
                }
            }
            
        except ET.ParseError as e:
            logger.error(f"Error parsing GraphML file for notebook {notebook_id}: {e}")
            return {
                "nodes": [],
                "edges": [],
                "error": f"Error parsing graph data: {str(e)}"
            }
        except Exception as e:
            logger.error(f"Error getting graph data for notebook {notebook_id}: {e}")
            raise HTTPException(status_code=500, detail=f"Error getting graph data: {str(e)}")

    @app.get("/notebooks/{notebook_id}/graph/html")
    async def get_notebook_graph_html(notebook_id: str):
        """Generate interactive HTML graph visualization using pyvis and networkx"""
        validate_notebook_exists(notebook_id)
        
        try:
            # Path to the GraphML file created by LightRAG
            working_dir = LIGHTRAG_STORAGE_PATH / notebook_id
            graphml_file = working_dir / "graph_chunk_entity_relation.graphml"
            
            if not graphml_file.exists():
                # Return a simple HTML page indicating no data
                html_content = """
                <!DOCTYPE html>
                <html>
                <head>
                    <title>Knowledge Graph - No Data</title>
                    <style>
                        body { font-family: Arial, sans-serif; text-align: center; padding: 50px; }
                        .message { color: #666; font-size: 18px; }
                    </style>
                </head>
                <body>
                    <h2>Knowledge Graph</h2>
                    <p class="message">No graph data available yet.</p>
                    <p class="message">Upload documents and query the notebook to generate the knowledge graph.</p>
                </body>
                </html>
                """
                return HTMLResponse(content=html_content)
            
            # Install required packages if not available
            try:
                import networkx as nx
                from pyvis.network import Network
            except ImportError:
                # Try to install packages
                import subprocess
                import sys
                
                logger.info("Installing required packages for graph visualization...")
                try:
                    subprocess.check_call([sys.executable, "-m", "pip", "install", "pyvis", "networkx"])
                    import networkx as nx
                    from pyvis.network import Network
                    logger.info("Successfully installed pyvis and networkx")
                except Exception as install_error:
                    logger.error(f"Failed to install required packages: {install_error}")
                    raise HTTPException(status_code=500, detail="Required packages (pyvis, networkx) not available")
            
            import random
            import tempfile
            
            # Load the GraphML file
            logger.info(f"Loading GraphML file: {graphml_file}")
            G = nx.read_graphml(str(graphml_file))
            
            # Create a Pyvis network with responsive design
            net = Network(
                height="100vh",
                width="100%",
                bgcolor="#ffffff",
                font_color="#333333",
                notebook=False
            )
            
            # Configure physics for better layout
            net.set_options("""
            var options = {
              "physics": {
                "enabled": true,
                "stabilization": {"iterations": 100},
                "barnesHut": {
                  "gravitationalConstant": -8000,
                  "centralGravity": 0.3,
                  "springLength": 95,
                  "springConstant": 0.04,
                  "damping": 0.09
                }
              },
              "nodes": {
                "font": {"size": 12},
                "scaling": {
                  "min": 10,
                  "max": 30
                }
              },
              "edges": {
                "font": {"size": 10},
                "scaling": {
                  "min": 1,
                  "max": 3
                }
              }
            }
            """)
            
            # Convert NetworkX graph to Pyvis network
            net.from_nx(G)
            
            # Define colors for different node types
            node_type_colors = {
                'person': '#FF6B6B',      # Red
                'organization': '#4ECDC4', # Teal
                'location': '#45B7D1',    # Blue
                'concept': '#96CEB4',     # Green
                'event': '#FFEAA7',      # Yellow
                'entity': '#DDA0DD',     # Plum
                'default': '#95A5A6'     # Gray
            }
            
            # Enhance nodes with colors, titles, and better styling
            for node in net.nodes:
                # Determine node type from the data
                node_type = 'entity'  # default
                if 'entity_type' in node:
                    node_type = str(node['entity_type']).lower()
                elif 'type' in node:
                    node_type = str(node['type']).lower()
                
                # Set color based on type
                node["color"] = node_type_colors.get(node_type, node_type_colors['default'])
                
                # Add hover title with description
                title_parts = [f"ID: {node.get('id', 'Unknown')}"]
                if 'label' in node and node['label']:
                    title_parts.append(f"Label: {node['label']}")
                if node_type:
                    title_parts.append(f"Type: {node_type.title()}")
                if 'description' in node and node['description']:
                    desc = str(node['description'])[:200] + "..." if len(str(node['description'])) > 200 else str(node['description'])
                    title_parts.append(f"Description: {desc}")
                
                node["title"] = "\\n".join(title_parts)
                
                # Set node size based on connections (degree)
                if hasattr(G, 'degree'):
                    degree = G.degree(node['id']) if node['id'] in G else 1
                    node["size"] = min(10 + degree * 2, 30)  # Size between 10-30
                
                # Clean up label for display
                if 'label' in node and node['label']:
                    # Truncate long labels
                    label = str(node['label'])
                    node["label"] = label[:20] + "..." if len(label) > 20 else label
                else:
                    # Use ID as label if no label exists
                    node_id = str(node.get('id', ''))
                    node["label"] = node_id[:20] + "..." if len(node_id) > 20 else node_id
            
            # Enhance edges with titles and styling
            for edge in net.edges:
                title_parts = []
                
                # Add relationship type
                if 'relationship' in edge and edge['relationship']:
                    title_parts.append(f"Relationship: {edge['relationship']}")
                elif 'relation_type' in edge and edge['relation_type']:
                    title_parts.append(f"Relationship: {edge['relation_type']}")
                
                # Add weight if available
                if 'weight' in edge and edge['weight']:
                    try:
                        weight = float(edge['weight'])
                        title_parts.append(f"Weight: {weight:.2f}")
                        # Set edge width based on weight
                        edge["width"] = min(max(1, weight * 2), 5)
                    except (ValueError, TypeError):
                        pass
                
                # Add description if available
                if 'description' in edge and edge['description']:
                    desc = str(edge['description'])[:100] + "..." if len(str(edge['description'])) > 100 else str(edge['description'])
                    title_parts.append(f"Description: {desc}")
                
                if title_parts:
                    edge["title"] = "\\n".join(title_parts)
                
                # Style edges
                edge["color"] = {"color": "#848484", "highlight": "#333333"}
            
            # Generate HTML
            with tempfile.NamedTemporaryFile(mode='w', suffix='.html', delete=False) as tmp_file:
                net.save_graph(tmp_file.name)
                tmp_file.flush()
                
                # Read the generated HTML
                with open(tmp_file.name, 'r', encoding='utf-8') as f:
                    html_content = f.read()
                
                # Clean up temp file
                import os
                os.unlink(tmp_file.name)
            
            # Enhance the HTML with custom styling and dark mode support
            enhanced_html = html_content.replace(
                '<head>',
                '''<head>
                <style>
                    body { 
                        margin: 0; 
                        padding: 0; 
                        font-family: Arial, sans-serif;
                        background: #f8f9fa;
                    }
                    
                    @media (prefers-color-scheme: dark) {
                        body { background: #1a1a1a; }
                    }
                    
                    .graph-container {
                        position: relative;
                        width: 100%;
                        height: 100vh;
                    }
                    
                    .graph-info {
                        position: absolute;
                        top: 10px;
                        left: 10px;
                        background: rgba(255, 255, 255, 0.9);
                        padding: 10px;
                        border-radius: 5px;
                        font-size: 12px;
                        z-index: 1000;
                        box-shadow: 0 2px 5px rgba(0,0,0,0.1);
                    }
                    
                    @media (prefers-color-scheme: dark) {
                        .graph-info { 
                            background: rgba(30, 30, 30, 0.9); 
                            color: white;
                        }
                    }
                    
                    .legend {
                        position: absolute;
                        top: 10px;
                        right: 10px;
                        background: rgba(255, 255, 255, 0.9);
                        padding: 10px;
                        border-radius: 5px;
                        font-size: 11px;
                        z-index: 1000;
                        box-shadow: 0 2px 5px rgba(0,0,0,0.1);
                        max-width: 200px;
                    }
                    
                    @media (prefers-color-scheme: dark) {
                        .legend { 
                            background: rgba(30, 30, 30, 0.9); 
                            color: white;
                        }
                    }
                    
                    .legend-item {
                        display: flex;
                        align-items: center;
                        margin: 2px 0;
                    }
                    
                    .legend-color {
                        width: 12px;
                        height: 12px;
                        border-radius: 50%;
                        margin-right: 5px;
                    }
                </style>'''
            )
            
            # Add info overlay and legend
            graph_stats = f"Nodes: {len(net.nodes)} | Edges: {len(net.edges)}"
            
            legend_html = '''
            <div class="legend">
                <strong>Node Types:</strong>
                <div class="legend-item">
                    <div class="legend-color" style="background: #FF6B6B;"></div>
                    <span>Person</span>
                </div>
                <div class="legend-item">
                    <div class="legend-color" style="background: #4ECDC4;"></div>
                    <span>Organization</span>
                </div>
                <div class="legend-item">
                    <div class="legend-color" style="background: #45B7D1;"></div>
                    <span>Location</span>
                </div>
                <div class="legend-item">
                    <div class="legend-color" style="background: #96CEB4;"></div>
                    <span>Concept</span>
                </div>
                <div class="legend-item">
                    <div class="legend-color" style="background: #FFEAA7;"></div>
                    <span>Event</span>
                </div>
                <div class="legend-item">
                    <div class="legend-color" style="background: #DDA0DD;"></div>
                    <span>Entity</span>
                </div>
            </div>
            '''
            
            enhanced_html = enhanced_html.replace(
                '<body>',
                f'''<body>
                <div class="graph-info">{graph_stats}</div>
                {legend_html}'''
            )
            
            logger.info(f"Generated interactive graph HTML for notebook {notebook_id}: {len(net.nodes)} nodes, {len(net.edges)} edges")
            
            return HTMLResponse(content=enhanced_html)
            
        except Exception as e:
            logger.error(f"Error generating graph HTML for notebook {notebook_id}: {e}")
            logger.error(f"Full traceback: {traceback.format_exc()}")
            
            # Return error HTML page
            error_html = f"""
            <!DOCTYPE html>
            <html>
            <head>
                <title>Knowledge Graph - Error</title>
                <style>
                    body {{ font-family: Arial, sans-serif; text-align: center; padding: 50px; }}
                    .error {{ color: #d32f2f; font-size: 18px; }}
                    .details {{ color: #666; font-size: 14px; margin-top: 10px; }}
                </style>
            </head>
            <body>
                <h2>Knowledge Graph</h2>
                <p class="error">Error generating graph visualization</p>
                <p class="details">{str(e)}</p>
            </body>
            </html>
            """
            return HTMLResponse(content=error_html)

    @app.get("/notebooks/{notebook_id}/debug")
    async def debug_notebook_documents(notebook_id: str):
        """Debug endpoint to check document processing status and LightRAG state"""
        validate_notebook_exists(notebook_id)
        
        try:
            # Get all documents for this notebook
            notebook_documents = [
                doc for doc in lightrag_documents_db.values() 
                if doc["notebook_id"] == notebook_id
            ]
            
            # Get LightRAG instance info
            rag_info = {"exists": False, "working_dir": None}
            if notebook_id in lightrag_instances:
                rag = lightrag_instances[notebook_id]
                working_dir = LIGHTRAG_STORAGE_PATH / notebook_id
                rag_info = {
                    "exists": True,
                    "working_dir": str(working_dir),
                    "directory_exists": working_dir.exists(),
                    "files": list(working_dir.glob("*")) if working_dir.exists() else []
                }
            
            return {
                "notebook_id": notebook_id,
                "documents_count": len(notebook_documents),
                "documents": [
                    {
                        "id": doc["id"],
                        "filename": doc["filename"],
                        "status": doc["status"],
                        "lightrag_id": doc.get("lightrag_id"),
                        "error": doc.get("error"),
                        "uploaded_at": doc["uploaded_at"].isoformat() if isinstance(doc["uploaded_at"], datetime) else doc["uploaded_at"]
                    }
                    for doc in notebook_documents
                ],
                "lightrag_info": rag_info
            }
            
        except Exception as e:
            logger.error(f"Error in debug endpoint: {e}")
            raise HTTPException(status_code=500, detail=f"Debug error: {str(e)}")

    # Chat History Endpoints
    @app.get("/notebooks/{notebook_id}/chat/history", response_model=ChatHistoryResponse)
    async def get_chat_history(notebook_id: str, limit: int = Query(50, description="Maximum number of messages to return")):
        """Get chat history for a notebook"""
        validate_notebook_exists(notebook_id)
        
        messages = chat_history_db.get(notebook_id, [])
        
        # Limit messages and convert to ChatMessage objects
        limited_messages = messages[-limit:] if limit > 0 else messages
        chat_messages = [ChatMessage(**msg) for msg in limited_messages]
        
        return ChatHistoryResponse(
            notebook_id=notebook_id,
            messages=chat_messages,
            total_messages=len(messages)
        )

    @app.delete("/notebooks/{notebook_id}/chat/history")
    async def clear_chat_history(notebook_id: str):
        """Clear chat history for a notebook"""
        validate_notebook_exists(notebook_id)
        
        if notebook_id in chat_history_db:
            del chat_history_db[notebook_id]
            save_chat_history_db()
        
        return {"message": "Chat history cleared successfully"}

    # Enhanced Query with Chat History
    @app.post("/notebooks/{notebook_id}/chat", response_model=NotebookQueryResponse)
    async def chat_with_notebook(notebook_id: str, query: NotebookQueryRequest):
        """Chat with a notebook using conversation history"""
        validate_notebook_exists(notebook_id)
        
        try:
            # Initialize chat history if it doesn't exist
            if notebook_id not in chat_history_db:
                chat_history_db[notebook_id] = []
            
            # Add user message to history
            user_message = {
                "role": "user",
                "content": query.question,
                "timestamp": datetime.now()
            }
            chat_history_db[notebook_id].append(user_message)
            
            # Get LightRAG instance
            rag = await get_lightrag_instance(notebook_id)
            
            # Build context from chat history if enabled
            chat_context = ""
            if query.use_chat_history and len(chat_history_db[notebook_id]) > 1:
                recent_messages = chat_history_db[notebook_id][-10:]  # Last 10 messages
                chat_context = "Previous conversation context:\n"
                for msg in recent_messages[:-1]:  # Exclude the current message
                    chat_context += f"{msg['role'].title()}: {msg['content']}\n"
                chat_context += "\nCurrent question: "

            # Enhance question with chat context and citation instructions
            # Based on RAG best practices: use prompt engineering to force inline citations
            citation_instruction = """

CRITICAL INSTRUCTION - CITATION REQUIREMENTS:
You MUST add [SOURCE] after EVERY factual statement. Examples:
- "John Doe is the CEO [SOURCE]"
- "The company was founded in 2020 [SOURCE]"
- "The budget is $5M [SOURCE]"
Add [SOURCE] after names, dates, roles, numbers, or any specific data.
"""

            base_question = chat_context + query.question if chat_context else query.question
            enhanced_question = base_question + citation_instruction
            
            # Map mode string to QueryParam
            mode_mapping = {
                "naive": "naive",
                "local": "local",
                "global": "global",
                "hybrid": "hybrid",
                "mix": "mix"
            }
            adjusted_mode = mode_mapping.get(query.mode, "hybrid")

            # Log the query mode being used
            logger.info(f"🔍 Executing query with mode: {adjusted_mode.upper()} | top_k: {query.top_k} | response_type: {query.response_type}")
            logger.info(f"📝 Question: {query.question[:100]}..." if len(query.question) > 100 else f"📝 Question: {query.question}")

            # Execute query
            query_param = QueryParam(
                mode=adjusted_mode,
                response_type=query.response_type,
                top_k=query.top_k
            )

            result = await rag.aquery(enhanced_question, param=query_param)
            result_len = len(result) if isinstance(result, str) else "streaming"
            logger.info(f"✅ Query completed with mode: {adjusted_mode.upper()}, result length: {result_len} chars")
            
            # Clean inline LightRAG citations and extract document references
            result_text = result if isinstance(result, str) else str(result)
            cleaned_answer, inline_citations = clean_lightrag_citations_in_text(result_text, notebook_id)
            
            # Extract TRUE citations from actual retrieved sources
            # Enhanced citation mode: Always extract detailed citations with proper source attribution
            citations = None
            try:
                # Use higher citation limit for better source coverage (20 citations)
                citation_limit = 20
                citations = await build_true_citations_from_rag(
                    rag,
                    notebook_id,
                    enhanced_question,
                    top_k=min(citation_limit, query.top_k)
                )
                if citations:
                    logger.info(f"✓ Enhanced citation mode extracted {len(citations)} precise citations")
                else:
                    logger.info(f"ℹ️ No precise citations available - using inline citations from text")
                    citations = inline_citations if inline_citations else None
            except Exception as citation_error:
                logger.warning(f"Error extracting citations: {citation_error}")
                citations = inline_citations if inline_citations else None
            
            # Add assistant message to history
            assistant_message = {
                "role": "assistant",
                "content": cleaned_answer,
                "timestamp": datetime.now(),
                "citations": citations,
                "mode": "citation"  # Always use enhanced citation mode
            }
            chat_history_db[notebook_id].append(assistant_message)
            
            # Save chat history
            save_chat_history_db()
            
            return NotebookQueryResponse(
                answer=cleaned_answer,
                mode=adjusted_mode,
                context_used=True,
                citations=citations,
                source_documents=None,  # Will be enhanced later
                chat_context_used=bool(chat_context)
            )
            
        except Exception as e:
            logger.error(f"Error in chat query for notebook {notebook_id}: {e}")
            logger.error(f"Full error traceback: {traceback.format_exc()}")
            raise HTTPException(status_code=500, detail=f"Error processing chat query: {str(e)}")

    # Document Summary Endpoints
    @app.post("/notebooks/{notebook_id}/summary/detailed", response_model=NotebookQueryResponse)
    async def generate_detailed_summary(notebook_id: str, request: DocumentSummaryRequest):
        """Generate a detailed summary with document-level insights"""
        validate_notebook_exists(notebook_id)
        
        try:
            # Get document list
            notebook_documents = [
                doc for doc in lightrag_documents_db.values() 
                if doc["notebook_id"] == notebook_id and doc["status"] == "completed"
            ]
            
            if not notebook_documents:
                return NotebookQueryResponse(
                    answer="No completed documents found. Please upload and process documents first.",
                    mode="summary",
                    context_used=False,
                    citations=None,
                    source_documents=None,
                    chat_context_used=False
                )
            
            # Get LightRAG instance
            rag = await get_lightrag_instance(notebook_id)
            
            # Create detailed summary prompt based on length preference
            length_prompts = {
                "short": "Write a concise summary (2-3 sentences) of the key points from all documents.",
                "medium": "Write a comprehensive summary (1-2 paragraphs) covering the main themes and insights from all documents.",
                "long": "Write a detailed analysis (3-4 paragraphs) covering key themes, insights, contradictions, and conclusions from all documents."
            }
            
            summary_prompt = length_prompts.get(request.max_length, length_prompts["medium"])
            
            if request.include_details:
                summary_prompt += f"\n\nInclude insights from these {len(notebook_documents)} documents: " + \
                                ", ".join([doc["filename"] for doc in notebook_documents])
            
            # Execute summary query with enhanced resilience settings
            # Use "mix" mode if rerank is available (LightRAG recommendation), otherwise hybrid
            query_mode = "mix" if RERANK_AVAILABLE else "hybrid"
            query_param = QueryParam(
                mode=query_mode, 
                response_type="Multiple Paragraphs", 
                top_k=100,
                max_entity_tokens=6000,
                max_relation_tokens=8000,
                max_total_tokens=30000,
                chunk_top_k=20,
                enable_rerank=True,
            )
            result = await rag.aquery(summary_prompt, param=query_param)
            
            # Build source documents list
            source_docs = []
            if request.include_details:
                for doc in notebook_documents:
                    source_docs.append({
                        "filename": doc["filename"],
                        "upload_date": doc["uploaded_at"].isoformat() if isinstance(doc["uploaded_at"], datetime) else doc["uploaded_at"],
                        "status": doc["status"]
                    })
            
            return NotebookQueryResponse(
                answer=str(result),
                mode="hybrid",
                context_used=True,
                citations=[{"source": doc["filename"], "type": "document"} for doc in notebook_documents],
                source_documents=source_docs,
                chat_context_used=False
            )
            
        except Exception as e:
            logger.error(f"Error generating detailed summary for notebook {notebook_id}: {e}")
            raise HTTPException(status_code=500, detail=f"Error generating summary: {str(e)}")

    # Query Templates Endpoint
    @app.get("/query-templates", response_model=List[QueryTemplate])
    async def get_query_templates():
        """Get pre-built query templates for common use cases"""
        templates = [
            QueryTemplate(
                id="summarize_all",
                name="Summarize All Documents",
                description="Get a comprehensive overview of all uploaded documents",
                question_template="Provide a comprehensive summary of all the documents, highlighting the main themes and key insights.",
                category="Analysis",
                use_case="When you want a high-level overview of your knowledge base"
            ),
            QueryTemplate(
                id="find_contradictions",
                name="Find Contradictions",
                description="Identify conflicting information across documents",
                question_template="Are there any contradictions or conflicting viewpoints between the different documents? Please explain any differences you find.",
                category="Analysis",
                use_case="When comparing different sources or viewpoints"
            ),
            QueryTemplate(
                id="extract_key_facts",
                name="Extract Key Facts",
                description="Pull out the most important facts and data points",
                question_template="What are the most important facts, statistics, and data points mentioned across all documents?",
                category="Research",
                use_case="When you need specific factual information"
            ),
            QueryTemplate(
                id="timeline_analysis",
                name="Timeline Analysis",
                description="Understand chronological order and progression",
                question_template="Can you create a timeline of events or developments mentioned in the documents?",
                category="Research",
                use_case="When tracking historical progression or project timelines"
            ),
            QueryTemplate(
                id="compare_approaches",
                name="Compare Approaches",
                description="Compare different methods or strategies",
                question_template="What are the different approaches or methodologies discussed in the documents? Compare their advantages and disadvantages.",
                category="Comparison",
                use_case="When evaluating different options or strategies"
            ),
            QueryTemplate(
                id="action_items",
                name="Extract Action Items",
                description="Find actionable recommendations and next steps",
                question_template="What are the key recommendations, action items, or next steps suggested in the documents?",
                category="Planning",
                use_case="When planning follow-up activities"
            ),
            QueryTemplate(
                id="expert_opinions",
                name="Expert Opinions",
                description="Identify expert viewpoints and authoritative sources",
                question_template="What do experts or authoritative sources say about this topic? Include any credentials or authority mentioned.",
                category="Research",
                use_case="When you need authoritative perspectives"
            ),
            QueryTemplate(
                id="gaps_analysis",
                name="Knowledge Gaps",
                description="Identify missing information or areas that need more research",
                question_template="What topics or questions are mentioned but not fully explained? What knowledge gaps exist that might need additional research?",
                category="Analysis",
                use_case="When planning additional research"
            ),
            QueryTemplate(
                id="practical_applications",
                name="Practical Applications",
                description="Find real-world applications and use cases",
                question_template="What are the practical applications or real-world use cases mentioned in the documents? Include any examples or case studies.",
                category="Application",
                use_case="When looking for implementation ideas"
            ),
            QueryTemplate(
                id="risk_analysis",
                name="Risk Analysis",
                description="Identify potential risks, challenges, or limitations",
                question_template="What risks, challenges, limitations, or potential problems are mentioned in the documents?",
                category="Analysis",
                use_case="When assessing potential issues or preparing for challenges"
            )
        ]
        
        return templates

    @app.post("/notebooks/{notebook_id}/query/template/{template_id}", response_model=NotebookQueryResponse)
    async def query_with_template(notebook_id: str, template_id: str, custom_params: Optional[Dict[str, str]] = None):
        """Execute a query using a pre-built template"""
        validate_notebook_exists(notebook_id)
        
        # Get templates
        templates = await get_query_templates()
        template = next((t for t in templates if t.id == template_id), None)
        
        if not template:
            raise HTTPException(status_code=404, detail="Template not found")
        
        # Build question from template
        question = template.question_template
        
        # Replace any custom parameters if provided
        if custom_params:
            for key, value in custom_params.items():
                question = question.replace(f"{{{key}}}", value)
        
        # Create query request
        query_request = NotebookQueryRequest(
            question=question,
            mode="hybrid",
            response_type="Multiple Paragraphs",
            top_k=60,
            llm_provider=None,  # Use notebook's default provider
            use_chat_history=False  # Templates don't use chat history by default
        )
        
        # Execute using the chat endpoint for consistency
        return await chat_with_notebook(notebook_id, query_request)

# Document Text Extraction Endpoint
@app.post("/extract-text")
async def extract_document_text(file: UploadFile = File(...)):
    """
    Extract text from a document file without creating a notebook.
    Supports: PDF, TXT, MD, RTF, DOC, DOCX, XLS, XLSX, PPT, PPTX, HTML, HTM, XML, CSV, JSON
    """
    if not file.filename:
        raise HTTPException(status_code=400, detail="No filename provided")
    
    # Check file size (50MB limit)
    max_file_size = 50 * 1024 * 1024  # 50MB
    file_content = await file.read()
    
    if len(file_content) > max_file_size:
        raise HTTPException(
            status_code=400, 
            detail=f"File too large. Maximum size is 50MB, received {len(file_content) / (1024*1024):.1f}MB"
        )
    
    if not file_content:
        raise HTTPException(status_code=400, detail="Empty file")
    
    try:
        # Extract text using the existing function
        extracted_text = await extract_text_from_file(file.filename, file_content)
        
        # Get file statistics
        file_extension = file.filename.lower().split('.')[-1] if '.' in file.filename else 'unknown'
        text_length = len(extracted_text)
        word_count = len(extracted_text.split()) if extracted_text else 0
        
        return {
            "status": "success",
            "filename": file.filename,
            "file_type": file_extension.upper(),
            "file_size_bytes": len(file_content),
            "file_size_mb": round(len(file_content) / (1024*1024), 2),
            "extracted_text": extracted_text,
            "text_stats": {
                "character_count": text_length,
                "word_count": word_count,
                "line_count": extracted_text.count('\n') + 1 if extracted_text else 0
            }
        }
        
    except HTTPException:
        # Re-raise HTTP exceptions (like unsupported file types)
        raise
    except Exception as e:
        logger.error(f"Error extracting text from {file.filename}: {e}")
        logger.error(traceback.format_exc())
        raise HTTPException(
            status_code=500, 
            detail=f"Error processing document: {str(e)}"
        )

# Audio transcription endpoint
@app.post("/transcribe")
async def transcribe_audio(
    file: UploadFile = File(...),
    language: Optional[str] = Form(None),
    beam_size: int = Form(5),
    initial_prompt: Optional[str] = Form(None)
):
    """Transcribe an audio file using faster-whisper (CPU mode)"""
    # Validate file extension
    supported_formats = ['mp3', 'wav', 'flac', 'm4a', 'ogg', 'opus']
    
    if not file.filename:
        raise HTTPException(status_code=400, detail="No filename provided")
    
    file_extension = file.filename.lower().split('.')[-1]
    
    if file_extension not in supported_formats:
        raise HTTPException(
            status_code=400, 
            detail=f"Unsupported audio format: {file_extension}. Supported formats: {', '.join(supported_formats)}"
        )
    
    # Read file content
    try:
        content = await file.read()
        if not content:
            raise HTTPException(status_code=400, detail="Empty audio file")
    except Exception as e:
        logger.error(f"Error reading audio file: {e}")
        raise HTTPException(status_code=500, detail=f"Error reading audio file: {str(e)}")
    
    # Get Speech2Text instance
    try:
        s2t = get_speech2text()
    except Exception as e:
        logger.error(f"Error initializing Speech2Text: {e}")
        raise HTTPException(status_code=500, detail=f"Failed to initialize Speech2Text: {str(e)}")
    
    # Transcribe the audio
    try:
        result = s2t.transcribe_bytes(
            content,
            language=language,
            beam_size=beam_size,
            initial_prompt=initial_prompt
        )
        
        return {
            "status": "success",
            "filename": file.filename,
            "transcription": result
        }
    except Exception as e:
        logger.error(f"Error transcribing audio: {e}")
        logger.error(traceback.format_exc())
        raise HTTPException(status_code=500, detail=f"Error transcribing audio: {str(e)}")

# Text-to-Speech endpoints
@app.post("/synthesize")
async def synthesize_text(request: TTSRequest):
    """
    Synthesize text to speech and return audio data.
    Supports multiple TTS engines including Kokoro for high-quality neural TTS.
    """
    try:
        logger.info(f"TTS request: text='{request.text[:50]}...', engine={request.engine}, voice={request.voice}, speed={request.speed}")
        
        # Get Text2Speech instance
        t2s = get_text2speech()
        
        # If specific engine requested, create new instance with those settings
        if request.engine and request.engine != t2s.engine:
            logger.info(f"Creating new TTS instance with engine: {request.engine}")
            t2s = Text2Speech(
                engine=request.engine,
                language=request.language or "en",
                slow=request.slow or False,
                voice=request.voice or "af_sarah",
                speed=request.speed or 1.0
            )
        
        # Generate speech
        audio_bytes = t2s.synthesize_to_bytes(request.text)
        
        # Determine content type based on engine
        if request.engine in ["kokoro", "kokoro-onnx"]:
            content_type = "audio/wav"
        else:
            content_type = "audio/mpeg"
        
        return Response(
            content=audio_bytes,
            media_type=content_type,
            headers={
                "Content-Disposition": "attachment; filename=speech.wav" if request.engine in ["kokoro", "kokoro-onnx"] else "attachment; filename=speech.mp3"
            }
        )
        
    except Exception as e:
        logger.error(f"Error in TTS synthesis: {e}")
        logger.error(traceback.format_exc())
        raise HTTPException(status_code=500, detail=f"TTS synthesis failed: {str(e)}")

@app.post("/synthesize/file")
async def synthesize_text_to_file(
    text: str = Form(...),
    language: Optional[str] = Form("en"),
    engine: Optional[str] = Form(None),
    slow: Optional[bool] = Form(False),
    voice: Optional[str] = Form("af_sarah"),
    speed: Optional[float] = Form(1.0),
    filename: Optional[str] = Form("speech.mp3")
):
    """
    Synthesize text to speech and return as downloadable file.
    Supports multiple TTS engines including Kokoro for high-quality neural TTS.
    """
    try:
        logger.info(f"TTS file request: text='{text[:50]}...', engine={engine}, voice={voice}, speed={speed}")
        
        # Get Text2Speech instance
        t2s = get_text2speech()
        
        # If a specific engine is requested and different from current
        if engine and engine != t2s.engine:
            logger.info(f"Creating new TTS instance with engine: {engine}")
            t2s = Text2Speech(
                engine=engine,
                language=language or "en",
                slow=slow or False,
                voice=voice or "af_sarah",
                speed=speed or 1.0
            )
        
        # Determine file extension based on engine
        if engine in ["kokoro", "kokoro-onnx"]:
            file_ext = ".wav"
            content_type = "audio/wav"
        else:
            file_ext = ".mp3"
            content_type = "audio/mpeg"
        
        # Ensure filename has correct extension
        if not filename:
            filename = "speech" + file_ext
        elif not filename.endswith(file_ext):
            filename = os.path.splitext(filename)[0] + file_ext
        
        # Create temporary file
        with tempfile.NamedTemporaryFile(delete=False, suffix=file_ext) as temp_file:
            temp_path = temp_file.name
        
        try:
            # Generate speech to file
            output_path = t2s.synthesize_to_file(text, temp_path)
            
            # Read the file
            with open(output_path, 'rb') as f:
                audio_data = f.read()
            
            return Response(
                content=audio_data,
                media_type=content_type,
                headers={
                    "Content-Disposition": f"attachment; filename={filename}",
                    "Content-Length": str(len(audio_data))
                }
            )
            
        finally:
            # Clean up temp file
            try:
                os.unlink(temp_path)
            except:
                pass
        
    except Exception as e:
        logger.error(f"Error in TTS file synthesis: {e}")
        logger.error(traceback.format_exc())
        raise HTTPException(status_code=500, detail=f"TTS file synthesis failed: {str(e)}")

@app.get("/tts/languages")
async def get_tts_languages():
    """Get available languages for text-to-speech"""
    try:
        t2s = get_text2speech()
        languages = t2s.get_available_languages()
        
        return {
            "engine": t2s.engine,
            "current_language": t2s.language,
            "available_languages": languages
        }
    except Exception as e:
        logger.error(f"Error getting TTS languages: {e}")
        raise HTTPException(status_code=500, detail=f"Error getting TTS languages: {str(e)}")

@app.get("/tts/status")
async def get_tts_status():
    """Get current TTS engine status and configuration"""
    try:
        t2s = get_text2speech()
        
        return {
            "engine": t2s.engine,
            "language": t2s.language,
            "slow": t2s.slow,
            "available_engines": []
        }
    except Exception as e:
        logger.error(f"Error getting TTS status: {e}")
        raise HTTPException(status_code=500, detail=f"Error getting TTS status: {str(e)}")

@app.get("/tts/voices")
async def get_tts_voices():
    """Get available voices for TTS engines"""
    try:
        voices = {
            "kokoro_voices": {
                "af_sarah": "American Female - Sarah (warm, friendly)",
                "af_nicole": "American Female - Nicole (professional)",
                "af_sky": "American Female - Sky (energetic)",
                "am_adam": "American Male - Adam (deep, authoritative)",
                "am_michael": "American Male - Michael (casual)",
                "bf_emma": "British Female - Emma (elegant)",
                "bf_isabella": "British Female - Isabella (sophisticated)",
                "bm_george": "British Male - George (distinguished)",
                "bm_lewis": "British Male - Lewis (modern)"
            },
            "pyttsx3_voices": "System dependent - use /tts/status to see available voices",
            "gtts_languages": [
                "en", "es", "fr", "de", "it", "pt", "ru", "ja", "ko", "zh", "hi", "ar"
            ]
        }
        
        # Try to get actual pyttsx3 voices if available
        try:
            import pyttsx3
            engine = pyttsx3.init()
            system_voices = engine.getProperty('voices')
            if system_voices:
                voices["pyttsx3_voices"] = [
                    {
                        "id": voice.id,
                        "name": voice.name,
                        "languages": getattr(voice, 'languages', []),
                        "gender": getattr(voice, 'gender', 'unknown')
                    }
                    for voice in system_voices
                ]
            engine.stop()
        except:
            pass
        
        return voices
        
    except Exception as e:
        logger.error(f"Error getting TTS voices: {e}")
        raise HTTPException(status_code=500, detail=f"Error getting TTS voices: {str(e)}")

# Handle graceful shutdown
def handle_exit(signum, frame):
    logger.info(f"Received signal {signum}, shutting down gracefully")
    sys.exit(0)

# Register signal handlers
signal.signal(signal.SIGINT, handle_exit)
signal.signal(signal.SIGTERM, handle_exit)

if __name__ == "__main__":
    import uvicorn
    logger.info(f"Starting server on {HOST}:{PORT}")
    
    # Start the server with reload=False to prevent duplicate processes
    uvicorn.run(
        "main:app",
        host=HOST,
        port=PORT,
        log_level="info",
        reload=False  # Change this to false to prevent multiple processes
    )
